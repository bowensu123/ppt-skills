"""Tests for run-level text formatting preservation:
  inspect_ppt._text_runs    →  per-shape runs[]
  apply_layout kind=rich_text →  multi-run rendering
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))


# ---- inspect_ppt _text_runs ----

def test_inspect_text_runs_captures_bold_color(tmp_path):
    """A text shape with two differently-styled runs produces 2 entries."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt
    import inspect_ppt as inspect_module

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Emu(457200), Emu(457200),
                                    Emu(8000000), Emu(457200))
    tf = box.text_frame
    tf.text = ""
    para = tf.paragraphs[0]
    r1 = para.add_run()
    r1.text = "Zero"
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xD9, 0x77, 0x57)
    r1.font.size = Pt(18)
    r2 = para.add_run()
    r2.text = "-shot"
    r2.font.bold = False
    r2.font.color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    r2.font.size = Pt(18)

    runs = inspect_module._text_runs(box)
    assert len(runs) == 2
    assert runs[0]["text"] == "Zero"
    assert runs[0]["bold"] is True
    assert runs[0]["color_hex"] == "D97757"
    assert runs[0]["size_pt"] == 18.0
    assert runs[1]["text"] == "-shot"
    assert runs[1]["bold"] is False
    assert runs[1]["color_hex"] == "F5F5F5"


def test_inspect_text_runs_empty_for_picture(tmp_path):
    """Picture shape returns no runs (no text frame)."""
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu
    import inspect_ppt as inspect_module

    Image.new("RGB", (8, 8), "white").save(tmp_path / "p.png", "PNG")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    pic = slide.shapes.add_picture(str(tmp_path / "p.png"),
                                    Emu(0), Emu(0), Emu(1000000), Emu(1000000))
    runs = inspect_module._text_runs(pic)
    assert runs == []


def test_inspect_runs_in_full_inspection_output(tmp_path):
    """End-to-end: inspect_ppt CLI emits text_runs in JSON."""
    import subprocess

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(8000000), Emu(500000))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Bold red"
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    r.font.size = Pt(20)

    deck = tmp_path / "d.pptx"; prs.save(deck)
    out = tmp_path / "ins.json"
    subprocess.run(
        [sys.executable, str(SKILL_ROOT / "scripts" / "inspect_ppt.py"),
         "--input", str(deck), "--output", str(out)],
        check=True,
    )
    data = json.loads(out.read_text(encoding="utf-8"))
    txt_obj = next(o for o in data["slides"][0]["objects"]
                    if o.get("text") == "Bold red")
    assert "text_runs" in txt_obj
    assert len(txt_obj["text_runs"]) == 1
    assert txt_obj["text_runs"][0]["bold"] is True
    assert txt_obj["text_runs"][0]["color_hex"] == "FF0000"


# ---- apply_layout kind=rich_text ----

def test_apply_layout_rich_text_inline_runs(tmp_path):
    """rich_text element with inline runs renders without error."""
    import apply_layout
    layout = {
        "elements": [{
            "kind": "rich_text",
            "bbox": [0, 0, 8000000, 500000],
            "runs": [
                {"text": "Bold ", "bold": True, "color_hex": "D97757", "size_pt": 18},
                {"text": "regular", "bold": False, "color_hex": "393939", "size_pt": 18},
            ],
            "default_size_pt": 18,
        }],
    }
    out = tmp_path / "out.pptx"
    result = apply_layout.render_layout(layout, {}, out)
    assert result["rendered"] == 1
    assert result["skipped"] == []


def test_apply_layout_rich_text_runs_ref(tmp_path):
    """rich_text element with `runs_ref` resolves a path into content.json."""
    import apply_layout
    content = {
        "items": [{
            "name_runs": [
                {"text": "Title ", "bold": True, "color_hex": "D97757"},
                {"text": "part", "bold": False, "color_hex": "F5F5F5"},
            ],
        }],
    }
    layout = {
        "elements": [{
            "kind": "rich_text",
            "bbox": [0, 0, 8000000, 500000],
            "runs_ref": "items.0.name_runs",
        }],
    }
    out = tmp_path / "out.pptx"
    result = apply_layout.render_layout(layout, content, out)
    assert result["rendered"] == 1


def test_apply_layout_rich_text_missing_runs_skipped(tmp_path):
    import apply_layout
    layout = {
        "elements": [{
            "kind": "rich_text",
            "bbox": [0, 0, 8000000, 500000],
            "runs_ref": "items.0.name_runs",   # not present in content
        }],
    }
    out = tmp_path / "out.pptx"
    result = apply_layout.render_layout(layout, {}, out)
    assert any(s["reason"] == "rich-text-no-runs" for s in result["skipped"])


def test_apply_layout_rich_text_multi_paragraph(tmp_path):
    """Runs with different paragraph_index land in different paragraphs."""
    import apply_layout
    layout = {
        "elements": [{
            "kind": "rich_text",
            "bbox": [0, 0, 8000000, 1000000],
            "runs": [
                {"text": "Para 1 line", "paragraph_index": 0, "size_pt": 14},
                {"text": "Para 2 line", "paragraph_index": 1, "size_pt": 14},
            ],
        }],
    }
    out = tmp_path / "out.pptx"
    result = apply_layout.render_layout(layout, {}, out)
    assert result["rendered"] == 1
    # Verify both paragraphs got rendered
    from pptx import Presentation
    prs = Presentation(str(out))
    box = next(s for s in prs.slides[0].shapes if s.has_text_frame
                and s.text_frame.paragraphs[0].text != "")
    assert len(box.text_frame.paragraphs) >= 2
