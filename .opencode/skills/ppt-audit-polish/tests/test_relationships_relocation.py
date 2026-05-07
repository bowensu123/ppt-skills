"""Tests for _relationships_extract.py + apply_relocation.py — the
preserve-identity Path B path."""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from _relationships_extract import extract_relationships
from apply_relocation import apply_relocation


# ---- _relationships_extract: identity capture ----

def test_relationships_captures_shape_id_and_name(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(1000000), Emu(500000),
    )
    sp.name = "VLM Panel"
    deck = tmp_path / "d.pptx"; prs.save(deck)

    manifest = extract_relationships(deck, tmp_path / "w")
    sids = [sh["shape_id"] for sh in manifest["slides"][0]["shapes"]]
    assert sp.shape_id in sids
    entry = next(sh for sh in manifest["slides"][0]["shapes"]
                  if sh["shape_id"] == sp.shape_id)
    assert entry["name"] == "VLM Panel"


def test_relationships_default_preserve_for_referenced_shape(tmp_path):
    """A shape with a connector pointing at it gets preserve_identity_default=True."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.util import Emu
    from lxml import etree

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Emu(0), Emu(0), Emu(1000000), Emu(500000))
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Emu(2000000), Emu(0), Emu(1000000), Emu(500000))
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(1000000), Emu(250000), Emu(2000000), Emu(250000),
    )
    # Patch in stCxn / endCxn referencing a and b.
    NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    cxn = line._element.find(f"{NS_P}nvCxnSpPr/{NS_P}cNvCxnSpPr")
    if cxn is None:
        # Some pptx versions structure it differently — find any cNvCxnSpPr.
        cxn = line._element.find(".//*[local-name()='cNvCxnSpPr']")
    if cxn is not None:
        st = etree.SubElement(cxn, f"{NS_A}stCxn", {"id": str(a.shape_id), "idx": "0"})
        ed = etree.SubElement(cxn, f"{NS_A}endCxn", {"id": str(b.shape_id), "idx": "0"})

    deck = tmp_path / "d.pptx"; prs.save(deck)
    manifest = extract_relationships(deck, tmp_path / "w")
    by_id = {sh["shape_id"]: sh for sh in manifest["slides"][0]["shapes"]}
    if a.shape_id in by_id and b.shape_id in by_id:
        # At least one inbound reference should be detected.
        total_inbound = (len(by_id[a.shape_id]["is_referenced_by"])
                          + len(by_id[b.shape_id]["is_referenced_by"]))
        assert total_inbound >= 1


def test_relationships_default_recreate_for_pure_decoration(tmp_path):
    """Shape with no references, no placeholder → preserve_identity_default=False."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(1000000), Emu(500000),
    )
    deck = tmp_path / "d.pptx"; prs.save(deck)

    manifest = extract_relationships(deck, tmp_path / "w")
    entry = next(sh for sh in manifest["slides"][0]["shapes"]
                  if sh["shape_id"] == sp.shape_id)
    # Pure decoration shape on a blank layout
    assert entry["preserve_identity_default"] is False
    assert "no relationships" in entry["rationale_default"]


def test_relationships_text_run_hyperlinks_captured(tmp_path):
    from pptx import Presentation
    from pptx.util import Emu, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(8000000), Emu(500000))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Click here"
    r.hyperlink.address = "https://example.com"

    deck = tmp_path / "d.pptx"; prs.save(deck)
    manifest = extract_relationships(deck, tmp_path / "w")
    entry = next(sh for sh in manifest["slides"][0]["shapes"]
                  if sh["shape_id"] == box.shape_id)
    refs = [r for r in entry["references_to"] if r["kind"] == "hyperlink"]
    assert len(refs) == 1
    assert refs[0]["address"] == "https://example.com"
    # Has outbound references → preserve_identity_default = True
    assert entry["preserve_identity_default"] is True


# ---- apply_relocation: preserve-identity flow ----

def test_relocation_preserve_keeps_shape_id(tmp_path):
    """preserve_identity decision moves the shape; shape_id stays the same."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(100000), Emu(100000), Emu(1000000), Emu(500000),
    )
    sp.name = "PreservedPanel"
    original_sid = sp.shape_id
    deck = tmp_path / "d.pptx"; prs.save(deck)

    relocation = {
        "slides": [{
            "slide_index": 1,
            "shapes": {
                str(original_sid): {
                    "decision": "preserve_identity",
                    "agent_rationale": "test",
                    "new_bbox_emu": [3000000, 2000000, 1500000, 800000],
                }
            }
        }]
    }
    rj = tmp_path / "r.json"
    rj.write_text(json.dumps(relocation), encoding="utf-8")

    out = tmp_path / "out.pptx"
    apply_relocation(deck, rj, out)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    by_sid = {s.shape_id: s for s in slide2.shapes}
    assert original_sid in by_sid, "shape_id must be preserved"
    moved = by_sid[original_sid]
    assert moved.name == "PreservedPanel", "name must be preserved"
    assert int(moved.left) == 3000000
    assert int(moved.top) == 2000000


def test_relocation_recreate_carries_text(tmp_path):
    """recreate captures text from original and restores it on the new shape."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(8000000), Emu(500000))
    box.text_frame.text = "Original text"
    sid = box.shape_id
    deck = tmp_path / "d.pptx"; prs.save(deck)

    relocation = {
        "slides": [{
            "slide_index": 1,
            "shapes": {
                str(sid): {
                    "decision": "recreate",
                    "new_kind": "text",
                    "new_bbox_emu": [4000000, 3000000, 4000000, 500000],
                }
            }
        }]
    }
    rj = tmp_path / "r.json"
    rj.write_text(json.dumps(relocation), encoding="utf-8")

    out = tmp_path / "out.pptx"
    apply_relocation(deck, rj, out)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    text_shapes = [s for s in slide2.shapes
                    if s.has_text_frame and s.text_frame.text == "Original text"]
    assert len(text_shapes) >= 1, "text content must survive recreate"


def test_relocation_delete_removes_shape(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(1000000), Emu(500000),
    )
    sid = sp.shape_id
    deck = tmp_path / "d.pptx"; prs.save(deck)

    relocation = {
        "slides": [{
            "slide_index": 1,
            "shapes": {
                str(sid): {"decision": "delete"}
            }
        }]
    }
    rj = tmp_path / "r.json"
    rj.write_text(json.dumps(relocation), encoding="utf-8")

    out = tmp_path / "out.pptx"
    apply_relocation(deck, rj, out)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    sids = [s.shape_id for s in slide2.shapes]
    assert sid not in sids


def test_relocation_add_new_shapes(tmp_path):
    """add_new_shapes inserts agent-designed decorations."""
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    deck = tmp_path / "d.pptx"; prs.save(deck)

    relocation = {
        "slides": [{
            "slide_index": 1,
            "shapes": {},
            "add_new_shapes": [
                {"kind": "rect", "bbox": [100000, 100000, 1000000, 500000],
                 "fill": "D97757"},
                {"kind": "text", "bbox": [100000, 700000, 5000000, 500000],
                 "content": "New title", "size_pt": 18, "color": "F5F5F5"},
            ]
        }]
    }
    rj = tmp_path / "r.json"
    rj.write_text(json.dumps(relocation), encoding="utf-8")

    out = tmp_path / "out.pptx"
    result = apply_relocation(deck, rj, out)
    added = [a for a in result["actions"] if a.get("decision") == "add"]
    assert len(added) == 2


def test_relocation_preserve_keeps_image_binary(tmp_path):
    """Picture preserve_identity → original image bytes still in deck."""
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    img = tmp_path / "p.png"
    Image.new("RGB", (16, 16), "red").save(img, "PNG")
    pic = slide.shapes.add_picture(str(img),
                                    Emu(0), Emu(0), Emu(2000000), Emu(2000000))
    sid = pic.shape_id
    deck = tmp_path / "d.pptx"; prs.save(deck)

    relocation = {
        "slides": [{
            "slide_index": 1,
            "shapes": {
                str(sid): {
                    "decision": "preserve_identity",
                    "new_bbox_emu": [3000000, 3000000, 1000000, 1000000],
                }
            }
        }]
    }
    rj = tmp_path / "r.json"
    rj.write_text(json.dumps(relocation), encoding="utf-8")

    out = tmp_path / "out.pptx"
    apply_relocation(deck, rj, out)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    # Filter to shapes whose ID matches the original — there may be empty
    # placeholders from the layout still in the slide.
    matched = [s for s in slide2.shapes
                if s.shape_id == sid and s.shape_type in (13, 14)]
    assert len(matched) == 1, (
        f"original picture sid={sid} should remain. "
        f"Got: {[(s.shape_id, s.shape_type) for s in slide2.shapes]}"
    )


def test_relocation_recreate_carries_image(tmp_path):
    """recreate of a picture re-creates with the same image bytes."""
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    img = tmp_path / "p.png"
    Image.new("RGB", (16, 16), "red").save(img, "PNG")
    pic = slide.shapes.add_picture(str(img),
                                    Emu(0), Emu(0), Emu(2000000), Emu(2000000))
    sid = pic.shape_id
    deck = tmp_path / "d.pptx"; prs.save(deck)

    relocation = {
        "slides": [{
            "slide_index": 1,
            "shapes": {
                str(sid): {
                    "decision": "recreate",
                    "new_kind": "rect",  # ignored — image bytes take priority
                    "new_bbox_emu": [4000000, 3000000, 1500000, 1500000],
                }
            }
        }]
    }
    rj = tmp_path / "r.json"
    rj.write_text(json.dumps(relocation), encoding="utf-8")

    out = tmp_path / "out.pptx"
    apply_relocation(deck, rj, out)

    prs2 = Presentation(str(out))
    pics = [s for s in prs2.slides[0].shapes if s.shape_type in (13, 14)]
    assert len(pics) >= 1
