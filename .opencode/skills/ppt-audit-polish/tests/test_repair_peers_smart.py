"""Unit tests for _repair_peers_smart.py — agent-driven peer repair."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from _repair_peers_smart import repair_peers_smart


@pytest.fixture
def deck_with_3_cards(tmp_path):
    """3 cards laid out horizontally with deliberately uneven sizes/gaps."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Card 1: 1000 × 800 at (500, 1000)
    c1 = slide.shapes.add_textbox(Emu(500000), Emu(1000000), Emu(1000000), Emu(800000))
    c1.text_frame.text = "card 1"
    # Card 2: 1100 × 800 at (1700, 1100)  — slightly bigger, 200K below row top
    c2 = slide.shapes.add_textbox(Emu(1700000), Emu(1100000), Emu(1100000), Emu(800000))
    c2.text_frame.text = "card 2"
    # Card 3: 900 × 800 at (3500, 1000)  — wider gap before this one
    c3 = slide.shapes.add_textbox(Emu(3500000), Emu(1000000), Emu(900000), Emu(800000))
    c3.text_frame.text = "card 3"
    return prs, [c1.shape_id, c2.shape_id, c3.shape_id]


def test_uniform_size_resizes_all_to_median(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "row",
            "shape_ids": sids,
            "axis": "horizontal",
            "uniform_size": True,
            "uniform_spacing": False,
            "uniform_alignment": False,
        }],
    }
    result = repair_peers_smart(prs, groups)
    # Median of [1000K, 1100K, 900K] = 1000K
    for shape in prs.slides[0].shapes:
        if shape.shape_id in sids:
            assert int(shape.width) == 1000000
            assert int(shape.height) == 800000
    # Each shape that wasn't already at the median produced an action
    assert any(a["action"] == "equalize-size" for a in result["actions"])


def test_uniform_alignment_aligns_tops(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "row",
            "shape_ids": sids,
            "axis": "horizontal",
            "uniform_size": False,
            "uniform_spacing": False,
            "uniform_alignment": True,
        }],
    }
    repair_peers_smart(prs, groups)
    tops = [
        int(s.top) for s in prs.slides[0].shapes if s.shape_id in sids
    ]
    assert len(set(tops)) == 1, f"tops should be equal, got {tops}"


def test_uniform_spacing_distributes_to_median_gap(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "row",
            "shape_ids": sids,
            "axis": "horizontal",
            "uniform_size": True,        # equalize first so gap math is clean
            "uniform_spacing": True,
            "uniform_alignment": False,
        }],
    }
    repair_peers_smart(prs, groups)
    shapes = [s for s in prs.slides[0].shapes if s.shape_id in sids]
    shapes.sort(key=lambda s: int(s.left))
    gap1 = int(shapes[1].left) - (int(shapes[0].left) + int(shapes[0].width))
    gap2 = int(shapes[2].left) - (int(shapes[1].left) + int(shapes[1].width))
    assert gap1 == gap2, f"gaps should be equal after distribute: {gap1} vs {gap2}"


def test_explicit_target_size_overrides_median(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "row",
            "shape_ids": sids,
            "axis": "horizontal",
            "uniform_size": True,
            "target_size": [2000000, 1500000],
            "uniform_spacing": False,
            "uniform_alignment": False,
        }],
    }
    repair_peers_smart(prs, groups)
    for s in prs.slides[0].shapes:
        if s.shape_id in sids:
            assert int(s.width) == 2000000
            assert int(s.height) == 1500000


def test_explicit_target_gap_overrides_median(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "row",
            "shape_ids": sids,
            "axis": "horizontal",
            "uniform_size": True,
            "uniform_spacing": True,
            "target_gap_emu": 500000,
        }],
    }
    repair_peers_smart(prs, groups)
    shapes = [s for s in prs.slides[0].shapes if s.shape_id in sids]
    shapes.sort(key=lambda s: int(s.left))
    gap = int(shapes[1].left) - (int(shapes[0].left) + int(shapes[0].width))
    assert gap == 500000


def test_children_move_with_parent(deck_with_3_cards):
    """Children attached to a peer move with the parent during distribution."""
    prs, sids = deck_with_3_cards
    from pptx.util import Emu
    slide = prs.slides[0]
    # Add a small "icon" inside card 2's bbox — relative offset (50K, 100K)
    icon = slide.shapes.add_textbox(Emu(1750000), Emu(1200000),
                                     Emu(200000), Emu(200000))
    icon.text_frame.text = "icon"

    groups = {
        "groups": [{
            "name": "row",
            "shape_ids": sids,
            "axis": "horizontal",
            "uniform_size": True,
            "uniform_spacing": True,
            "uniform_alignment": True,
            "children_per_peer": [[], [icon.shape_id], []],
        }],
    }
    # Track relative offset before
    card2 = next(s for s in slide.shapes if s.shape_id == sids[1])
    rel_x_before = int(icon.left) - int(card2.left)
    rel_y_before = int(icon.top) - int(card2.top)

    repair_peers_smart(prs, groups)

    rel_x_after = int(icon.left) - int(card2.left)
    rel_y_after = int(icon.top) - int(card2.top)
    assert rel_x_before == rel_x_after, "child x-offset should be preserved"
    assert rel_y_before == rel_y_after, "child y-offset should be preserved"


def test_unknown_shape_ids_skipped_silently(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "row",
            "shape_ids": sids + [99999, 88888],
            "axis": "horizontal",
            "uniform_size": True,
        }],
    }
    # Should not raise — just process the real ones.
    result = repair_peers_smart(prs, groups)
    assert result["groups_processed"] == 1


def test_invalid_axis_skips_group(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "weird",
            "shape_ids": sids,
            "axis": "diagonal",
        }],
    }
    result = repair_peers_smart(prs, groups)
    assert result["groups_processed"] == 0
    assert any("invalid-axis" in s["reason"] for s in result["skipped"])


def test_fewer_than_2_shapes_skips(deck_with_3_cards):
    prs, sids = deck_with_3_cards
    groups = {
        "groups": [{
            "name": "lonely",
            "shape_ids": [sids[0]],
            "axis": "horizontal",
        }],
    }
    result = repair_peers_smart(prs, groups)
    assert any("fewer-than-2" in s["reason"] for s in result["skipped"])


def test_vertical_axis_distributes_top(deck_with_3_cards):
    """Same logic but on the vertical axis."""
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    s1 = slide.shapes.add_textbox(Emu(1000000), Emu(500000), Emu(2000000), Emu(800000))
    s2 = slide.shapes.add_textbox(Emu(1000000), Emu(1500000), Emu(2000000), Emu(800000))
    s3 = slide.shapes.add_textbox(Emu(1000000), Emu(3000000), Emu(2000000), Emu(800000))   # uneven gap
    sids = [s1.shape_id, s2.shape_id, s3.shape_id]

    groups = {
        "groups": [{
            "name": "col", "shape_ids": sids, "axis": "vertical",
            "uniform_spacing": True,
        }],
    }
    repair_peers_smart(prs, groups)
    shapes = sorted(
        [s for s in slide.shapes if s.shape_id in sids],
        key=lambda s: int(s.top),
    )
    gap1 = int(shapes[1].top) - (int(shapes[0].top) + int(shapes[0].height))
    gap2 = int(shapes[2].top) - (int(shapes[1].top) + int(shapes[1].height))
    assert gap1 == gap2
