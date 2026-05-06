"""End-to-end tests for the mutate.py CLI on a synthesized fixture deck.

We build a tiny deck with three text boxes, run a chain of mutate
subcommands, and re-open the result with python-pptx to assert the
expected XML state. This is the closest thing to a contract test the L4
orchestrator and OpenCode-driven model code see.
"""
from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


SKILL = Path(__file__).resolve().parents[1]
MUTATE = SKILL / "scripts" / "mutate.py"


@pytest.fixture
def fixture_deck(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    a.text_frame.text = "A"
    b = slide.shapes.add_textbox(Inches(3.0), Inches(0.7), Inches(2), Inches(0.5))
    b.text_frame.text = "B"
    c = slide.shapes.add_textbox(Inches(5.5), Inches(0.55), Inches(2), Inches(0.5))
    c.text_frame.text = "C"
    out = tmp_path / "fixture.pptx"
    prs.save(str(out))
    return out


def _run(out_path: Path, *argv: str) -> dict:
    result = subprocess.run(
        [sys.executable, str(MUTATE), *argv, "--out", str(out_path)],
        check=True,
        capture_output=True,
        text=True,
    )
    return json.loads(result.stdout.strip().splitlines()[-1])


def _shape_by_text(prs: Presentation, text: str):
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text == text:
            return shape
    raise AssertionError(f"shape with text {text!r} not found")


def test_list_ops_json(tmp_path: Path):
    result = subprocess.run(
        [sys.executable, str(MUTATE), "list-ops", "--json"],
        check=True,
        capture_output=True,
        text=True,
    )
    payload = json.loads(result.stdout)
    assert "ops" in payload
    assert len(payload["ops"]) >= 30


def test_set_fill(fixture_deck: Path, tmp_path: Path):
    prs0 = Presentation(str(fixture_deck))
    sid = prs0.slides[0].shapes[0].shape_id
    out = tmp_path / "v1.pptx"
    payload = _run(out, "set-fill", "--in", str(fixture_deck), "--shape-id", str(sid), "--color", "0F62FE")
    assert payload["op"] == "set-fill"
    prs = Presentation(str(out))
    fill = prs.slides[0].shapes[0].fill
    assert str(fill.fore_color.rgb) == "0F62FE"


def test_set_font_size(fixture_deck: Path, tmp_path: Path):
    prs0 = Presentation(str(fixture_deck))
    sid = prs0.slides[0].shapes[0].shape_id
    out = tmp_path / "v1.pptx"
    _run(out, "set-font-size", "--in", str(fixture_deck), "--shape-id", str(sid), "--size-pt", "28")
    prs = Presentation(str(out))
    para = prs.slides[0].shapes[0].text_frame.paragraphs[0]
    sizes = [r.font.size for r in para.runs]
    assert any(s and s.pt == 28 for s in sizes)


def test_align_left(fixture_deck: Path, tmp_path: Path):
    prs0 = Presentation(str(fixture_deck))
    ids = [s.shape_id for s in prs0.slides[0].shapes]
    out = tmp_path / "v1.pptx"
    _run(out, "align", "--in", str(fixture_deck), "--shape-ids", ",".join(map(str, ids)), "--edge", "left", "--target", "457200")
    prs = Presentation(str(out))
    lefts = [int(s.left) for s in prs.slides[0].shapes]
    assert all(left == 457200 for left in lefts)


def test_distribute_horizontal(fixture_deck: Path, tmp_path: Path):
    prs0 = Presentation(str(fixture_deck))
    ids = [s.shape_id for s in prs0.slides[0].shapes]
    out = tmp_path / "v1.pptx"
    _run(out, "distribute", "--in", str(fixture_deck), "--shape-ids", ",".join(map(str, ids)), "--axis", "horizontal")
    prs = Presentation(str(out))
    shapes = sorted(prs.slides[0].shapes, key=lambda s: int(s.left))
    gaps = [
        int(b.left) - (int(a.left) + int(a.width))
        for a, b in zip(shapes, shapes[1:])
    ]
    assert max(gaps) - min(gaps) <= 1  # equal within rounding


def test_z_order_to_back(fixture_deck: Path, tmp_path: Path):
    prs0 = Presentation(str(fixture_deck))
    target = prs0.slides[0].shapes[2]  # currently last
    sid = target.shape_id
    out = tmp_path / "v1.pptx"
    _run(out, "z-order", "--in", str(fixture_deck), "--shape-ids", str(sid), "--position", "back")
    prs = Presentation(str(out))
    first_shape = prs.slides[0].shapes[0]
    assert first_shape.shape_id == sid


def test_apply_typography_title(fixture_deck: Path, tmp_path: Path):
    prs0 = Presentation(str(fixture_deck))
    sid = prs0.slides[0].shapes[0].shape_id
    out = tmp_path / "v1.pptx"
    payload = _run(out, "apply-typography", "--in", str(fixture_deck), "--shape-id", str(sid), "--role", "title")
    assert payload["op"] == "apply-typography"
    prs = Presentation(str(out))
    para = prs.slides[0].shapes[0].text_frame.paragraphs[0]
    runs = list(para.runs)
    sizes = [r.font.size for r in runs if r.font.size is not None]
    assert any(s.pt == 28 for s in sizes)
    assert any(r.font.bold for r in runs)


def test_chain_two_mutations(fixture_deck: Path, tmp_path: Path):
    prs0 = Presentation(str(fixture_deck))
    sid = prs0.slides[0].shapes[0].shape_id
    out_a = tmp_path / "a.pptx"
    out_b = tmp_path / "b.pptx"
    _run(out_a, "set-fill", "--in", str(fixture_deck), "--shape-id", str(sid), "--color", "FF0000")
    _run(out_b, "set-line", "--in", str(out_a), "--shape-id", str(sid), "--color", "00FF00", "--width-pt", "1.5")
    prs = Presentation(str(out_b))
    s = prs.slides[0].shapes[0]
    assert str(s.fill.fore_color.rgb) == "FF0000"
    assert str(s.line.color.rgb) == "00FF00"
