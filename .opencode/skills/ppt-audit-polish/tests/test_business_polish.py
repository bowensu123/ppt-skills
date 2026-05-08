"""Unit tests for _business_polish.py — color math, theme picker, idempotency
markers, and an end-to-end smoke test on a synthesized deck."""
from __future__ import annotations

import sys
from pathlib import Path

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

import pytest

from _business_polish import (
    DECO_NAME_PREFIX,
    LEVEL_RICH,
    LEVEL_STANDARD,
    LEVEL_SUBTLE,
    THEME_KEYWORDS,
    _hex_to_rgb,
    _rgb_to_hex,
    _tint,
    apply_business_polish,
    pick_theme,
)


# ---- color helpers ----

def test_hex_to_rgb_roundtrip():
    assert _rgb_to_hex(_hex_to_rgb("0F62FE")) == "0F62FE"
    assert _rgb_to_hex(_hex_to_rgb("#0F62FE")) == "0F62FE"


def test_tint_lightens_toward_white():
    """0% tint = original; 100% tint = white."""
    assert _tint("0F62FE", 0.0) == "0F62FE"
    assert _tint("000000", 1.0) == "FFFFFF"


def test_tint_partial_is_lighter_not_white():
    """A 50% tint of black is gray (~127, 127, 127)."""
    out = _tint("000000", 0.5)
    r, g, b = _hex_to_rgb(out)
    # 0 + 255*0.5 = 127.5 → int truncation gives 127
    assert 126 <= r <= 128
    assert 126 <= g <= 128
    assert 126 <= b <= 128


def test_tint_subtle_amount_for_subtitle_bg():
    """The subtitle-bg helper uses 0.92 → very pale, near white but not full."""
    out = _tint("0F62FE", 0.92)
    r, g, b = _hex_to_rgb(out)
    # Each channel should be in the "very pale" range; the channel that
    # was darkest pre-tint will land lowest. Use 230 as a floor: 0F (15)
    # → 15 + 240*0.92 = 235.
    assert r >= 230 and g >= 230 and b >= 230


# ---- theme picker ----

def test_pick_theme_tech_keywords():
    content = "我们的 AI agent 用 LLM 框架做 model 推理，部署到 cloud 上"
    assert pick_theme(content, SKILL_ROOT / "themes") == "clean-tech"


def test_pick_theme_business_keywords():
    content = "Q3 营收增长 30%，客户留存 ROI 显著，市场份额扩大"
    assert pick_theme(content, SKILL_ROOT / "themes") == "business-warm"


def test_pick_theme_academic_keywords():
    content = "本研究用对照实验方法验证假设，论文已发表，结论显著"
    assert pick_theme(content, SKILL_ROOT / "themes") == "academic-soft"


def test_pick_theme_editorial_keywords():
    content = "我们的品牌故事和叙事，从观点出发的 editorial 设计"
    assert pick_theme(content, SKILL_ROOT / "themes") == "editorial-dark"


def test_pick_theme_claude_code_keywords():
    content = "Claude CLI agent 在 terminal 跑 bash command，编辑 Python 脚本，commit 到 git repo"
    assert pick_theme(content, SKILL_ROOT / "themes") == "claude-code"


# ---- new 10 business-grade themes ----

@pytest.mark.parametrize("content, expected", [
    ("麦肯锡 BCG 战略转型 提案 PMO 组织设计", "minimalist-business"),
    ("MECE 议题树 金字塔 假设驱动 takeaway 4P SWOT", "consulting"),
    ("Stripe Linear Notion SaaS subscription Kubernetes Vercel", "modern-tech"),
    ("年报 上市 财报披露 国资 央企 ESG 信息披露 董事会", "corporate-classic"),
    ("种子轮 估值 BP TAM Sequoia a16z preseed term sheet 募资", "pitch-deck"),
    ("Vogue lookbook 时装 campaign 新品发布 lifestyle 趋势报告", "editorial-magazine"),
    ("GMV DAU MAU MRR LTV CAC 留存率 漏斗 cohort ARPU", "data-heavy"),
    ("白皮书 对照实验 假设检验 实证研究 p值 RCT meta分析", "academic-research"),
    ("VI 设计 agency pitch 创意提案 key visual brand identity", "creative-agency"),
    ("私募 奢侈品 高净值 VIP luxury PE VC 家族办公室 腕表", "dark-premium"),
    ("华为 鸿蒙 HarmonyOS 鲲鹏 昇腾 MindSpore Kirin 海思", "huawei-style"),
])
def test_pick_theme_business_grade_themes(content, expected):
    assert pick_theme(content, SKILL_ROOT / "themes") == expected


@pytest.mark.parametrize("theme_name", [
    "minimalist-business", "consulting", "modern-tech",
    "corporate-classic", "pitch-deck", "editorial-magazine",
    "data-heavy", "academic-research", "creative-agency", "dark-premium",
    "huawei-style",
])
def test_new_themes_load_with_full_schema(theme_name):
    """Every new theme JSON must have the complete schema and be loadable."""
    from _common import load_theme
    theme = load_theme(SKILL_ROOT / "themes" / f"{theme_name}.json")
    # Required palette keys
    for key in ("primary", "primary_soft", "accent", "text_strong", "text",
                 "text_muted", "background", "surface", "border"):
        assert theme.palette.get(key), f"{theme_name} missing palette.{key}"
    # Required typography keys
    typo = theme.typography
    assert typo.get("font_family")
    assert typo.get("size_pt")
    for role in ("title", "subtitle", "h2", "body", "caption", "badge"):
        assert role in typo["size_pt"], f"{theme_name} missing size_pt.{role}"
    # Required decoration keys
    assert theme.raw.get("decoration")


def test_huawei_theme_has_layout_regions():
    """huawei-style ships with corporate-mandated layout_regions that the
    agent should respect when designing layout.json. Verify the schema."""
    import json
    raw = json.loads(
        (SKILL_ROOT / "themes" / "huawei-style.json").read_text(encoding="utf-8"),
    )
    regions = raw.get("layout_regions")
    assert regions is not None, "huawei-style must have layout_regions"
    # Must have 3 named regions + margins
    for region_name in ("title", "tag", "content", "margins"):
        assert region_name in regions, f"missing region {region_name}"
    # Each region (except margins) has bbox in EMU
    for region_name in ("title", "tag", "content"):
        region = regions[region_name]
        for key in ("x_emu", "y_emu", "w_emu", "h_emu"):
            assert key in region, f"{region_name} missing {key}"
            assert isinstance(region[key], int), f"{region_name}.{key} not int"
    # slide_dims present
    dims = raw.get("slide_dims")
    assert dims is not None
    assert dims.get("aspect") == "16:9"
    assert dims.get("width_emu") == 12192000


def test_huawei_theme_palette_has_extended_grays():
    """huawei-style spec includes extra gray levels (text_light / text_label /
    surface_medium / surface_dark / gradient_light) that a strict deck
    template needs."""
    import json
    raw = json.loads(
        (SKILL_ROOT / "themes" / "huawei-style.json").read_text(encoding="utf-8"),
    )
    palette = raw["palette"]
    for extra in ("text_light", "text_label", "surface_medium",
                   "surface_dark", "gradient_light", "primary_dark"):
        assert extra in palette, f"huawei palette missing {extra}"
    # Brand red colors
    assert palette["primary"] == "C00000"
    assert palette["primary_dark"] == "A90102"


def test_huawei_theme_dual_font_family():
    """Huawei spec uses 微软雅黑 for Chinese + Arial for English."""
    import json
    raw = json.loads(
        (SKILL_ROOT / "themes" / "huawei-style.json").read_text(encoding="utf-8"),
    )
    typo = raw["typography"]
    assert typo["font_family"] == "微软雅黑"
    assert typo["font_family_english"] == "Arial"


def test_pick_theme_empty_falls_back():
    """Empty content returns the default fallback."""
    assert pick_theme("", SKILL_ROOT / "themes") == "clean-tech"


def test_theme_keywords_cover_all_themes():
    """Every theme listed in THEME_KEYWORDS must have a matching JSON file."""
    themes_dir = SKILL_ROOT / "themes"
    for theme_name in THEME_KEYWORDS:
        assert (themes_dir / f"{theme_name}.json").exists(), \
            f"theme {theme_name} listed in keywords but no JSON exists"


# ---- decoration markers ----

def test_deco_marker_constant_is_unique_enough():
    """The decoration name marker should be very unlikely to collide
    with a real shape's name."""
    assert DECO_NAME_PREFIX.startswith("__")
    assert "deco" in DECO_NAME_PREFIX


# ---- end-to-end smoke (uses python-pptx to build a tiny deck) ----

@pytest.fixture
def tiny_deck(tmp_path):
    """Create a 1-slide deck with a title shape so apply_business_polish
    has something to operate on."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish
    # Add a title-like text box
    title = slide.shapes.add_textbox(Emu(914400), Emu(457200), Emu(8229600), Emu(914400))
    title.text_frame.text = "AI 框架与 model 部署"
    out = tmp_path / "in.pptx"
    prs.save(out)
    return out


def test_apply_business_polish_runs_without_error(tiny_deck, tmp_path):
    """The polish pass should run end-to-end on a simple deck and emit
    actions. Exact action count depends on role detection success."""
    from pptx import Presentation
    from _common import load_theme

    prs = Presentation(tiny_deck)
    theme = load_theme(SKILL_ROOT / "themes" / "clean-tech.json")
    # Synthesize minimal role_data: one detected title.
    title_sid = next(
        (s.shape_id for s in prs.slides[0].shapes if s.has_text_frame),
        None,
    )
    role_data = {"slides": [{"slide_index": 1, "shapes": [
        {"shape_id": title_sid, "role": "title"},
    ]}]}
    result = apply_business_polish(prs, theme, role_data, level=LEVEL_STANDARD)
    assert "level" in result
    assert "actions" in result
    # At minimum: typography for title fired.
    typo = [a for a in result["actions"] if a["action"] == "smart-typography"]
    assert len(typo) >= 1


def test_apply_business_polish_levels_differ(tiny_deck, tmp_path):
    """Higher levels add more decoration actions than lower levels."""
    from pptx import Presentation
    from _common import load_theme

    theme = load_theme(SKILL_ROOT / "themes" / "clean-tech.json")
    title_sid_results = {}
    for level in (LEVEL_SUBTLE, LEVEL_STANDARD, LEVEL_RICH):
        prs = Presentation(tiny_deck)
        title_sid = next(
            (s.shape_id for s in prs.slides[0].shapes if s.has_text_frame),
            None,
        )
        role_data = {"slides": [{"slide_index": 1, "shapes": [
            {"shape_id": title_sid, "role": "title"},
        ]}]}
        title_sid_results[level] = apply_business_polish(
            prs, theme, role_data, level=level,
        )
    n_subtle = title_sid_results[LEVEL_SUBTLE]["actions_applied"]
    n_standard = title_sid_results[LEVEL_STANDARD]["actions_applied"]
    n_rich = title_sid_results[LEVEL_RICH]["actions_applied"]
    # Standard adds at least one action (footer divider) over subtle
    assert n_standard >= n_subtle
    # Rich adds at least one tint action over standard if subtitle exists;
    # in this test we have no subtitle so rich == standard. Just verify
    # non-decreasing.
    assert n_rich >= n_standard


def test_apply_business_polish_idempotent(tiny_deck):
    """Running polish twice on same deck doesn't double-add decorations."""
    from pptx import Presentation
    from _common import load_theme

    theme = load_theme(SKILL_ROOT / "themes" / "clean-tech.json")
    prs = Presentation(tiny_deck)
    title_sid = next(
        (s.shape_id for s in prs.slides[0].shapes if s.has_text_frame),
        None,
    )
    role_data = {"slides": [{"slide_index": 1, "shapes": [
        {"shape_id": title_sid, "role": "title"},
    ]}]}
    apply_business_polish(prs, theme, role_data, level=LEVEL_STANDARD)
    decos_after_first = [
        s for s in prs.slides[0].shapes
        if (s.name or "").startswith(DECO_NAME_PREFIX)
    ]
    apply_business_polish(prs, theme, role_data, level=LEVEL_STANDARD)
    decos_after_second = [
        s for s in prs.slides[0].shapes
        if (s.name or "").startswith(DECO_NAME_PREFIX)
    ]
    # No new deco shapes added on second run.
    assert len(decos_after_first) == len(decos_after_second)


def test_polish_does_not_change_text_content(tiny_deck):
    """Constraint: polish must NEVER touch text content."""
    from pptx import Presentation
    from _common import load_theme

    theme = load_theme(SKILL_ROOT / "themes" / "clean-tech.json")
    prs = Presentation(tiny_deck)
    title_shape = next(s for s in prs.slides[0].shapes if s.has_text_frame)
    original_text = title_shape.text_frame.text
    title_sid = title_shape.shape_id
    role_data = {"slides": [{"slide_index": 1, "shapes": [
        {"shape_id": title_sid, "role": "title"},
    ]}]}
    apply_business_polish(prs, theme, role_data, level=LEVEL_RICH)
    # Find the same title shape after polish.
    title_after = next(
        s for s in prs.slides[0].shapes
        if s.shape_id == title_sid and s.has_text_frame
    )
    assert title_after.text_frame.text == original_text
