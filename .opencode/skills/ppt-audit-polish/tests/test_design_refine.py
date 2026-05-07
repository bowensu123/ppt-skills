"""Tests for _design_refine.py — background-aware text contrast + theme luminance pick."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from _design_refine import (
    WCAG_AA_NORMAL,
    _bbox_contains_center,
    _best_contrast_color,
    _find_effective_background,
    _is_large_text,
    detect_background_luminance,
    pick_theme_for_background,
    refine_text_contrast,
)


# ---- bbox containment ----

def test_bbox_contains_center_inside():
    outer = {"left": 100, "top": 100, "width": 1000, "height": 500}
    inner = {"left": 300, "top": 200, "width": 100, "height": 50}
    assert _bbox_contains_center(outer, inner)


def test_bbox_contains_center_outside():
    """Use big EMU values so the default 0.1" slack (91440) doesn't bridge them."""
    outer = {"left": 100000, "top": 100000, "width": 1000000, "height": 500000}
    inner = {"left": 5000000, "top": 5000000, "width": 100000, "height": 50000}
    assert not _bbox_contains_center(outer, inner)


# ---- background discovery ----

def test_find_effective_background_returns_innermost():
    text = {"shape_id": 99, "kind": "text",
            "left": 1100000, "top": 1100000, "width": 200000, "height": 100000,
            "text": "hi"}
    objects = [
        text,
        {"shape_id": 1, "kind": "container", "fill_hex": "0F0F0F",
         "left": 0, "top": 0, "width": 12000000, "height": 6000000},
        {"shape_id": 2, "kind": "container", "fill_hex": "FF0000",
         "left": 1000000, "top": 1000000, "width": 1000000, "height": 500000},
    ]
    bg = _find_effective_background(text, objects, slide_background="FFFFFF")
    assert bg == "FF0000"   # innermost = smallest containing


def test_find_effective_background_falls_back_to_slide():
    text = {"shape_id": 99, "kind": "text",
            "left": 100, "top": 100, "width": 100, "height": 50}
    bg = _find_effective_background(text, [text], slide_background="0F0F0F")
    assert bg == "0F0F0F"


# ---- best contrast color ----

def test_best_contrast_color_picks_white_on_dark():
    color, ratio = _best_contrast_color("0F0F0F", ["161616", "FFFFFF", "FF0000"])
    assert color == "FFFFFF"
    assert ratio > 15.0


def test_best_contrast_color_picks_dark_on_light():
    color, ratio = _best_contrast_color("FFFFFF", ["161616", "FFFFFF", "F5F5F5"])
    assert color == "161616"
    assert ratio > 10.0


# ---- large text classification ----

def test_is_large_text_yes():
    """font_sizes are EMU; 18pt = 18 * 12700 = 228600 EMU."""
    assert _is_large_text({"font_sizes": [228600]})


def test_is_large_text_no():
    assert not _is_large_text({"font_sizes": [127000]})  # 10pt


# ---- background luminance detection ----

def test_detect_background_luminance_dark():
    inspection = {"objects": [
        {"fill_hex": "0F0F0F", "width": 12000000, "height": 6000000,
         "anomalous": False},
    ]}
    lum = detect_background_luminance(inspection)
    assert lum < 0.05


def test_detect_background_luminance_light():
    inspection = {"objects": [
        {"fill_hex": "FFFFFF", "width": 12000000, "height": 6000000,
         "anomalous": False},
    ]}
    lum = detect_background_luminance(inspection)
    assert lum > 0.95


def test_detect_background_luminance_no_fills():
    inspection = {"objects": [{"width": 1000, "height": 1000, "anomalous": False}]}
    lum = detect_background_luminance(inspection)
    assert lum == 1.0   # safe default


# ---- theme picker by background ----

def test_theme_picker_dark_bg_overrides_keyword_pick():
    """A tech-keyword content on a dark background should NOT pick clean-tech
    (light theme) — should switch to claude-code."""
    content = "AI 框架 model agent cloud"
    theme = pick_theme_for_background(content, SKILL_ROOT / "themes",
                                       background_lum=0.05)
    assert theme == "claude-code"


def test_theme_picker_light_bg_keeps_keyword_pick():
    content = "AI 框架 model agent cloud"
    theme = pick_theme_for_background(content, SKILL_ROOT / "themes",
                                       background_lum=0.95)
    assert theme == "clean-tech"


def test_theme_picker_no_bg_lum_falls_back_to_keyword():
    content = "AI 框架 model"
    theme = pick_theme_for_background(content, SKILL_ROOT / "themes",
                                       background_lum=None)
    assert theme == "clean-tech"


# ---- end-to-end refine_text_contrast ----

def test_refine_text_contrast_swaps_low_contrast_text(tmp_path):
    """Red text on dark navy panel → swap to white (or text_strong)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Pt
    from _common import load_theme

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Dark panel covering most of the slide
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(914400), Emu(914400), Emu(8000000), Emu(4000000),
    )
    panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x0F)
    # Red text in the middle of the panel
    txt = slide.shapes.add_textbox(
        Emu(2000000), Emu(2000000), Emu(4000000), Emu(500000),
    )
    p = txt.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Important"
    r.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
    r.font.size = Pt(14)

    deck = tmp_path / "d.pptx"; prs.save(deck)
    # Build minimal inspection-like dict from the deck.
    text_sid = txt.shape_id
    panel_sid = panel.shape_id
    insp_slide = {"objects": [
        {"shape_id": panel_sid, "kind": "container",
         "left": int(panel.left), "top": int(panel.top),
         "width": int(panel.width), "height": int(panel.height),
         "fill_hex": "0F0F0F", "anomalous": False},
        {"shape_id": text_sid, "kind": "text",
         "left": int(txt.left), "top": int(txt.top),
         "width": int(txt.width), "height": int(txt.height),
         "text": "Important", "text_color": "CC0000",
         "font_sizes": [14 * 12700],   # 14pt in EMU
         "anomalous": False},
    ]}
    theme = load_theme(SKILL_ROOT / "themes" / "clean-tech.json")
    result = refine_text_contrast(prs, theme, [insp_slide])
    assert result["applied"] >= 1
    action = result["actions"][0]
    assert action["shape_id"] == text_sid
    # New color must give better contrast
    assert action["ratio_after"] >= 4.5


def test_refine_text_contrast_skips_passing_text(tmp_path):
    """White text on dark panel already passes → no action."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt
    from _common import load_theme

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txt = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(8000000), Emu(500000))
    p = txt.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Already fine"
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.size = Pt(14)
    deck = tmp_path / "d.pptx"; prs.save(deck)

    insp_slide = {"objects": [{
        "shape_id": txt.shape_id, "kind": "text",
        "left": 0, "top": 0, "width": 8000000, "height": 500000,
        "text": "Already fine", "text_color": "FFFFFF",
        "font_sizes": [14 * 12700], "anomalous": False,
    }], "__background_hex": "0F0F0F"}
    theme = load_theme(SKILL_ROOT / "themes" / "clean-tech.json")
    result = refine_text_contrast(prs, theme, [insp_slide])
    assert result["applied"] == 0
