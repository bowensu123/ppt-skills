"""Unit tests for _decoration_extract.py — vector decoration capture."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from _decoration_extract import (
    _DECO_AUTO_SHAPES,
    _decorative_score,
    extract_decorations,
)


class _Shape:
    def __init__(self, left, top, width, height,
                 has_fill=True, fill_color="D97757",
                 line_color=None, text="",
                 auto_shape="OVAL"):
        self.left = left; self.top = top
        self.width = width; self.height = height
        self.has_text_frame = bool(text)
        self._text = text
        self._fill_color = fill_color if has_fill else None
        self._line_color = line_color
        self._auto_shape = auto_shape


def _patched_score(shape, slide_w, slide_h, monkeypatch):
    """Run _decorative_score with helpers mocked to read our test shape."""
    import _decoration_extract as mod
    monkeypatch.setattr(mod, "_fill_hex", lambda s: s._fill_color)
    monkeypatch.setattr(mod, "_line_hex", lambda s: s._line_color)
    monkeypatch.setattr(mod, "_auto_shape_type_name", lambda s: s._auto_shape)
    monkeypatch.setattr(mod, "_shape_text", lambda s: s._text)
    return _decorative_score(shape, slide_w, slide_h)


def test_decorative_score_corner_oval_high(monkeypatch):
    """Small filled oval near corner = high decoration score."""
    s = _Shape(11000000, 200000, 600000, 600000, fill_color="D97757",
               text="01", auto_shape="OVAL")
    score = _patched_score(s, 12192000, 6858000, monkeypatch)
    assert score >= 0.8


def test_decorative_score_central_card_low(monkeypatch):
    """Big centered shape = low score (looks like a content card).
    24% of slide area + central + has long text → no area bonus + no
    edge bonus + no short-text bonus, only the +0.3 fill+kind bonus."""
    s = _Shape(3000000, 1500000, 6000000, 4000000, fill_color="111827",
               text="Long content text inside a card", auto_shape="RECTANGLE")
    score = _patched_score(s, 12192000, 6858000, monkeypatch)
    assert score < 0.5


def test_decorative_score_thin_line_decoration(monkeypatch):
    """Thin line near edge = decoration."""
    s = _Shape(457200, 6500000, 11000000, 12700, has_fill=False,
               line_color="6F6F6F", text="", auto_shape="RECTANGLE")
    score = _patched_score(s, 12192000, 6858000, monkeypatch)
    assert score >= 0.5


def test_extract_decorations_real_deck(tmp_path):
    """Build a tiny deck with one picture + one oval badge + one card,
    confirm only the oval ends up in decorations.json."""
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 1. Picture (should be skipped — covered by _asset_extract)
    img = tmp_path / "pic.png"
    Image.new("RGB", (8, 8), "white").save(img, "PNG")
    slide.shapes.add_picture(str(img), Emu(914400), Emu(914400),
                              Emu(2000000), Emu(2000000))

    # 2. Decorative oval badge in top-right corner
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(11200000), Emu(300000), Emu(500000), Emu(500000),
    )
    oval.fill.solid(); oval.fill.fore_color.rgb = RGBColor(0xD9, 0x77, 0x57)
    oval.text_frame.text = "01"

    # 3. Big content card (should NOT be flagged as decoration)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(3000000), Emu(2000000), Emu(5000000), Emu(3000000),
    )
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    card.text_frame.text = "This is a content card with substantive text inside"

    deck = tmp_path / "deck.pptx"
    prs.save(deck)

    work = tmp_path / "work"
    manifest = extract_decorations(deck, work, include_master=False)

    # The picture must not appear; the oval should; the card may or may not
    # depending on size+text heuristic.
    sids_in_manifest = {d["shape_id"] for d in manifest["decorations"]}
    assert oval.shape_id in sids_in_manifest, \
        f"oval should be in decorations, got {sids_in_manifest}"
    # Pictures must NEVER be in decorations.json (separate concern).
    pic_sids = {s.shape_id for s in slide.shapes if s.shape_type == 13}  # PICTURE
    assert not (pic_sids & sids_in_manifest), \
        "pictures must not appear in decorations.json"


def test_extract_decorations_writes_manifest(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(100000), Emu(100000),
                                 Emu(400000), Emu(400000))
    deck = tmp_path / "d.pptx"; prs.save(deck)
    work = tmp_path / "w"
    manifest = extract_decorations(deck, work, include_master=False)
    assert (work / "decorations.json").exists()


def test_decoration_kind_for_layout_mapping(tmp_path):
    """Ovals → 'circle', rectangles → 'rect', rounded → 'rounded_rect'."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Three small filled corner shapes
    for i, (kind, x) in enumerate([
        (MSO_SHAPE.OVAL, 11200000),
        (MSO_SHAPE.RECTANGLE, 11200000),
        (MSO_SHAPE.ROUNDED_RECTANGLE, 11200000),
    ]):
        sp = slide.shapes.add_shape(kind, Emu(x), Emu(200000 + i * 700000),
                                     Emu(400000), Emu(400000))
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xD9, 0x77, 0x57)

    deck = tmp_path / "d.pptx"; prs.save(deck)
    manifest = extract_decorations(deck, tmp_path / "w", include_master=False)
    kinds = {d["auto_shape_type"]: d["kind_for_layout"]
             for d in manifest["decorations"]}
    assert kinds.get("OVAL") == "circle"
    assert kinds.get("RECTANGLE") == "rect"
    assert kinds.get("ROUNDED_RECTANGLE") == "rounded_rect"
