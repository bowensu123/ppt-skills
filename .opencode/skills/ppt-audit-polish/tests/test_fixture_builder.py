from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
FIXTURE_SCRIPT = ROOT / "tests" / "fixtures" / "build_fixtures.py"


def _build_fixtures(tmp_path: Path) -> None:
    result = subprocess.run(
        [sys.executable, str(FIXTURE_SCRIPT), "--output-dir", str(tmp_path)],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr


def test_fixture_builder_generates_expected_pptx_files(tmp_path: Path) -> None:
    _build_fixtures(tmp_path)

    overflow_path = tmp_path / "overflow-case.pptx"
    alignment_path = tmp_path / "alignment-case.pptx"
    chart_risk_path = tmp_path / "chart-risk-case.pptx"

    assert overflow_path.exists()
    assert alignment_path.exists()
    assert chart_risk_path.exists()


def test_overflow_case_contains_shape_beyond_slide_width(tmp_path: Path) -> None:
    _build_fixtures(tmp_path)

    prs = Presentation(tmp_path / "overflow-case.pptx")
    slide = prs.slides[0]

    assert any(shape.left + shape.width > prs.slide_width for shape in slide.shapes)


def test_alignment_case_has_misaligned_cards_and_larger_font(tmp_path: Path) -> None:
    _build_fixtures(tmp_path)

    prs = Presentation(tmp_path / "alignment-case.pptx")
    slide = prs.slides[0]
    cards = [shape for shape in slide.shapes if getattr(shape, "text", "") in {"Card 1", "Card 2", "Card 3"}]

    assert len(cards) == 3
    assert len({card.top for card in cards}) > 1
    assert any(card.text == "Card 2" and card.text_frame.paragraphs[0].font.size == Pt(22) for card in cards)


def test_chart_risk_case_contains_a_chart_object(tmp_path: Path) -> None:
    _build_fixtures(tmp_path)

    prs = Presentation(tmp_path / "chart-risk-case.pptx")
    slide = prs.slides[0]

    assert any(shape.shape_type == MSO_SHAPE_TYPE.CHART for shape in slide.shapes)
