from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def _save_overflow_case(output_dir: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(7.3), Inches(0.4), Inches(3.2), Inches(0.8))
    frame = title.text_frame
    frame.text = "Overflow headline"
    frame.paragraphs[0].font.size = Pt(28)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.0), Inches(2.0))
    body.text_frame.text = "This slide deliberately pushes the title beyond the right edge."
    prs.save(output_dir / "overflow-case.pptx")


def _save_alignment_case(output_dir: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for index, left in enumerate((0.8, 3.4, 6.4), start=1):
        card = slide.shapes.add_textbox(Inches(left), Inches(1.5 + (0.12 * index)), Inches(2.0), Inches(1.2))
        frame = card.text_frame
        frame.text = f"Card {index}"
        frame.paragraphs[0].font.size = Pt(22 if index == 2 else 18)
    prs.save(output_dir / "alignment-case.pptx")


def _save_peer_gap_case(output_dir: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for index, (left, top, font_size) in enumerate(
        ((0.8, 1.5, 18), (3.1, 1.63, 22), (6.7, 1.43, 18)),
        start=1,
    ):
        card = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(1.8), Inches(1.0))
        frame = card.text_frame
        frame.text = f"Peer {index}"
        frame.paragraphs[0].font.size = Pt(font_size)
    prs.save(output_dir / "peer-gap-case.pptx")


def _save_chart_risk_case(output_dir: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Revenue", (5.2, 7.4, 6.8))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0),
        Inches(1.0),
        Inches(6.8),
        Inches(4.0),
        chart_data,
    )
    prs.save(output_dir / "chart-risk-case.pptx")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", required=True)
    args = parser.parse_args()

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    _save_overflow_case(output_dir)
    _save_alignment_case(output_dir)
    _save_peer_gap_case(output_dir)
    _save_chart_risk_case(output_dir)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
