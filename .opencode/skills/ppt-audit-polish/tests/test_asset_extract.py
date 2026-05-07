"""Unit tests for _asset_extract.py — deterministic image extraction."""
from __future__ import annotations

import sys
from pathlib import Path

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

import pytest

from _asset_extract import _decorative_hint, extract_assets


# ---- _decorative_hint ----

class _Shape:
    def __init__(self, left, top, width, height):
        self.left = left; self.top = top
        self.width = width; self.height = height


def test_decorative_hint_corner_logo():
    """Small shape near top-right corner = looks like a logo."""
    slide_w, slide_h = 12192000, 6858000
    # 800K wide × 800K tall, at top-right (10M, 200K)
    s = _Shape(10000000, 200000, 800000, 800000)
    assert _decorative_hint(s, slide_w, slide_h) is True


def test_decorative_hint_central_image_not_decorative():
    """Big shape in the middle = content, not chrome."""
    slide_w, slide_h = 12192000, 6858000
    s = _Shape(3000000, 1500000, 6000000, 4000000)
    assert _decorative_hint(s, slide_w, slide_h) is False


def test_decorative_hint_small_centered_not_decorative():
    """Small shape but in the middle = could be an item icon, not chrome."""
    slide_w, slide_h = 12192000, 6858000
    s = _Shape(5000000, 3000000, 500000, 500000)
    assert _decorative_hint(s, slide_w, slide_h) is False


def test_decorative_hint_invalid_geometry():
    s = _Shape(0, 0, 0, 0)
    assert _decorative_hint(s, 12192000, 6858000) is False


# ---- extract_assets end-to-end (synthesized deck with one image) ----

@pytest.fixture
def deck_with_picture(tmp_path):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu

    # Generate a valid 8x8 PNG via Pillow.
    png_path = tmp_path / "tiny.png"
    Image.new("RGB", (8, 8), "white").save(png_path, "PNG")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(str(png_path),
                             Emu(914400), Emu(914400),
                             Emu(2000000), Emu(2000000))
    out = tmp_path / "deck.pptx"
    prs.save(out)
    return out


def test_extract_assets_finds_picture(tmp_path, deck_with_picture):
    work_dir = tmp_path / "work"
    manifest = extract_assets(deck_with_picture, work_dir)
    assert manifest["extracted_count"] == 1
    assert (work_dir / "assets-manifest.json").exists()
    asset = manifest["assets"][0]
    assert asset["kind"] == "picture"
    assert asset["slide_index"] == 1
    assert (work_dir / asset["path"]).exists()


def test_extract_assets_writes_correct_binary(tmp_path, deck_with_picture):
    work_dir = tmp_path / "work"
    manifest = extract_assets(deck_with_picture, work_dir)
    asset = manifest["assets"][0]
    written = (work_dir / asset["path"]).read_bytes()
    # Round-tripped binary must start with PNG magic
    assert written.startswith(b"\x89PNG\r\n\x1a\n")
    assert asset["size_bytes"] == len(written)


def test_extract_assets_includes_geometry_and_zorder(tmp_path, deck_with_picture):
    work_dir = tmp_path / "work"
    manifest = extract_assets(deck_with_picture, work_dir)
    asset = manifest["assets"][0]
    # bbox we set: left=914400, top=914400, w=2000000, h=2000000
    assert asset["bbox_emu"] == [914400, 914400, 2000000, 2000000]
    assert asset["z_index"] >= 0


def test_extract_assets_empty_deck(tmp_path):
    """Deck with no pictures → manifest with extracted_count 0."""
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    deck = tmp_path / "empty.pptx"
    prs.save(deck)
    work_dir = tmp_path / "w"
    manifest = extract_assets(deck, work_dir)
    assert manifest["extracted_count"] == 0
    assert manifest["assets"] == []


def test_extract_assets_creates_assets_dir(tmp_path, deck_with_picture):
    """assets/ folder is created even if work_dir doesn't exist."""
    work_dir = tmp_path / "fresh-work"
    extract_assets(deck_with_picture, work_dir)
    assert (work_dir / "assets").is_dir()
