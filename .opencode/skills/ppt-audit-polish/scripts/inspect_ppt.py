from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


_TEXT_KIND_CONTAINER = "container"
_TEXT_KIND_TEXT = "text"
_TEXT_KIND_PICTURE = "picture"
_TEXT_KIND_CHART = "chart"
_TEXT_KIND_TABLE = "table"
_TEXT_KIND_CONNECTOR = "connector"
_TEXT_KIND_GROUP = "group"
_TEXT_KIND_OTHER = "shape"


def _shape_kind(shape) -> str:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.LINE:
        return _TEXT_KIND_CONNECTOR
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return _TEXT_KIND_GROUP
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return _TEXT_KIND_PICTURE
    if shape_type == MSO_SHAPE_TYPE.CHART:
        return _TEXT_KIND_CHART
    if shape_type == MSO_SHAPE_TYPE.TABLE:
        return _TEXT_KIND_TABLE
    has_tf = getattr(shape, "has_text_frame", False)
    if has_tf:
        text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
        return _TEXT_KIND_TEXT if text else _TEXT_KIND_CONTAINER
    return _TEXT_KIND_OTHER


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def _font_sizes(shape) -> list[int]:
    sizes: list[int] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        paragraph_sizes: list[int] = []
        for run in paragraph.runs:
            if run.font.size is not None:
                paragraph_sizes.append(int(run.font.size))
        if paragraph_sizes:
            sizes.extend(paragraph_sizes)
        elif paragraph.font.size is not None:
            sizes.append(int(paragraph.font.size))
    return sizes


def _font_families(shape) -> list[str]:
    """Distinct declared font family names across all runs in this shape.

    Used by the SVG-based font-fallback detector: compares declared
    families against what the renderer actually used.
    """
    families: list[str] = []
    seen: set[str] = set()
    if not getattr(shape, "has_text_frame", False):
        return families
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            name = getattr(run.font, "name", None)
            if name and name not in seen:
                seen.add(name)
                families.append(name)
        para_name = getattr(paragraph.font, "name", None)
        if para_name and para_name not in seen:
            seen.add(para_name)
            families.append(para_name)
    return families


def _text_runs(shape) -> list[dict]:
    """Per-run text content + formatting (font/size/bold/italic/color).

    Captures everything needed to faithfully reconstruct mixed-format
    text in a regenerated layout. Each run entry:
      {text, font_family, size_pt, bold, italic, color_hex,
       paragraph_index, run_index}
    """
    runs: list[dict] = []
    if not getattr(shape, "has_text_frame", False):
        return runs
    try:
        from pptx.util import Pt
    except ImportError:
        return runs
    for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
        for r_idx, run in enumerate(paragraph.runs):
            try:
                text = run.text or ""
            except (AttributeError, ValueError):
                continue
            if not text:
                continue
            font = run.font
            size_pt: float | None = None
            try:
                if font.size is not None:
                    size_pt = float(font.size) / Pt(1)
                elif paragraph.font.size is not None:
                    size_pt = float(paragraph.font.size) / Pt(1)
            except (AttributeError, ValueError):
                pass
            color_hex: str | None = None
            try:
                if font.color is not None and font.color.type is not None \
                        and font.color.rgb is not None:
                    color_hex = str(font.color.rgb)
            except (AttributeError, ValueError, KeyError):
                pass
            # Hyperlink (run.hyperlink.address)
            hyperlink: str | None = None
            try:
                hl = run.hyperlink
                if hl is not None and hl.address:
                    hyperlink = str(hl.address)
            except (AttributeError, ValueError, KeyError):
                pass
            runs.append({
                "text": text,
                "font_family": getattr(font, "name", None),
                "size_pt": size_pt,
                "bold": bool(getattr(font, "bold", False)) if font.bold is not None else None,
                "italic": bool(getattr(font, "italic", False)) if font.italic is not None else None,
                "color_hex": color_hex,
                "hyperlink": hyperlink,
                "paragraph_index": p_idx,
                "run_index": r_idx,
            })
    return runs


def _fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        fore = fill.fore_color
        if fore.type is None:
            return None
        rgb = fore.rgb
        if rgb is None:
            return None
        return str(rgb)
    except (AttributeError, ValueError, KeyError, TypeError):
        return None


def _line_hex(shape) -> str | None:
    try:
        line = shape.line
        color = line.color
        if color.type is None:
            return None
        rgb = color.rgb
        if rgb is None:
            return None
        return str(rgb)
    except (AttributeError, ValueError, KeyError, TypeError):
        return None


def _normalize_geometry(left: int, top: int, width: int, height: int) -> tuple[int, int, int, int, bool]:
    flipped = width < 0 or height < 0
    if width < 0:
        left = left + width
        width = -width
    if height < 0:
        top = top + height
        height = -height
    return left, top, width, height, flipped


def _text_first_color(shape) -> str | None:
    """Return the hex color of the first run with an explicit color, if any."""
    if not getattr(shape, "has_text_frame", False):
        return None
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                color = run.font.color
                if color.type is not None and color.rgb is not None:
                    return str(color.rgb)
        for paragraph in shape.text_frame.paragraphs:
            color = paragraph.font.color
            if color.type is not None and color.rgb is not None:
                return str(color.rgb)
    except (AttributeError, ValueError, KeyError, TypeError):
        return None
    return None


_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _picture_crop(shape) -> dict | None:
    """For pictures, return per-100000 src crop fractions {l, t, r, b}."""
    try:
        el = shape._element
        for child in el.iter():
            if child.tag.endswith("}srcRect"):
                return {
                    "l": int(child.get("l", "0") or 0),
                    "t": int(child.get("t", "0") or 0),
                    "r": int(child.get("r", "0") or 0),
                    "b": int(child.get("b", "0") or 0),
                }
    except (AttributeError, ValueError):
        pass
    return None


def _table_info(shape) -> dict | None:
    if not getattr(shape, "has_table", False):
        return None
    try:
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        empty = 0
        for r in table.rows:
            for c in r.cells:
                if not (c.text or "").strip():
                    empty += 1
        return {"rows": rows, "cols": cols, "empty_cells": empty}
    except (AttributeError, ValueError):
        return None


def _chart_info(shape) -> dict | None:
    if not getattr(shape, "has_chart", False):
        return None
    try:
        chart = shape.chart
        return {
            "chart_type": str(chart.chart_type),
            "series_count": len(list(chart.series)),
            "has_legend": chart.has_legend,
        }
    except (AttributeError, ValueError):
        return None


def _emit_object(slide_index: int, object_index: int, shape) -> dict:
    raw_left = int(shape.left) if shape.left is not None else 0
    raw_top = int(shape.top) if shape.top is not None else 0
    raw_width = int(shape.width) if shape.width is not None else 0
    raw_height = int(shape.height) if shape.height is not None else 0
    left, top, width, height, flipped = _normalize_geometry(raw_left, raw_top, raw_width, raw_height)

    kind = _shape_kind(shape)
    text = _shape_text(shape)
    font_sizes = _font_sizes(shape)
    font_families = _font_families(shape)
    text_runs = _text_runs(shape)
    is_anomalous = width <= 0 or height <= 0

    out = {
        "object_index": object_index,
        "shape_id": getattr(shape, "shape_id", object_index),
        "name": getattr(shape, "name", f"Shape {object_index}"),
        "kind": kind,
        "shape_type": int(shape.shape_type) if getattr(shape, "shape_type", None) is not None else None,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "raw_geometry": {"left": raw_left, "top": raw_top, "width": raw_width, "height": raw_height},
        "geometry_flipped": flipped,
        "anomalous": is_anomalous,
        "rotation": float(getattr(shape, "rotation", 0.0) or 0.0),
        "text": text,
        "font_sizes": font_sizes,
        "font_families": font_families,
        "text_runs": text_runs,
        "fill_hex": _fill_hex(shape),
        "line_hex": _line_hex(shape),
        "text_color": _text_first_color(shape),
    }
    if kind == "picture":
        crop = _picture_crop(shape)
        if crop:
            out["crop"] = crop
    elif kind == "table":
        ti = _table_info(shape)
        if ti:
            out["table_info"] = ti
    elif kind == "chart":
        ci = _chart_info(shape)
        if ci:
            out["chart_info"] = ci
    return out


def inspect_presentation(input_path: Path) -> dict:
    prs = Presentation(str(input_path))
    slides: list[dict] = []

    # Lazy imports so tests / minimal envs without these helpers still parse.
    try:
        from _master_inherit import collect_inheritance_info
    except ImportError:
        collect_inheritance_info = None  # type: ignore
    try:
        from _group_recurse import walk_group_children
    except ImportError:
        walk_group_children = None  # type: ignore

    for slide_index, slide in enumerate(prs.slides, start=1):
        objects: list[dict] = []
        group_children: list[dict] = []
        for object_index, shape in enumerate(slide.shapes, start=1):
            obj = _emit_object(slide_index, object_index, shape)
            objects.append(obj)
            # If this is a group, also collect transformed children for opt-in detectors.
            if obj.get("kind") == "group" and walk_group_children:
                group_children.extend(walk_group_children(shape))

        slide_payload = {
            "slide_index": slide_index,
            "width_emu": int(prs.slide_width),
            "height_emu": int(prs.slide_height),
            "objects": objects,
            "group_children": group_children,
        }
        if collect_inheritance_info:
            slide_payload["inheritance"] = collect_inheritance_info(slide)

        slides.append(slide_payload)

    return {"input": str(input_path), "slides": slides}


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    payload = inspect_presentation(Path(args.input))
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
