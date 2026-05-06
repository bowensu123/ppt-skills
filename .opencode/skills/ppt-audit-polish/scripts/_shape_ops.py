"""Low-level shape mutation primitives.

Every operation in this module:
  * takes a python-pptx shape (and sometimes auxiliary args)
  * mutates only that shape (no slide-wide side effects)
  * returns a small dict describing exactly what happened (or None if no-op)

The L2 ``mutate`` CLI binds CLI arguments to these primitives. The L4
orchestrator scores deck snapshots before/after and accepts/rejects
based on the result. Every primitive is deliberately small and pure-ish so
those higher layers stay simple.
"""
from __future__ import annotations

from copy import deepcopy
from typing import Any

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt

from _common import hex_to_rgb


_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


# =================== Geometry ===================

def set_position(shape, left: int | None = None, top: int | None = None) -> dict | None:
    changes: dict[str, Any] = {}
    if left is not None and int(shape.left) != left:
        changes["left"] = {"from": int(shape.left), "to": left}
        shape.left = Emu(left)
    if top is not None and int(shape.top) != top:
        changes["top"] = {"from": int(shape.top), "to": top}
        shape.top = Emu(top)
    return changes or None


def set_size(shape, width: int | None = None, height: int | None = None) -> dict | None:
    changes: dict[str, Any] = {}
    if width is not None and int(shape.width) != width:
        changes["width"] = {"from": int(shape.width), "to": width}
        shape.width = Emu(width)
    if height is not None and int(shape.height) != height:
        changes["height"] = {"from": int(shape.height), "to": height}
        shape.height = Emu(height)
    return changes or None


def nudge(shape, dx_emu: int = 0, dy_emu: int = 0) -> dict | None:
    if dx_emu == 0 and dy_emu == 0:
        return None
    new_left = int(shape.left) + dx_emu
    new_top = int(shape.top) + dy_emu
    return set_position(shape, left=new_left, top=new_top)


def set_rotation(shape, degrees: float) -> dict | None:
    current = float(getattr(shape, "rotation", 0.0) or 0.0)
    if abs(current - degrees) < 0.01:
        return None
    shape.rotation = degrees
    return {"rotation": {"from": current, "to": degrees}}


# =================== Fill ===================

def set_solid_fill(shape, hex_color: str) -> dict | None:
    try:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(hex_color)
        return {"fill": {"to": hex_color, "kind": "solid"}}
    except (AttributeError, ValueError, KeyError, TypeError):
        return None


def clear_fill(shape) -> dict | None:
    try:
        shape.fill.background()
        return {"fill": {"to": "none"}}
    except (AttributeError, ValueError, KeyError, TypeError):
        return None


def set_opacity(shape, alpha: float) -> dict | None:
    """Set fill alpha 0.0-1.0 by injecting <a:alpha> on solidFill."""
    if not (0.0 <= alpha <= 1.0):
        raise ValueError("alpha must be in [0, 1]")
    sp = shape._element
    nsmap = {"a": _NS_A}
    solidFill = sp.find(".//a:solidFill", nsmap)
    if solidFill is None:
        return None
    srgb = solidFill.find("a:srgbClr", nsmap)
    if srgb is None:
        return None
    for existing in srgb.findall(f"{{{_NS_A}}}alpha"):
        srgb.remove(existing)
    alpha_el = etree.SubElement(srgb, f"{{{_NS_A}}}alpha")
    alpha_el.set("val", str(int(alpha * 100000)))
    return {"opacity": {"to": alpha}}


# =================== Line ===================

def set_line(shape, hex_color: str | None = None, width_pt: float | None = None) -> dict | None:
    changes: dict[str, Any] = {}
    try:
        line = shape.line
        if hex_color is not None:
            line.color.rgb = hex_to_rgb(hex_color)
            changes["line_color"] = hex_color
        if width_pt is not None:
            line.width = Pt(width_pt)
            changes["line_width_pt"] = width_pt
    except (AttributeError, ValueError, KeyError, TypeError):
        return None
    return changes or None


def clear_line(shape) -> dict | None:
    """Remove the outline by setting it to noFill."""
    sp = shape._element
    nsmap = {"a": _NS_A}
    spPr = sp.find(".//p:spPr", {"p": _NS_P}) or sp.find(".//a:spPr", nsmap)
    if spPr is None:
        # Fallback: every shape has an spPr; search any namespace.
        for child in sp.iter():
            if child.tag.endswith("}spPr"):
                spPr = child
                break
    if spPr is None:
        return None
    for ln in spPr.findall(f"{{{_NS_A}}}ln"):
        spPr.remove(ln)
    ln = etree.SubElement(spPr, f"{{{_NS_A}}}ln")
    etree.SubElement(ln, f"{{{_NS_A}}}noFill")
    return {"line": {"to": "none"}}


# =================== Z-order ===================

def move_z_order(shape, position: str) -> dict | None:
    """position: 'back' | 'front' | 'up' | 'down'."""
    sp = shape._element
    parent = sp.getparent()
    if parent is None:
        return None
    siblings = list(parent)
    cur = siblings.index(sp)
    parent.remove(sp)

    if position == "back":
        # Find first non-metadata child to insert before.
        target = 0
        for idx, child in enumerate(parent):
            tag = etree.QName(child.tag).localname
            if tag not in ("nvGrpSpPr", "grpSpPr"):
                target = idx
                break
        else:
            target = len(parent)
        parent.insert(target, sp)
    elif position == "front":
        parent.append(sp)
    elif position == "up":
        new_idx = min(cur + 1, len(parent))
        parent.insert(new_idx, sp)
    elif position == "down":
        new_idx = max(cur - 1, 0)
        parent.insert(new_idx, sp)
    else:
        # Restore and bail.
        parent.insert(cur, sp)
        raise ValueError(f"position must be back|front|up|down, got {position!r}")
    return {"z_order": {"to": position}}


# =================== Shadow ===================

def set_shadow(shape, hex_color: str = "888888", blur_pt: float = 4.0, dist_pt: float = 2.0, alpha: float = 0.3) -> dict | None:
    sp = shape._element
    spPr = None
    for child in sp.iter():
        if child.tag.endswith("}spPr"):
            spPr = child
            break
    if spPr is None:
        return None
    nsmap = {"a": _NS_A}
    for existing in spPr.findall(f"{{{_NS_A}}}effectLst"):
        spPr.remove(existing)
    eff = etree.SubElement(spPr, f"{{{_NS_A}}}effectLst")
    outer = etree.SubElement(eff, f"{{{_NS_A}}}outerShdw")
    outer.set("blurRad", str(int(blur_pt * 12700)))
    outer.set("dist", str(int(dist_pt * 12700)))
    outer.set("dir", "2700000")  # 45 degrees
    outer.set("rotWithShape", "0")
    srgb = etree.SubElement(outer, f"{{{_NS_A}}}srgbClr")
    srgb.set("val", hex_color.lstrip("#"))
    alpha_el = etree.SubElement(srgb, f"{{{_NS_A}}}alpha")
    alpha_el.set("val", str(int(alpha * 100000)))
    return {"shadow": {"color": hex_color, "blur_pt": blur_pt, "dist_pt": dist_pt, "alpha": alpha}}


def clear_shadow(shape) -> dict | None:
    sp = shape._element
    removed = False
    for child in sp.iter():
        if child.tag.endswith("}effectLst"):
            child.getparent().remove(child)
            removed = True
    return {"shadow": {"to": "none"}} if removed else None


# =================== Corner radius (for ROUNDED_RECTANGLE) ===================

def set_corner_radius(shape, ratio: float) -> dict | None:
    """ratio in [0, 0.5]; only meaningful for shapes with adj1 (rounded rect)."""
    if not (0.0 <= ratio <= 0.5):
        raise ValueError("ratio must be in [0, 0.5]")
    sp = shape._element
    nsmap = {"a": _NS_A}
    prstGeom = None
    for child in sp.iter():
        if child.tag.endswith("}prstGeom"):
            prstGeom = child
            break
    if prstGeom is None:
        return None
    for av in prstGeom.findall(f"{{{_NS_A}}}avLst"):
        prstGeom.remove(av)
    av = etree.SubElement(prstGeom, f"{{{_NS_A}}}avLst")
    gd = etree.SubElement(av, f"{{{_NS_A}}}gd")
    gd.set("name", "adj")
    gd.set("fmla", f"val {int(ratio * 100000)}")
    return {"corner_radius_ratio": ratio}


# =================== Text ===================

def _iter_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        runs = list(paragraph.runs)
        if runs:
            for run in runs:
                yield "run", run, paragraph
        else:
            yield "paragraph", paragraph.font, paragraph


def set_font_family(shape, family: str) -> dict | None:
    n = 0
    for kind, target, _para in _iter_runs(shape):
        if kind == "run":
            target.font.name = family
        else:
            target.name = family
        n += 1
    return {"font_family": family, "run_count": n} if n > 0 else None


def set_font_size(shape, size_pt: float) -> dict | None:
    n = 0
    for kind, target, _para in _iter_runs(shape):
        if kind == "run":
            target.font.size = Pt(size_pt)
        else:
            target.size = Pt(size_pt)
        n += 1
    return {"font_size_pt": size_pt, "run_count": n} if n > 0 else None


def set_font_bold(shape, bold: bool) -> dict | None:
    n = 0
    for kind, target, _para in _iter_runs(shape):
        if kind == "run":
            target.font.bold = bold
        else:
            target.bold = bold
        n += 1
    return {"font_bold": bold, "run_count": n} if n > 0 else None


def set_font_italic(shape, italic: bool) -> dict | None:
    n = 0
    for kind, target, _para in _iter_runs(shape):
        if kind == "run":
            target.font.italic = italic
        else:
            target.italic = italic
        n += 1
    return {"font_italic": italic, "run_count": n} if n > 0 else None


def set_font_color(shape, hex_color: str) -> dict | None:
    rgb = hex_to_rgb(hex_color)
    n = 0
    for kind, target, _para in _iter_runs(shape):
        try:
            if kind == "run":
                target.font.color.rgb = rgb
            else:
                target.color.rgb = rgb
            n += 1
        except (AttributeError, ValueError):
            continue
    return {"font_color": hex_color, "run_count": n} if n > 0 else None


def set_text_align(shape, align: str) -> dict | None:
    """align: left|center|right|justify."""
    if not getattr(shape, "has_text_frame", False):
        return None
    from pptx.enum.text import PP_ALIGN

    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    target = mapping.get(align)
    if target is None:
        raise ValueError(f"align must be left|center|right|justify, got {align!r}")
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = target
    return {"text_align": align}


def set_text_v_align(shape, anchor: str) -> dict | None:
    """anchor: top|middle|bottom."""
    if not getattr(shape, "has_text_frame", False):
        return None
    from pptx.enum.text import MSO_ANCHOR

    mapping = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    target = mapping.get(anchor)
    if target is None:
        raise ValueError(f"anchor must be top|middle|bottom, got {anchor!r}")
    shape.text_frame.vertical_anchor = target
    return {"text_v_align": anchor}


def set_text_margin(shape, left_emu: int | None = None, right_emu: int | None = None,
                    top_emu: int | None = None, bottom_emu: int | None = None) -> dict | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    changes: dict[str, Any] = {}
    if left_emu is not None:
        tf.margin_left = Emu(left_emu)
        changes["margin_left_emu"] = left_emu
    if right_emu is not None:
        tf.margin_right = Emu(right_emu)
        changes["margin_right_emu"] = right_emu
    if top_emu is not None:
        tf.margin_top = Emu(top_emu)
        changes["margin_top_emu"] = top_emu
    if bottom_emu is not None:
        tf.margin_bottom = Emu(bottom_emu)
        changes["margin_bottom_emu"] = bottom_emu
    return changes or None


def set_line_spacing(shape, ratio: float) -> dict | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        paragraph.line_spacing = ratio
    return {"line_spacing": ratio}


def set_text_content(shape, text: str) -> dict | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    old = shape.text_frame.text
    if old == text:
        return None
    shape.text_frame.text = text
    return {"text": {"from": old[:60], "to": text[:60]}}


# =================== Slide-level shape add/remove ===================

def delete_shape(shape) -> dict:
    sp = shape._element
    parent = sp.getparent()
    if parent is not None:
        parent.remove(sp)
    return {"deleted": True}


def duplicate_shape(slide, shape) -> tuple[Any, dict]:
    """Append a deep-copy of shape to the slide; return (new_shape, info)."""
    sp = shape._element
    new_el = deepcopy(sp)
    spTree = slide.shapes._spTree
    spTree.append(new_el)
    # python-pptx wraps via factory; reuse by name lookup.
    new_shape = slide.shapes[-1]
    return new_shape, {"duplicated_from_shape_id": getattr(shape, "shape_id", None)}


# =================== Helpers ===================

def is_connector_like(shape) -> bool:
    """Return True if the shape behaves as a connector line.

    Either MSO_SHAPE_TYPE.LINE (genuine connector) OR an AUTO_SHAPE whose
    width or height is zero/negative (lines drawn as auto shapes — common in
    user-authored diagrams).
    """
    st = getattr(shape, "shape_type", None)
    if st == MSO_SHAPE_TYPE.LINE:
        return True
    try:
        w = int(shape.width)
        h = int(shape.height)
    except (AttributeError, TypeError):
        return False
    return w == 0 or h == 0
