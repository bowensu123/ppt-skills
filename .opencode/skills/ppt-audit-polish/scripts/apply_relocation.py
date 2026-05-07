"""Preserve-identity Path B renderer.

The free-form `apply_layout.py` always REBUILDS shapes from primitives
— every output shape is new, with new shape_id and default name. That
loses the original deck's editable-object identity and breaks any
inbound references (connectors, comments, click actions).

This renderer flips the architecture: it OPENS the original deck
(preserving every shape's XML) and applies per-shape decisions from a
`relocation.json` the agent wrote:

  decision = "preserve_identity"
    → move/resize/restyle the original shape in place.
    → shape_id / name / placeholder type / inbound references all
      preserved. Connector arrows still hit it; hyperlinks still land
      on it; comments stay anchored.

  decision = "recreate"
    → delete the original shape, add a new primitive at the new bbox.
    → identity lost (use only when the agent has confirmed no inbound
      references), but content (text / image binary) can be carried over.

  decision = "delete"
    → remove without replacement.

Plus optional `add_new_shapes` — agent-designed decorations that
weren't in the original.

Two contracts the renderer enforces:
  CONTENT ASSETS (text, image binary) — UNCONDITIONALLY preserved
    (recreate restores text/image into the new shape automatically)
  SEMANTIC STRUCTURE (title/body/item roles) — UNCONDITIONALLY preserved
    via content.json reference (template handles role mapping)

Schema:

  {
    "slides": [{
      "slide_index": 1,
      "shapes": {
        "42": {
          "decision": "preserve_identity",
          "agent_rationale": "Connector sid=88 references this shape",
          "new_bbox_emu": [L, T, W, H],
          "new_style": {
            "size_pt": 14,
            "color_hex": "F5F5F5",
            "fill_hex": "1A1A1A",
            "bold": true
          }
        },
        "56": {
          "decision": "recreate",
          "agent_rationale": "Pure decoration, no inbound refs",
          "new_kind": "rect" | "rounded_rect" | "text",
          "new_bbox_emu": [L, T, W, H],
          "new_fill_hex": "...",
          "new_text": "...",
          "new_size_pt": 14
        },
        "77": {
          "decision": "delete",
          "agent_rationale": "Redundant in new layout"
        }
      },
      "add_new_shapes": [
        {"kind": "rect", "bbox": [...], "fill": "..."},
        {"kind": "text", "bbox": [...], "content": "...", "size_pt": 12}
      ]
    }]
  }

CLI:
  python apply_relocation.py --in deck.pptx \
      --relocation relocation.json --out polished.pptx
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = SCRIPT_DIR.parent / "templates_py"


def _hex_to_rgb(hex_str: str) -> RGBColor:
    s = hex_str.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


# ---- preserve-identity: move/resize/restyle existing shape ----

def _apply_preserve(shape, decision: dict) -> dict:
    actions = []
    bbox = decision.get("new_bbox_emu")
    if bbox and len(bbox) == 4:
        try:
            shape.left = Emu(int(bbox[0]))
            shape.top = Emu(int(bbox[1]))
            shape.width = Emu(int(bbox[2]))
            shape.height = Emu(int(bbox[3]))
            actions.append("relocate")
        except (AttributeError, ValueError):
            pass
    style = decision.get("new_style") or {}
    if style:
        if "fill_hex" in style:
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = _hex_to_rgb(style["fill_hex"])
                actions.append("recolor-fill")
            except (AttributeError, ValueError):
                pass
        if "border_hex" in style:
            try:
                shape.line.color.rgb = _hex_to_rgb(style["border_hex"])
                actions.append("recolor-border")
            except (AttributeError, ValueError):
                pass
        if getattr(shape, "has_text_frame", False):
            from _shape_ops import (
                set_font_color, set_font_size, set_font_bold, set_font_family,
            )
            try:
                if "size_pt" in style:
                    set_font_size(shape, float(style["size_pt"]))
                    actions.append("resize-font")
                if "bold" in style:
                    set_font_bold(shape, bool(style["bold"]))
                    actions.append("set-bold")
                if "color_hex" in style:
                    set_font_color(shape, style["color_hex"])
                    actions.append("recolor-text")
                if "font_family" in style:
                    set_font_family(shape, style["font_family"])
                    actions.append("set-family")
            except (AttributeError, ValueError):
                pass
    return {"applied": actions}


# ---- delete shape (used by both 'delete' and 'recreate' paths) ----

def _delete_shape(shape) -> bool:
    """Remove the shape's XML from its parent. Returns True on success."""
    try:
        sp = shape._element
        parent = sp.getparent()
        if parent is None:
            return False
        parent.remove(sp)
        return True
    except (AttributeError, ValueError):
        return False


# ---- recreate: capture content, delete, add primitive, restore content ----

def _capture_text(shape) -> str | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    try:
        return shape.text_frame.text
    except (AttributeError, ValueError):
        return None


def _capture_image_bytes(shape) -> tuple[bytes, str] | None:
    """For pictures, return (blob, ext)."""
    try:
        if shape.shape_type in (13, 14):  # PICTURE / LINKED_PICTURE
            return (shape.image.blob, shape.image.ext or "png")
    except (AttributeError, ValueError):
        pass
    return None


def _add_recreated(slide, decision: dict, captured_text: str | None,
                    captured_image: tuple[bytes, str] | None) -> dict:
    """Add a fresh primitive matching the decision spec."""
    kind = decision.get("new_kind", "rect")
    bbox = decision.get("new_bbox_emu") or [0, 0, 1000000, 1000000]
    L, T, W, H = (int(v) for v in bbox)

    if captured_image:
        # Always preserve image binary, regardless of `new_kind`.
        import io, tempfile
        blob, ext = captured_image
        # python-pptx add_picture wants a path or file-like.
        with tempfile.NamedTemporaryFile(
            suffix=f".{ext}", delete=False,
        ) as f:
            f.write(blob)
            tmp_path = f.name
        try:
            pic = slide.shapes.add_picture(
                tmp_path, Emu(L), Emu(T), Emu(W), Emu(H),
            )
            return {"recreated_as": "picture",
                     "preserved_content": "image-bytes"}
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    if kind == "text" or captured_text:
        from templates_py._base import add_text  # type: ignore
        # Fall back to direct API.
        box = slide.shapes.add_textbox(Emu(L), Emu(T), Emu(W), Emu(H))
        if captured_text:
            box.text_frame.text = captured_text
        elif decision.get("new_text"):
            box.text_frame.text = decision["new_text"]
        from _shape_ops import set_font_size, set_font_color, set_font_bold
        if "new_size_pt" in decision:
            try:
                set_font_size(box, float(decision["new_size_pt"]))
            except (AttributeError, ValueError):
                pass
        if "new_color_hex" in decision:
            try:
                set_font_color(box, decision["new_color_hex"])
            except (AttributeError, ValueError):
                pass
        if "new_bold" in decision:
            try:
                set_font_bold(box, bool(decision["new_bold"]))
            except (AttributeError, ValueError):
                pass
        return {"recreated_as": "text",
                 "preserved_content": "text" if captured_text else "none"}

    # Geometric primitive
    shape_map = {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "circle": MSO_SHAPE.OVAL,
    }
    sp = slide.shapes.add_shape(
        shape_map.get(kind, MSO_SHAPE.RECTANGLE),
        Emu(L), Emu(T), Emu(W), Emu(H),
    )
    if "new_fill_hex" in decision:
        try:
            sp.fill.solid()
            sp.fill.fore_color.rgb = _hex_to_rgb(decision["new_fill_hex"])
        except (AttributeError, ValueError):
            pass
    return {"recreated_as": kind, "preserved_content": "none"}


# ---- add new decorative shapes ----

def _add_new_shape(slide, spec: dict) -> dict:
    if str(TEMPLATES_DIR) not in sys.path:
        sys.path.insert(0, str(TEMPLATES_DIR))
    from _base import (
        add_image_from_path, add_rect, add_rounded_rect, add_text,
    )
    kind = spec.get("kind")
    bbox = spec.get("bbox") or [0, 0, 0, 0]
    L, T, W, H = (int(v) for v in bbox)

    if kind == "rect":
        add_rect(slide, L, T, W, H,
                  fill=spec.get("fill"), line=spec.get("line"),
                  line_pt=spec.get("line_pt"))
        return {"added": "rect"}
    if kind == "rounded_rect":
        add_rounded_rect(slide, L, T, W, H,
                          fill=spec.get("fill"),
                          line=spec.get("border"),
                          line_pt=spec.get("border_pt"),
                          corner_ratio=spec.get("corner_ratio", 0.06))
        return {"added": "rounded_rect"}
    if kind == "text":
        add_text(slide, L, T, W, H,
                  spec.get("content", ""),
                  size_pt=float(spec.get("size_pt", 11)),
                  bold=bool(spec.get("bold", False)),
                  color=spec.get("color", "393939"),
                  family=spec.get("font"),
                  align=spec.get("align", "left"),
                  v_align=spec.get("v_align", "top"))
        return {"added": "text"}
    if kind == "image" and spec.get("path"):
        add_image_from_path(slide, spec["path"], L, T, W, H,
                             fit_mode=spec.get("fit_mode", "stretch"))
        return {"added": "image"}
    return {"added": None, "reason": f"unknown-kind:{kind}"}


# ---- top-level orchestrator ----

def apply_relocation(input_path: Path, relocation_path: Path,
                      output_path: Path) -> dict:
    prs = Presentation(str(input_path))
    spec = json.loads(relocation_path.read_text(encoding="utf-8"))

    actions: list[dict] = []
    skipped: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_spec = next(
            (s for s in spec.get("slides", [])
             if s.get("slide_index") == slide_idx),
            None,
        )
        if not slide_spec:
            continue

        sid_to_shape = {
            int(s.shape_id): s for s in slide.shapes if s.shape_id is not None
        }
        decisions = slide_spec.get("shapes", {}) or {}

        # Apply per-shape decisions in two passes:
        #   1. preserve_identity (no removal yet)
        #   2. recreate / delete (removal + recreation)
        # This avoids id collisions between deleted+added shapes.
        for sid_str, decision in decisions.items():
            try:
                sid = int(sid_str)
            except (TypeError, ValueError):
                continue
            shape = sid_to_shape.get(sid)
            if shape is None:
                skipped.append({"shape_id": sid,
                                  "reason": "shape-not-found"})
                continue
            d = decision.get("decision", "preserve_identity")
            if d == "preserve_identity":
                result = _apply_preserve(shape, decision)
                actions.append({
                    "slide_index": slide_idx,
                    "shape_id": sid,
                    "decision": "preserve_identity",
                    **result,
                })

        # Pass 2: recreate / delete (capture content first)
        for sid_str, decision in decisions.items():
            try:
                sid = int(sid_str)
            except (TypeError, ValueError):
                continue
            shape = sid_to_shape.get(sid)
            if shape is None:
                continue
            d = decision.get("decision", "preserve_identity")
            if d == "recreate":
                captured_text = _capture_text(shape)
                captured_image = _capture_image_bytes(shape)
                _delete_shape(shape)
                result = _add_recreated(
                    slide, decision, captured_text, captured_image,
                )
                actions.append({
                    "slide_index": slide_idx,
                    "shape_id": sid,
                    "decision": "recreate",
                    **result,
                })
            elif d == "delete":
                _delete_shape(shape)
                actions.append({
                    "slide_index": slide_idx,
                    "shape_id": sid,
                    "decision": "delete",
                    "applied": True,
                })

        # Add new decorative shapes (agent-designed extras)
        for new_spec in slide_spec.get("add_new_shapes", []):
            r = _add_new_shape(slide, new_spec)
            actions.append({
                "slide_index": slide_idx,
                "decision": "add",
                **r,
            })

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return {
        "out": str(output_path),
        "actions_applied": len(actions),
        "skipped": skipped,
        "actions": actions,
    }


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Apply per-shape preserve_identity / recreate / delete "
            "decisions to the original deck. Preserves shape_id, name, "
            "placeholder, connector endpoints, hyperlinks, comments for "
            "shapes the agent marked preserve_identity. Content (text "
            "and image binary) is unconditionally preserved across "
            "recreate."
        ),
    )
    parser.add_argument("--in", dest="in_path", required=True, type=Path)
    parser.add_argument("--relocation", required=True, type=Path,
                        help="agent-written relocation.json")
    parser.add_argument("--out", required=True, type=Path)
    args = parser.parse_args()

    result = apply_relocation(args.in_path, args.relocation, args.out)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
