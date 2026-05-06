from __future__ import annotations

import argparse
import json
from pathlib import Path
from statistics import median

from pptx import Presentation


ALIGN_TOLERANCE = 120000
FONT_TOLERANCE = 30000
SPACING_TOLERANCE = 140000
ROW_TOLERANCE = 220000
MIN_DIMENSION_TOLERANCE = 180000
DIMENSION_RATIO_TOLERANCE = 0.35


def _record(action_log: list[dict], slide_index: int, action: str, target: str) -> None:
    action_log.append({"slide_index": slide_index, "action": action, "target": target})


def _shape_font_sizes(shape) -> list[int]:
    if not getattr(shape, "has_text_frame", False):
        return []

    sizes: list[int] = []
    for paragraph in shape.text_frame.paragraphs:
        paragraph_sizes = [int(run.font.size) for run in paragraph.runs if run.font.size is not None]
        if paragraph_sizes:
            sizes.extend(paragraph_sizes)
        elif paragraph.font.size is not None:
            sizes.append(int(paragraph.font.size))
    return sizes


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text.strip()


def _dimension_close(value: int, target: int) -> bool:
    tolerance = max(MIN_DIMENSION_TOLERANCE, int(abs(target) * DIMENSION_RATIO_TOLERANCE))
    return abs(value - target) <= tolerance


def _row_groups(text_shapes: list) -> list[list]:
    rows: list[list] = []
    for shape in sorted(text_shapes, key=lambda item: int(item.top)):
        for row in rows:
            row_top = sum(int(item.top) for item in row) / len(row)
            if abs(int(shape.top) - row_top) <= ROW_TOLERANCE:
                row.append(shape)
                break
        else:
            rows.append([shape])
    return rows


def _peer_groups(text_shapes: list) -> list[list]:
    groups: list[list] = []
    for row in _row_groups(text_shapes):
        row_groups: list[list] = []
        for shape in sorted(row, key=lambda item: (int(item.width), int(item.height))):
            for group in row_groups:
                avg_width = int(sum(int(item.width) for item in group) / len(group))
                avg_height = int(sum(int(item.height) for item in group) / len(group))
                if _dimension_close(int(shape.width), avg_width) and _dimension_close(int(shape.height), avg_height):
                    group.append(shape)
                    break
            else:
                row_groups.append([shape])
        groups.extend(group for group in row_groups if len(group) >= 2)
    return groups


def _fix_boundary_overflow(shape, slide_index: int, slide_width: int, action_log: list[dict]) -> bool:
    if shape.shape_type == 3:
        return False

    changed = False
    right = int(shape.left) + int(shape.width)
    overflow = right - slide_width
    if overflow > 0:
        shape.left = max(0, int(shape.left) - overflow)
        changed = True
        _record(action_log, slide_index, "move-within-slide-bounds", getattr(shape, "name", "Shape"))

    return changed


def _normalize_group_tops(shapes: list, slide_index: int, action_log: list[dict]) -> None:
    target_top = int(median([int(shape.top) for shape in shapes]))
    for shape in shapes:
        if abs(int(shape.top) - target_top) > int(ALIGN_TOLERANCE / 2):
            shape.top = target_top
            _record(action_log, slide_index, "align-peer-row", getattr(shape, "name", "Text"))


def _normalize_group_fonts(shapes: list, slide_index: int, action_log: list[dict]) -> None:
    font_sizes = [size for shape in shapes for size in _shape_font_sizes(shape)]
    if not font_sizes:
        return
    target_font_size = int(median(font_sizes))
    for shape in shapes:
        for paragraph in shape.text_frame.paragraphs:
            paragraph_changed = False
            for run in paragraph.runs:
                if (
                    run.font.size is not None
                    and abs(int(run.font.size) - target_font_size) > FONT_TOLERANCE
                ):
                    run.font.size = target_font_size
                    paragraph_changed = True
            if (
                not paragraph_changed
                and paragraph.font.size is not None
                and abs(int(paragraph.font.size) - target_font_size) > FONT_TOLERANCE
            ):
                paragraph.font.size = target_font_size
                paragraph_changed = True
            if paragraph_changed:
                _record(action_log, slide_index, "normalize-font-hierarchy", getattr(shape, "name", "Text"))


def _normalize_group_gaps(shapes: list, slide_index: int, action_log: list[dict]) -> None:
    if len(shapes) < 3:
        return
    ordered = sorted(shapes, key=lambda shape: int(shape.left))
    gaps = [
        int(current.left) - (int(previous.left) + int(previous.width))
        for previous, current in zip(ordered, ordered[1:])
    ]
    if not gaps or max(gaps) - min(gaps) <= SPACING_TOLERANCE:
        return

    left_bound = min(int(shape.left) for shape in ordered)
    right_bound = max(int(shape.left) + int(shape.width) for shape in ordered)
    total_width = sum(int(shape.width) for shape in ordered)
    available_gap = right_bound - left_bound - total_width
    if available_gap < 0:
        return

    gap = int(available_gap / (len(ordered) - 1))
    cursor = left_bound
    for shape in ordered:
        if abs(int(shape.left) - cursor) > ALIGN_TOLERANCE:
            shape.left = cursor
            _record(action_log, slide_index, "normalize-peer-gaps", getattr(shape, "name", "Text"))
        cursor += int(shape.width) + gap


def apply_fixes(input_path: Path, output_path: Path, actions_output: Path | None) -> None:
    prs = Presentation(str(input_path))
    action_log: list[dict] = []
    skipped_items: list[dict] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        text_shapes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and _shape_text(shape)
            and int(shape.width) > 0
            and int(shape.height) > 0
        ]

        for shape in slide.shapes:
            if shape.shape_type == 3:
                skipped_items.append(
                    {
                        "slide_index": slide_index,
                        "target": getattr(shape, "name", "Chart"),
                        "reason": "chart requires manual review",
                    }
                )
                continue
            _fix_boundary_overflow(shape, slide_index, int(prs.slide_width), action_log)

        for peer_group in _peer_groups(text_shapes):
            _normalize_group_tops(peer_group, slide_index, action_log)
            _normalize_group_fonts(peer_group, slide_index, action_log)
            _normalize_group_gaps(peer_group, slide_index, action_log)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    if actions_output is not None:
        actions_output.parent.mkdir(parents=True, exist_ok=True)
        actions_output.write_text(
            json.dumps({"applied_actions": action_log, "skipped_items": skipped_items}, indent=2),
            encoding="utf-8",
        )


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--actions-output")
    args = parser.parse_args()

    actions_output = Path(args.actions_output) if args.actions_output else None
    apply_fixes(Path(args.input), Path(args.output), actions_output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
