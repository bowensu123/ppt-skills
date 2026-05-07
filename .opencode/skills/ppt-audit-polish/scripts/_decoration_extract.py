"""Capture vector decorations (custom shapes used as icons / accents).

Path B's existing extractors cover:
  * extract_content.py     → text content (title, items, footer)
  * _asset_extract.py      → picture/icon binaries (PNG/JPEG/SVG/EMF)

But many decks render "icons" as CUSTOM SHAPES rather than embedded
images: a colored oval with "01" inside is an MSO_SHAPE.OVAL + a text
frame, NOT an MSO_SHAPE_TYPE.PICTURE. Those slip through the picture
extractor. Same for accent bars (RECTANGLE), divider lines, dot
markers, badges — all geometric primitives that visually function as
decorations.

This module walks every slide (and the inheritance chain — slide
layout + slide master) and emits a `decorations.json` manifest of
non-picture, non-text shapes that look decorative. The agent reads
this when designing layout.json and re-creates each decoration with
`kind: rect / circle / rounded_rect / line` elements.

Output (decorations.json):
{
  "input": "...",
  "extracted_count": N,
  "decorations": [
    {
      "decoration_id": "d01",
      "shape_id": 7,
      "slide_index": 1,
      "source": "slide" | "layout" | "master",
      "auto_shape_type": "OVAL",
      "kind_for_layout": "circle",
      "bbox_emu": [L, T, W, H],
      "fill_hex": "D97757",
      "line_hex": "FFFFFF",
      "line_pt": 0.75,
      "rotation": 0.0,
      "text": "01",
      "z_index": 7,
      "decorative_score": 0.85    // 0-1; high = strongly looks like decoration
    }
  ]
}

Decoration heuristic (decorative_score):
  small area (< 10% slide)                 +0.4
  has fill + simple geometric type         +0.3
  near edge or in corner                   +0.2
  has very short text or none              +0.1
  → 0..1; > 0.5 included by default
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Auto-shape types that are commonly used as decorations.
_DECO_AUTO_SHAPES = {
    "OVAL", "RECTANGLE", "ROUNDED_RECTANGLE",
    "RIGHT_TRIANGLE", "DIAMOND", "ISOCELES_TRIANGLE",
    "PARALLELOGRAM", "TRAPEZOID", "PENTAGON", "HEXAGON",
    "CHEVRON", "RIGHT_ARROW", "LEFT_ARROW", "UP_ARROW", "DOWN_ARROW",
    "STAR_5_POINT", "STAR_6_POINT",
    "PIE", "CHORD",
}

# auto_shape_type → friendly kind for apply_layout.py.
_KIND_MAP = {
    "OVAL": "circle",
    "RECTANGLE": "rect",
    "ROUNDED_RECTANGLE": "rounded_rect",
}

# Score thresholds.
SCORE_INCLUDE = 0.5
SLIDE_AREA_DECORATIVE_FRAC = 0.10        # ≤ 10% of slide = small enough to be decoration
EDGE_FRAC = 0.18                          # within 18% of any edge
MAX_DECORATION_TEXT_CHARS = 12            # short or no text = more decoration-like


def _safe_str(x):
    return str(x) if x is not None else None


def _shape_is_picture(shape) -> bool:
    try:
        return shape.shape_type in (MSO_SHAPE_TYPE.PICTURE,
                                     MSO_SHAPE_TYPE.LINKED_PICTURE)
    except (AttributeError, ValueError):
        return False


def _shape_is_pure_text(shape) -> bool:
    """A shape that's just a text box (no fill, no border)."""
    try:
        if not getattr(shape, "has_text_frame", False):
            return False
        # Text frames with no fill / no line are pure text containers.
        # If they have visible fill AND substantive text → could be a card,
        # not a decoration.
        return shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    except (AttributeError, ValueError):
        return False


def _auto_shape_type_name(shape) -> str | None:
    try:
        ast = shape.auto_shape_type
        if ast is None:
            return None
        return ast.name if hasattr(ast, "name") else str(ast)
    except (AttributeError, ValueError):
        return None


def _fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        rgb = fill.fore_color.rgb
        return str(rgb) if rgb is not None else None
    except (AttributeError, ValueError, KeyError, TypeError):
        return None


def _line_hex(shape) -> str | None:
    try:
        rgb = shape.line.color.rgb
        return str(rgb) if rgb is not None else None
    except (AttributeError, ValueError, KeyError, TypeError):
        return None


def _line_pt(shape) -> float | None:
    try:
        from pptx.util import Pt
        w = shape.line.width
        if w is None:
            return None
        return float(w) / Pt(1)
    except (AttributeError, ValueError, KeyError, TypeError):
        return None


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return "\n".join(
            p.text for p in shape.text_frame.paragraphs
        ).strip()
    except (AttributeError, ValueError):
        return ""


def _decorative_score(shape, slide_w: int, slide_h: int) -> float:
    """0..1 how "decoration-like" this shape is."""
    score = 0.0
    try:
        sw = int(shape.width or 0); sh = int(shape.height or 0)
        sl = int(shape.left or 0); st = int(shape.top or 0)
    except (TypeError, ValueError):
        return 0.0
    if sw <= 0 or sh <= 0:
        return 0.0
    area_ratio = (sw * sh) / max(slide_w * slide_h, 1)
    if area_ratio < SLIDE_AREA_DECORATIVE_FRAC:
        score += 0.4
    elif area_ratio < SLIDE_AREA_DECORATIVE_FRAC * 2:
        score += 0.2

    # Has fill and a recognizable simple shape
    if _fill_hex(shape) is not None:
        ast = _auto_shape_type_name(shape)
        if ast in _DECO_AUTO_SHAPES:
            score += 0.3
    elif _line_hex(shape) is not None:
        # Stroked-only shapes (dividers, accent strokes)
        score += 0.2

    # Edge / corner positioning
    edge_x = slide_w * EDGE_FRAC
    edge_y = slide_h * EDGE_FRAC
    near_edge = (
        sl < edge_x or st < edge_y
        or (sl + sw) > (slide_w - edge_x)
        or (st + sh) > (slide_h - edge_y)
    )
    if near_edge:
        score += 0.2

    # Short text / no text
    text = _shape_text(shape)
    if len(text) <= MAX_DECORATION_TEXT_CHARS:
        score += 0.1

    return min(score, 1.0)


def _iter_shapes_with_z(shapes):
    """Yield (shape, z_index) recursing into groups."""
    z = 0
    def walk(parent_shapes):
        nonlocal z
        for shape in parent_shapes:
            try:
                stype = shape.shape_type
            except (AttributeError, ValueError):
                stype = None
            if stype == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes)
                continue
            yield shape, z
            z += 1
    yield from walk(shapes)


def _classify_shape(shape, slide_w: int, slide_h: int) -> dict | None:
    """Return a decoration entry if shape qualifies, else None."""
    if _shape_is_picture(shape):
        return None
    # Skip pure text-frame boxes (already in content.json).
    # But ALLOW shapes that have BOTH a visible fill AND text — those are
    # often badges (oval + "01") — exactly what we want to capture.
    has_text = bool(_shape_text(shape))
    has_fill = _fill_hex(shape) is not None
    if has_text and not has_fill:
        # Text-only — covered by extract_content.py.
        return None

    score = _decorative_score(shape, slide_w, slide_h)
    if score < SCORE_INCLUDE:
        return None

    auto_type = _auto_shape_type_name(shape)
    sid = getattr(shape, "shape_id", None)
    if sid is None:
        return None

    return {
        "shape_id": int(sid),
        "auto_shape_type": auto_type,
        "kind_for_layout": _KIND_MAP.get(auto_type, "rect"),
        "bbox_emu": [
            int(shape.left or 0), int(shape.top or 0),
            int(shape.width or 0), int(shape.height or 0),
        ],
        "fill_hex": _fill_hex(shape),
        "line_hex": _line_hex(shape),
        "line_pt": _line_pt(shape),
        "rotation": float(getattr(shape, "rotation", 0.0) or 0.0),
        "text": _safe_str(_shape_text(shape)) or "",
        "decorative_score": round(score, 2),
    }


def extract_decorations(input_path: Path, work_dir: Path,
                         include_master: bool = True) -> dict:
    """Walk slides + (optionally) layout + master, capture decorations."""
    prs = Presentation(str(input_path))
    slide_w = int(prs.slide_width); slide_h = int(prs.slide_height)
    work_dir.mkdir(parents=True, exist_ok=True)

    decorations: list[dict] = []
    deco_counter = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape, z in _iter_shapes_with_z(slide.shapes):
            entry = _classify_shape(shape, slide_w, slide_h)
            if entry is None:
                continue
            deco_counter += 1
            entry["decoration_id"] = f"d{deco_counter:02d}"
            entry["slide_index"] = slide_idx
            entry["source"] = "slide"
            entry["z_index"] = z
            decorations.append(entry)

        if include_master:
            # Walk slide layout shapes (chrome / placeholders / decorations
            # inherited from layout level).
            try:
                for shape, z in _iter_shapes_with_z(slide.slide_layout.shapes):
                    entry = _classify_shape(shape, slide_w, slide_h)
                    if entry is None:
                        continue
                    # Avoid duplicate IDs across slides — layout shapes are
                    # shared, but each slide's render counts them once.
                    deco_counter += 1
                    entry["decoration_id"] = f"d{deco_counter:02d}"
                    entry["slide_index"] = slide_idx
                    entry["source"] = "layout"
                    entry["z_index"] = z
                    decorations.append(entry)
            except (AttributeError, ValueError):
                pass

            # Walk slide master shapes (logos / global decorations).
            try:
                for shape, z in _iter_shapes_with_z(
                    slide.slide_layout.slide_master.shapes
                ):
                    entry = _classify_shape(shape, slide_w, slide_h)
                    if entry is None:
                        continue
                    deco_counter += 1
                    entry["decoration_id"] = f"d{deco_counter:02d}"
                    entry["slide_index"] = slide_idx
                    entry["source"] = "master"
                    entry["z_index"] = z
                    decorations.append(entry)
            except (AttributeError, ValueError):
                pass

    manifest = {
        "input": str(input_path),
        "slide_count": len(prs.slides),
        "extracted_count": len(decorations),
        "decorations": decorations,
    }
    (work_dir / "decorations.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Extract vector-shape decorations (ovals / rects / rounded "
            "rects / lines used as icons) into a manifest the agent can "
            "read when designing layout.json. Includes shapes inherited "
            "from slide layout / slide master unless --no-master."
        ),
    )
    parser.add_argument("--in", dest="in_path", required=True, type=Path)
    parser.add_argument("--work-dir", required=True, type=Path)
    parser.add_argument("--no-master", action="store_true",
                        help="skip slide layout + master shape walk")
    args = parser.parse_args()

    manifest = extract_decorations(
        args.in_path, args.work_dir, include_master=not args.no_master,
    )
    print(json.dumps({
        "extracted_count": manifest["extracted_count"],
        "manifest": str(args.work_dir / "decorations.json"),
    }, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
