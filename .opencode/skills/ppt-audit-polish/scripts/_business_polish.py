"""Business-grade visual polish that DOES NOT change content.

Path A's existing detectors find structural bugs (overlap, alignment,
spacing). Once those are fixed the deck still typically looks "OK but
not polished" — flat colors, inconsistent shadows, no visual hierarchy
beyond what the original author put in. This module is a single-pass
elevator from "OK" to "business-grade":

  1. Smart typography     — apply theme type-scale to every detected role
  2. Geometric consistency — unify corner radius and shadow across cards
  3. Hierarchy decoration — accent bar above title, divider before footer
  4. Subtle surface tint  — lighter background behind subtitle/section heads

Hard constraints:
  * NEVER modify any text content (run.text stays untouched)
  * NEVER delete or hide existing shapes
  * Newly added decorative shapes get an `__bp_deco` flag in the name so
    repeat runs can detect and skip them (idempotent)
  * Decorations go BEHIND existing content via z-order so they can't
    accidentally cover anything

The orchestrator function `apply_business_polish(prs, ...)` is wired
into mutate.py as the `polish-business` op. Each helper is also
exposed for fine-grained agent control.
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

from _shape_ops import (
    set_corner_radius,
    set_font_bold,
    set_font_color,
    set_font_family,
    set_font_size,
    set_shadow,
    set_solid_fill,
)


def _theme_dict(theme) -> dict:
    """Accept either a dict or _common.Theme dataclass."""
    if hasattr(theme, "raw"):
        return theme.raw
    return theme


# Decoration shape name marker — used for idempotency.
DECO_NAME_PREFIX = "__bp_deco"

# Type scale ratios. We use the major-third ratio (1.25), recognizable
# as the "standard editorial scale". theme.typography.size_pt values
# already encode this, so we pass them through.
DEFAULT_SCALE_PT = {
    "title":    32,
    "subtitle": 16,
    "h2":       16,
    "body":     12,
    "caption":  10,
    "badge":    10,
}

# Visual polish levels — agent picks intensity.
LEVEL_SUBTLE   = 1   # typography + corner-radius unification + shadow consistency
LEVEL_STANDARD = 2   # + accent bar above title + divider above footer
LEVEL_RICH     = 3   # + tinted surface behind subtitle


# ---- color helpers ----

def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    s = hex_str.lstrip("#")
    return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)


def _rgb_to_hex(rgb: tuple[int, int, int]) -> str:
    return "{:02X}{:02X}{:02X}".format(*rgb)


def _tint(hex_str: str, amount: float) -> str:
    """Lighten a color by `amount` (0..1). 0.05 = 5% closer to white."""
    r, g, b = _hex_to_rgb(hex_str)
    r = int(r + (255 - r) * amount)
    g = int(g + (255 - g) * amount)
    b = int(b + (255 - b) * amount)
    return _rgb_to_hex((r, g, b))


# ---- role discovery (lightweight; doesn't run detect_roles again) ----

def _shape_by_id(prs, shape_id: int):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_id", None) == shape_id:
                return slide, shape
    return None, None


def _shapes_by_role(prs, role_data: dict, role: str) -> list:
    """Return python-pptx shapes whose detected role matches."""
    sids = []
    for slide in role_data.get("slides", []):
        for entry in slide.get("shapes", []):
            if entry.get("role") == role:
                sids.append(entry["shape_id"])
    out = []
    for sid in sids:
        _, shape = _shape_by_id(prs, sid)
        if shape is not None:
            out.append(shape)
    return out


# ---- 1. smart typography ----

def apply_smart_typography(prs, theme, role_data: dict) -> list[dict]:
    """For every detected role, apply theme typography (size + bold +
    color + family). Returns one log entry per shape touched.
    """
    theme = _theme_dict(theme)
    log: list[dict] = []
    typ = theme.get("typography", {})
    sizes = {**DEFAULT_SCALE_PT, **typ.get("size_pt", {})}
    bolds = typ.get("bold", {})
    color_role = typ.get("color_role", {})
    palette = theme.get("palette", {})
    family = typ.get("font_family", "Microsoft YaHei")

    for role in ("title", "subtitle", "h2", "heading", "body", "caption", "badge"):
        # Roles like 'heading' should map to h2 sizing if not declared.
        size_key = role if role in sizes else ("h2" if role == "heading" else "body")
        size_pt = sizes.get(size_key)
        if size_pt is None:
            continue
        for shape in _shapes_by_role(prs, role_data, role):
            try:
                set_font_size(shape, float(size_pt))
                set_font_family(shape, family, include_eastasia=True)
                set_font_bold(shape, bool(bolds.get(size_key, False)))
                color_key = color_role.get(size_key)
                if color_key and palette.get(color_key):
                    set_font_color(shape, palette[color_key])
                log.append({
                    "action": "smart-typography",
                    "shape_id": shape.shape_id,
                    "role": role,
                    "size_pt": size_pt,
                })
            except (AttributeError, ValueError):
                continue
    return log


# ---- 2. geometric consistency ----

def unify_corner_radius(prs, ratio: float = 0.06) -> list[dict]:
    """Set every rounded-rectangle container to the same corner ratio.

    Skips connectors, lines, and shapes without a corner radius adjustment.
    """
    log: list[dict] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                # Only proceed if the shape supports a corner adjustment.
                # set_corner_radius returns None if the shape isn't a
                # rounded-rect type.
                result = set_corner_radius(shape, ratio)
                if result is not None:
                    log.append({
                        "action": "unify-corner-radius",
                        "shape_id": shape.shape_id,
                        "ratio": ratio,
                    })
            except (AttributeError, ValueError, KeyError):
                continue
    return log


def unify_card_shadow(prs, theme, min_w_emu: int = 1200000,
                      min_h_emu: int = 800000) -> list[dict]:
    """Apply consistent subtle shadow to every container shape that's
    big enough to be a 'card'. Subtler than default: low alpha, small blur.
    """
    theme = _theme_dict(theme)
    log: list[dict] = []
    palette = theme.get("palette", {})
    shadow_color = palette.get("text_muted", "6F6F6F")
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if not getattr(shape, "fill", None):
                    continue
                w = int(shape.width or 0); h = int(shape.height or 0)
                if w < min_w_emu or h < min_h_emu:
                    continue
                set_shadow(shape, hex_color=shadow_color, blur_pt=3.0,
                           dist_pt=1.5, alpha=0.18)
                log.append({"action": "unify-shadow", "shape_id": shape.shape_id})
            except (AttributeError, ValueError, KeyError):
                continue
    return log


# ---- 3. hierarchy decorations (added shapes, idempotent) ----

def _has_existing_deco(slide, marker: str) -> bool:
    for shape in slide.shapes:
        if (getattr(shape, "name", "") or "").startswith(marker):
            return True
    return False


def _slide_dims(slide):
    """Return (width_emu, height_emu) of the slide."""
    prs = slide.part.package.presentation_part.presentation
    return int(prs.slide_width), int(prs.slide_height)


def add_title_accent_bar(prs, theme, role_data: dict,
                          height_emu: int = 36000) -> list[dict]:
    """Add a thin primary-color rectangle above each detected title.

    Goes BEHIND the title in z-order. Skips slides where there's no
    space above the title (top - height_emu < 0).
    """
    theme = _theme_dict(theme)
    log: list[dict] = []
    primary = theme.get("palette", {}).get("primary", "0F62FE")
    marker = f"{DECO_NAME_PREFIX}-accent-bar"

    for slide_idx, slide in enumerate(prs.slides, start=1):
        if _has_existing_deco(slide, marker):
            continue
        # Find this slide's title.
        slide_role = next(
            (s for s in role_data.get("slides", []) if s.get("slide_index") == slide_idx),
            None,
        )
        if not slide_role:
            continue
        title_sid = next(
            (e["shape_id"] for e in slide_role.get("shapes", [])
             if e.get("role") == "title"),
            None,
        )
        if title_sid is None:
            continue
        _, title_shape = _shape_by_id(prs, title_sid)
        if title_shape is None or title_shape.top is None:
            continue
        # Bar sits 1.5x bar-height above title.
        top = int(title_shape.top) - int(height_emu * 1.5)
        if top < 0:
            continue
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(title_shape.left or 0),
            top,
            min(int(title_shape.width or 0), 2000000),  # cap width to ~2.2"
            height_emu,
        )
        bar.name = marker
        try:
            set_solid_fill(bar, primary)
        except Exception:
            pass
        # Send to back so it can never cover content.
        try:
            sp = bar._element
            spTree = sp.getparent()
            spTree.remove(sp)
            spTree.insert(2, sp)  # after nvGrpSpPr + grpSpPr
        except Exception:
            pass
        log.append({
            "action": "add-accent-bar",
            "slide_index": slide_idx,
            "above_shape_id": title_sid,
        })
    return log


def add_footer_divider(prs, theme, role_data: dict) -> list[dict]:
    """Add a thin divider line above the bottom band of each slide,
    spanning the slide width with a margin. Skips slides with existing
    divider deco."""
    from pptx.enum.shapes import MSO_CONNECTOR

    theme = _theme_dict(theme)
    log: list[dict] = []
    marker = f"{DECO_NAME_PREFIX}-footer-divider"
    border_color = theme.get("palette", {}).get("border", "DDE1E6")
    margin = 457200  # 0.5"

    for slide_idx, slide in enumerate(prs.slides, start=1):
        if _has_existing_deco(slide, marker):
            continue
        sw, sh = _slide_dims(slide)
        # Place divider 12% above slide bottom.
        y = int(sh * 0.88)
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, margin, y, sw - margin, y,
        )
        line.name = marker
        try:
            line.line.color.rgb = _rgb_from_hex(border_color)
            line.line.width = Pt(0.5)
        except Exception:
            pass
        try:
            sp = line._element
            spTree = sp.getparent()
            spTree.remove(sp)
            spTree.insert(2, sp)
        except Exception:
            pass
        log.append({"action": "add-footer-divider", "slide_index": slide_idx})
    return log


def _rgb_from_hex(hex_str: str):
    from pptx.dml.color import RGBColor
    s = hex_str.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


# ---- 4. subtle surface tint ----

def tint_subtitle_background(prs, theme, role_data: dict) -> list[dict]:
    """Add a 5% tint rectangle behind each subtitle. The tint is the
    primary color, very lightened — gives a subtle 'eyebrow' band
    without competing with the main content."""
    theme = _theme_dict(theme)
    log: list[dict] = []
    primary = theme.get("palette", {}).get("primary", "0F62FE")
    tint_color = _tint(primary, 0.92)  # 92% closer to white = very pale
    marker = f"{DECO_NAME_PREFIX}-subtitle-tint"

    for slide_idx, slide in enumerate(prs.slides, start=1):
        if _has_existing_deco(slide, marker):
            continue
        slide_role = next(
            (s for s in role_data.get("slides", []) if s.get("slide_index") == slide_idx),
            None,
        )
        if not slide_role:
            continue
        sub_sids = [e["shape_id"] for e in slide_role.get("shapes", [])
                    if e.get("role") == "subtitle"]
        for sid in sub_sids:
            _, sub_shape = _shape_by_id(prs, sid)
            if sub_shape is None:
                continue
            tint = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(sub_shape.left or 0) - 91440,
                int(sub_shape.top or 0) - 45720,
                int(sub_shape.width or 0) + 182880,
                int(sub_shape.height or 0) + 91440,
            )
            tint.name = marker
            try:
                set_solid_fill(tint, tint_color)
                tint.line.fill.background()  # no border
            except Exception:
                pass
            try:
                sp = tint._element
                spTree = sp.getparent()
                spTree.remove(sp)
                spTree.insert(2, sp)
            except Exception:
                pass
            log.append({
                "action": "tint-subtitle-bg",
                "slide_index": slide_idx,
                "behind_shape_id": sid,
            })
    return log


# ---- 5. content-aware theme picker ----

# Domain keyword fingerprints for the 4 built-in themes. Hand-tuned;
# the agent can override with --theme.
THEME_KEYWORDS = {
    # ── original 5 ────────────────────────────────────────────────
    "clean-tech":     ["框架", "模型", "AI", "算法", "API", "data", "model",
                       "tensor", "neural", "vision", "embed",
                       "tech", "infrastructure", "deploy"],
    "business-warm":  ["营收", "增长", "客户", "市场", "ROI", "growth",
                       "revenue", "customer", "market", "sales", "strategy",
                       "Q1", "Q2", "Q3", "Q4", "KPI", "OKR", "成本", "利润"],
    "academic-soft":  ["研究", "方法", "实验", "论文", "假设", "research",
                       "method", "experiment", "hypothesis", "analysis",
                       "literature", "theory", "study", "结论", "数据集"],
    "editorial-dark": ["故事", "观点", "品牌", "叙事", "narrative", "brand",
                       "story", "voice", "perspective", "editorial",
                       "manifesto", "vision", "design"],
    "claude-code":    ["claude", "anthropic", "cli", "terminal", "command",
                       "bash", "shell", "code", "agent", "developer", "dev",
                       "cursor", "vscode", "ide", "repo", "git", "commit",
                       "function", "method", "script", "subprocess",
                       "stdin", "stdout", "Microsoft", "build", "deploy",
                       "compile", "runtime", "env", "config", "yaml", "json",
                       "Python", "TypeScript", "Rust", "Go", "Node",
                       "工具链", "脚本", "CLI", "终端", "命令行"],

    # ── new 10 (per user request) ────────────────────────────────
    "minimalist-business":
        ["麦肯锡", "贝恩", "BCG", "提案", "transformation", "并购",
         "M&A", "重组", "战略选择", "McKinsey", "Bain", "consulting firm",
         "PMO", "战略转型", "组织设计"],
    "consulting":
        ["金字塔", "瀑布图", "MECE", "takeaway", "矩阵图",
         "议题树", "假设驱动", "敏感性分析", "结构化分析",
         "盈利模型", "5C", "4P", "SWOT"],
    "modern-tech":
        ["Stripe", "Linear", "Vercel", "Notion", "SaaS",
         "subscription", "platform", "微服务", "云原生",
         "DevOps", "Kubernetes", "container", "serverless"],
    "corporate-classic":
        ["年报", "上市", "财报披露", "国资", "央企", "信用评级",
         "审计", "合规", "董事会", "股东大会", "监管", "ESG",
         "信息披露", "投资者关系"],
    "pitch-deck":
        ["融资", "估值", "种子轮", "天使轮", "BP", "TAM", "SAM", "SOM",
         "Y Combinator", "Sequoia", "a16z", "valuation", "term sheet",
         "preseed", "Series", "投资人", "募资"],
    "editorial-magazine":
        ["杂志", "lifestyle", "campaign", "时尚", "新品发布",
         "品牌故事", "Vogue", "GQ", "Elle", "时装", "潮流",
         "buyer", "lookbook", "趋势报告"],
    "data-heavy":
        ["GMV", "DAU", "MAU", "MRR", "LTV", "CAC", "运营复盘",
         "财务复盘", "市场研究", "数据看板", "ARPU", "留存率",
         "渠道分析", "漏斗", "cohort"],
    "academic-research":
        ["白皮书", "对照实验", "假设检验", "实证研究", "p值",
         "回归分析", "样本", "置信区间", "RCT", "meta分析",
         "学术演讲", "research paper", "abstract"],
    "creative-agency":
        ["VI 设计", "agency pitch", "创意提案", "campaign",
         "advertising", "广告创意", "品牌策划", "brand identity",
         "key visual", "tagline", "art direction"],
    "dark-premium":
        ["私募", "奢侈品", "高净值", "premium", "VIP", "luxury",
         "PE", "VC", "家族办公室", "private banking", "高端定制",
         "藏品", "winemaker", "腕表", "yacht"],
}


def pick_theme(content_text: str, themes_dir: Path) -> str:
    """Score each theme by keyword frequency in content, return the
    name (matching themes/<name>.json). Falls back to clean-tech.
    """
    content_lower = content_text.lower()
    best_theme = "clean-tech"
    best_score = 0
    for theme_name, kws in THEME_KEYWORDS.items():
        score = sum(content_lower.count(k.lower()) for k in kws)
        if score > best_score:
            best_score = score
            best_theme = theme_name
    return best_theme


# ---- top-level orchestrator ----

def apply_business_polish(prs, theme, role_data: dict,
                           level: int = LEVEL_STANDARD) -> dict:
    """Run the full polish pass.

    level 1 (subtle):   typography + corner radius + shadow
    level 2 (standard): + accent bar + footer divider
    level 3 (rich):     + subtitle background tint
    """
    actions: list[dict] = []

    actions.extend(apply_smart_typography(prs, theme, role_data))
    actions.extend(unify_corner_radius(prs, ratio=0.06))
    actions.extend(unify_card_shadow(prs, theme))

    if level >= LEVEL_STANDARD:
        actions.extend(add_title_accent_bar(prs, theme, role_data))
        actions.extend(add_footer_divider(prs, theme, role_data))

    if level >= LEVEL_RICH:
        actions.extend(tint_subtitle_background(prs, theme, role_data))

    return {"level": level, "actions_applied": len(actions), "actions": actions}
