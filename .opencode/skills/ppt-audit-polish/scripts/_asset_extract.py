"""Walk a PPTX and extract every picture / icon to a side-car assets/
folder, emitting a structured manifest the agent can read.

This is the DETERMINISTIC half of the visual-asset preservation flow.
The SEMANTIC question — "this image belongs to item N / it's a slide
decoration / it's junk" — is left to the agent. The agent reads:
  * content.json      (text-only output of extract_content.py)
  * assets-manifest.json (this module's output)
  * annotated render   (state_summary's slide-NNN.annotated.png)
…and updates content.json with `items[].image` attributions and a
`decorations[]` list. Templates then render accordingly.

The split keeps rules in their lane (binary extraction is exact and
fast; bbox / size / position metadata is exact) while leaving the
fundamentally judgment-based question of "what does this image MEAN"
to the multimodal model.

Output structure (assets-manifest.json):
{
  "assets": [
    {
      "asset_id": "a01",
      "shape_id": 42,
      "slide_index": 1,
      "kind": "picture",
      "path": "assets/sid_42.png",
      "ext": "png",
      "bbox_emu": [L, T, W, H],
      "z_index": 7,
      "alt": "Shape Label",
      "size_bytes": 23456,
      "decorative_hint": false   // small + edge position = likely logo
    },
    ...
  ],
  "input": "...",
  "extracted_count": N
}
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Reasonable default — anything in this set is a "picture-like" asset.
PICTURE_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}

# Threshold for "decorative" hint: shape area <= this fraction of slide area
# AND positioned within this fraction of an edge → likely a logo / chrome.
DECORATIVE_AREA_RATIO = 0.04
DECORATIVE_EDGE_RATIO = 0.12


def _iter_picture_shapes(slide):
    """Yield (shape, parent_chain) for every picture, recursing into groups.

    parent_chain is a list of group shape names (empty for top-level).
    Useful for the agent to disambiguate nested icons.
    """
    def walk(shapes, chain):
        for shape in shapes:
            try:
                stype = shape.shape_type
            except (AttributeError, ValueError):
                stype = None
            if stype == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes, chain + [shape.name or ""])
                continue
            if stype in PICTURE_TYPES:
                yield shape, list(chain)
    yield from walk(slide.shapes, [])


def _picture_crop(shape) -> dict | None:
    """Return srcRect as fractions 0..1 if the picture has one, else None.

    PPTX stores srcRect with values in 1000ths of a percent. blipFill
    sits in either `p:` (slide pic) or `a:` (drawingml pic in
    graphicFrame) namespace — try both.
    """
    from lxml import etree
    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    src = None
    try:
        for path in (".//p:blipFill/a:srcRect", ".//a:blipFill/a:srcRect"):
            src = shape._element.find(path, nsmap)
            if src is not None:
                break
    except (AttributeError, etree.LxmlError):
        return None
    if src is None:
        return None
    out: dict[str, float] = {}
    for side in ("left", "top", "right", "bottom"):
        raw = src.get(side)
        if raw is not None:
            try:
                out[side] = float(raw) / 100000.0
            except (TypeError, ValueError):
                continue
    return out if out else None


def _decorative_hint(shape, slide_w: int, slide_h: int) -> bool:
    """Heuristic: small area + close to slide edge → probably a logo/chrome."""
    try:
        sw = int(shape.width or 0); sh = int(shape.height or 0)
        sl = int(shape.left or 0); st = int(shape.top or 0)
    except (TypeError, ValueError):
        return False
    if sw <= 0 or sh <= 0:
        return False
    area_ratio = (sw * sh) / max(slide_w * slide_h, 1)
    if area_ratio > DECORATIVE_AREA_RATIO:
        return False
    edge_threshold_x = slide_w * DECORATIVE_EDGE_RATIO
    edge_threshold_y = slide_h * DECORATIVE_EDGE_RATIO
    near_left = sl < edge_threshold_x
    near_right = (sl + sw) > (slide_w - edge_threshold_x)
    near_top = st < edge_threshold_y
    near_bottom = (st + sh) > (slide_h - edge_threshold_y)
    return near_left or near_right or near_top or near_bottom


def extract_assets(input_path: Path, work_dir: Path) -> dict:
    """Walk every slide, write each picture's binary to assets/, return manifest.

    Files: <work_dir>/assets/sid_<N>.<ext> — keyed by shape_id so the
    agent can correlate manifest entries with the annotated render.
    """
    prs = Presentation(str(input_path))
    slide_w = int(prs.slide_width); slide_h = int(prs.slide_height)
    assets_dir = work_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    manifest_entries: list[dict] = []
    asset_counter = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Build a z-order map keyed by shape_id (stable across iterations
        # — id(shape) isn't, because python-pptx returns fresh wrappers).
        z_lookup: dict[int, int] = {}
        for idx, shape in enumerate(slide.shapes):
            sid = getattr(shape, "shape_id", None)
            if sid is not None:
                z_lookup[int(sid)] = idx
        for shape, parent_chain in _iter_picture_shapes(slide):
            try:
                blob = shape.image.blob
                ext = (shape.image.ext or "png").lower().lstrip(".")
            except (AttributeError, ValueError, KeyError):
                continue
            sid = getattr(shape, "shape_id", None)
            if sid is None:
                continue
            asset_counter += 1
            asset_id = f"a{asset_counter:02d}"
            fname = f"sid_{sid}.{ext}"
            (assets_dir / fname).write_bytes(blob)
            manifest_entries.append({
                "asset_id": asset_id,
                "shape_id": int(sid),
                "slide_index": slide_idx,
                "kind": "picture",
                "path": f"assets/{fname}",
                "ext": ext,
                "bbox_emu": [
                    int(shape.left or 0), int(shape.top or 0),
                    int(shape.width or 0), int(shape.height or 0),
                ],
                "z_index": z_lookup.get(int(sid), -1),
                "alt": shape.name or "",
                "parent_groups": parent_chain,
                "size_bytes": len(blob),
                "decorative_hint": _decorative_hint(shape, slide_w, slide_h),
                "crop": _picture_crop(shape),
            })

    manifest = {
        "input": str(input_path),
        "slide_count": len(prs.slides),
        "extracted_count": len(manifest_entries),
        "assets": manifest_entries,
    }
    (work_dir / "assets-manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Extract every picture/icon from a PPTX into a side-car "
            "assets/ folder + manifest.json. The semantic question of "
            "'which item does each image belong to' is left to the agent."
        ),
    )
    parser.add_argument("--in", dest="in_path", required=True, type=Path)
    parser.add_argument("--work-dir", required=True, type=Path)
    args = parser.parse_args()

    args.work_dir.mkdir(parents=True, exist_ok=True)
    manifest = extract_assets(args.in_path, args.work_dir)
    print(json.dumps({
        "extracted_count": manifest["extracted_count"],
        "manifest": str(args.work_dir / "assets-manifest.json"),
    }, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
