"""Unified granular mutation CLI for ppt-audit-polish.

Every subcommand follows the same contract:
  * reads --in <pptx> and writes --out <pptx>
  * mutates only the shapes/slides explicitly referenced
  * emits a structured JSON record on stdout describing what changed
  * exits 0 on success, 2 on usage error, 3 on runtime error

This contract is what allows the L4 orchestrator to chain calls and what
allows OpenCode-driven models to do "render → look → mutate → render"
loops without any hidden state.

For the full op catalog, run:  python mutate.py list-ops --json
"""
from __future__ import annotations

import argparse
import inspect
import json
import sys
from pathlib import Path
from statistics import median
from typing import Any, Callable

from pptx import Presentation
from pptx.util import Emu, Pt

from _common import (
    JsonlLogger,
    Theme,
    hex_to_rgb,
    load_theme,
    parse_shape_ids,
    THEMES_DIR,
)
from _shape_ops import (
    clear_fill,
    clear_line,
    clear_shadow,
    delete_shape,
    duplicate_shape,
    is_connector_like,
    move_z_order,
    nudge,
    set_corner_radius,
    set_font_bold,
    set_font_color,
    set_font_family,
    set_font_italic,
    set_font_size,
    set_line,
    set_line_spacing,
    set_opacity,
    set_position,
    set_rotation,
    set_shadow,
    set_size,
    set_solid_fill,
    set_text_align,
    set_text_content,
    set_text_margin,
    set_text_v_align,
)


LOG = JsonlLogger.from_env(component="mutate")


# ============================================================
#  Registry & helpers
# ============================================================

OP_CATALOG: list[dict[str, Any]] = []


def op(category: str, summary: str, example: str):
    """Decorator that registers a subcommand into OP_CATALOG."""

    def decorator(fn: Callable):
        OP_CATALOG.append(
            {
                "name": fn.__name__.removeprefix("cmd_").replace("_", "-"),
                "category": category,
                "summary": summary,
                "example": example,
            }
        )
        return fn

    return decorator


def _open(path: Path) -> Presentation:
    return Presentation(str(path))


def _save(prs: Presentation, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def _shapes_by_id(prs: Presentation) -> dict[int, tuple[int, Any]]:
    """Map shape_id -> (slide_index, shape)."""
    out: dict[int, tuple[int, Any]] = {}
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            sid = getattr(shape, "shape_id", None)
            if sid is not None:
                out[int(sid)] = (slide_index, shape)
    return out


def _emit(payload: dict) -> None:
    sys.stdout.write(json.dumps(payload, ensure_ascii=False) + "\n")


def _resolve_targets(prs: Presentation, args) -> list[tuple[int, int, Any]]:
    """Return [(shape_id, slide_index, shape)] for --shape-id or --shape-ids."""
    by_id = _shapes_by_id(prs)
    ids: list[int] = []
    if getattr(args, "shape_id", None) is not None:
        ids.append(int(args.shape_id))
    if getattr(args, "shape_ids", None):
        ids.extend(parse_shape_ids(args.shape_ids))
    if not ids:
        raise SystemExit("error: provide --shape-id or --shape-ids")
    out = []
    for sid in ids:
        if sid not in by_id:
            raise SystemExit(f"error: shape_id {sid} not found")
        slide_idx, shape = by_id[sid]
        out.append((sid, slide_idx, shape))
    return out


def _slide_dims(prs: Presentation) -> tuple[int, int]:
    return int(prs.slide_width), int(prs.slide_height)


# ============================================================
#  GEOMETRY OPS
# ============================================================

@op("geometry", "Set absolute position (left/top in EMU).", "mutate move --in X --out Y --shape-id 5 --left 502920 --top 228600")
def cmd_move(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_position(shape, left=args.left, top=args.top)
    _save(prs, args.out_path)
    LOG.event("move", shape_id=sid, change=change)
    _emit({"op": "move", "shape_id": sid, "change": change})
    return 0


@op("geometry", "Move shape by relative offset (dx/dy in EMU).", "mutate nudge --in X --out Y --shape-id 5 --dx 50800 --dy 0")
def cmd_nudge(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = nudge(shape, dx_emu=args.dx, dy_emu=args.dy)
    _save(prs, args.out_path)
    LOG.event("nudge", shape_id=sid, change=change)
    _emit({"op": "nudge", "shape_id": sid, "change": change})
    return 0


@op("geometry", "Set absolute size (width/height in EMU).", "mutate resize --in X --out Y --shape-id 5 --width 1828800 --height 228600")
def cmd_resize(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_size(shape, width=args.width, height=args.height)
    _save(prs, args.out_path)
    LOG.event("resize", shape_id=sid, change=change)
    _emit({"op": "resize", "shape_id": sid, "change": change})
    return 0


@op("geometry", "Align multiple shapes to a common edge (left/right/top/bottom/center-h/center-v).",
    "mutate align --in X --out Y --shape-ids 5,6,7 --edge left --target 502920")
def cmd_align(args):
    prs = _open(args.in_path)
    targets = _resolve_targets(prs, args)
    edge = args.edge
    shapes = [t[2] for t in targets]

    if args.target is not None:
        ref = int(args.target)
    else:
        # Auto: pick the most common anchor (median).
        if edge in ("left", "right", "center-h"):
            vals = [int(s.left) + (int(s.width) // 2 if edge == "center-h" else (int(s.width) if edge == "right" else 0)) for s in shapes]
        else:
            vals = [int(s.top) + (int(s.height) // 2 if edge == "center-v" else (int(s.height) if edge == "bottom" else 0)) for s in shapes]
        ref = int(median(vals))

    changes = []
    for sid, _slide, shape in targets:
        if edge == "left":
            change = set_position(shape, left=ref)
        elif edge == "right":
            change = set_position(shape, left=ref - int(shape.width))
        elif edge == "center-h":
            change = set_position(shape, left=ref - int(shape.width) // 2)
        elif edge == "top":
            change = set_position(shape, top=ref)
        elif edge == "bottom":
            change = set_position(shape, top=ref - int(shape.height))
        elif edge == "center-v":
            change = set_position(shape, top=ref - int(shape.height) // 2)
        else:
            raise SystemExit(f"error: unknown edge {edge!r}")
        changes.append({"shape_id": sid, "change": change})
    _save(prs, args.out_path)
    LOG.event("align", edge=edge, target=ref, changes=changes)
    _emit({"op": "align", "edge": edge, "target": ref, "changes": changes})
    return 0


@op("geometry", "Distribute shapes evenly along an axis (horizontal/vertical).",
    "mutate distribute --in X --out Y --shape-ids 5,6,7,8 --axis horizontal")
def cmd_distribute(args):
    prs = _open(args.in_path)
    targets = _resolve_targets(prs, args)
    if len(targets) < 3:
        raise SystemExit("error: distribute needs at least 3 shapes")
    if args.axis == "horizontal":
        ordered = sorted(targets, key=lambda t: int(t[2].left))
        first = ordered[0][2]
        last = ordered[-1][2]
        first_left = int(first.left)
        last_right = int(last.left) + int(last.width)
        total_w = sum(int(t[2].width) for t in ordered)
        gap_budget = last_right - first_left - total_w
        if gap_budget < 0:
            raise SystemExit("error: shapes overlap; cannot distribute")
        gap = gap_budget // (len(ordered) - 1)
        cursor = first_left
        changes = []
        for sid, _slide, shape in ordered:
            change = set_position(shape, left=cursor)
            changes.append({"shape_id": sid, "change": change})
            cursor += int(shape.width) + gap
    else:
        ordered = sorted(targets, key=lambda t: int(t[2].top))
        first = ordered[0][2]
        last = ordered[-1][2]
        first_top = int(first.top)
        last_bottom = int(last.top) + int(last.height)
        total_h = sum(int(t[2].height) for t in ordered)
        gap_budget = last_bottom - first_top - total_h
        if gap_budget < 0:
            raise SystemExit("error: shapes overlap; cannot distribute")
        gap = gap_budget // (len(ordered) - 1)
        cursor = first_top
        changes = []
        for sid, _slide, shape in ordered:
            change = set_position(shape, top=cursor)
            changes.append({"shape_id": sid, "change": change})
            cursor += int(shape.height) + gap
    _save(prs, args.out_path)
    LOG.event("distribute", axis=args.axis, gap=gap, changes=changes)
    _emit({"op": "distribute", "axis": args.axis, "gap_emu": gap, "changes": changes})
    return 0


@op("geometry", "Set every shape's gap to a fixed value or auto-equalized one (axis horizontal|vertical).",
    "mutate equalize-gaps --in X --out Y --shape-ids 5,6,7,8 --axis horizontal")
def cmd_equalize_gaps(args):
    return cmd_distribute(args)


@op("geometry", "Equalize widths and/or heights of multiple shapes to median.",
    "mutate equalize-size --in X --out Y --shape-ids 5,6,7 --dimension both")
def cmd_equalize_size(args):
    prs = _open(args.in_path)
    targets = _resolve_targets(prs, args)
    shapes = [t[2] for t in targets]
    target_w = int(median(int(s.width) for s in shapes)) if args.dimension in ("width", "both") else None
    target_h = int(median(int(s.height) for s in shapes)) if args.dimension in ("height", "both") else None
    changes = []
    for sid, _slide, shape in targets:
        change = set_size(shape, width=target_w, height=target_h)
        changes.append({"shape_id": sid, "change": change})
    _save(prs, args.out_path)
    LOG.event("equalize-size", changes=changes)
    _emit({"op": "equalize-size", "target_w": target_w, "target_h": target_h, "changes": changes})
    return 0


@op("geometry", "Snap shape's position to the nearest grid step.",
    "mutate snap-to-grid --in X --out Y --shape-id 5 --grid-emu 91440")
def cmd_snap_to_grid(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    g = args.grid_emu
    new_left = (int(shape.left) // g) * g
    new_top = (int(shape.top) // g) * g
    change = set_position(shape, left=new_left, top=new_top)
    _save(prs, args.out_path)
    LOG.event("snap-to-grid", shape_id=sid, grid=g, change=change)
    _emit({"op": "snap-to-grid", "shape_id": sid, "grid_emu": g, "change": change})
    return 0


@op("geometry", "Center a shape horizontally, vertically, or both.",
    "mutate center-on-slide --in X --out Y --shape-id 5 --axis horizontal")
def cmd_center_on_slide(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    sw, sh = _slide_dims(prs)
    new_left = (sw - int(shape.width)) // 2 if args.axis in ("horizontal", "both") else None
    new_top = (sh - int(shape.height)) // 2 if args.axis in ("vertical", "both") else None
    change = set_position(shape, left=new_left, top=new_top)
    _save(prs, args.out_path)
    LOG.event("center-on-slide", shape_id=sid, axis=args.axis, change=change)
    _emit({"op": "center-on-slide", "shape_id": sid, "axis": args.axis, "change": change})
    return 0


@op("geometry", "Pull a shape inside slide bounds with the given padding (no-op if already inside).",
    "mutate fit-to-slide --in X --out Y --shape-id 5 --pad-emu 457200")
def cmd_fit_to_slide(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    sw, sh = _slide_dims(prs)
    pad = args.pad_emu
    new_left = max(pad, min(int(shape.left), sw - int(shape.width) - pad))
    new_top = max(pad, min(int(shape.top), sh - int(shape.height) - pad))
    change = set_position(shape, left=new_left, top=new_top)
    _save(prs, args.out_path)
    LOG.event("fit-to-slide", shape_id=sid, change=change)
    _emit({"op": "fit-to-slide", "shape_id": sid, "change": change})
    return 0


@op("geometry", "Rotate a shape by absolute degrees.", "mutate rotate --in X --out Y --shape-id 5 --degrees 0")
def cmd_rotate(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_rotation(shape, args.degrees)
    _save(prs, args.out_path)
    LOG.event("rotate", shape_id=sid, change=change)
    _emit({"op": "rotate", "shape_id": sid, "change": change})
    return 0


# ============================================================
#  STYLING OPS
# ============================================================

@op("style", "Set solid fill color (hex RRGGBB).", "mutate set-fill --in X --out Y --shape-id 5 --color 0F62FE")
def cmd_set_fill(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_solid_fill(shape, args.color)
    _save(prs, args.out_path)
    LOG.event("set-fill", shape_id=sid, change=change)
    _emit({"op": "set-fill", "shape_id": sid, "change": change})
    return 0


@op("style", "Clear fill (transparent).", "mutate clear-fill --in X --out Y --shape-id 5")
def cmd_clear_fill(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = clear_fill(shape)
    _save(prs, args.out_path)
    LOG.event("clear-fill", shape_id=sid, change=change)
    _emit({"op": "clear-fill", "shape_id": sid, "change": change})
    return 0


@op("style", "Set border color and/or width (pt).", "mutate set-line --in X --out Y --shape-id 5 --color DDE1E6 --width-pt 0.75")
def cmd_set_line(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_line(shape, hex_color=args.color, width_pt=args.width_pt)
    _save(prs, args.out_path)
    LOG.event("set-line", shape_id=sid, change=change)
    _emit({"op": "set-line", "shape_id": sid, "change": change})
    return 0


@op("style", "Remove border.", "mutate clear-line --in X --out Y --shape-id 5")
def cmd_clear_line(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = clear_line(shape)
    _save(prs, args.out_path)
    LOG.event("clear-line", shape_id=sid, change=change)
    _emit({"op": "clear-line", "shape_id": sid, "change": change})
    return 0


@op("style", "Add an outer drop shadow.", "mutate set-shadow --in X --out Y --shape-id 5 --color 888888 --blur-pt 4 --dist-pt 2 --alpha 0.3")
def cmd_set_shadow(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_shadow(shape, hex_color=args.color, blur_pt=args.blur_pt, dist_pt=args.dist_pt, alpha=args.alpha)
    _save(prs, args.out_path)
    LOG.event("set-shadow", shape_id=sid, change=change)
    _emit({"op": "set-shadow", "shape_id": sid, "change": change})
    return 0


@op("style", "Remove drop shadow.", "mutate clear-shadow --in X --out Y --shape-id 5")
def cmd_clear_shadow(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = clear_shadow(shape)
    _save(prs, args.out_path)
    LOG.event("clear-shadow", shape_id=sid, change=change)
    _emit({"op": "clear-shadow", "shape_id": sid, "change": change})
    return 0


@op("style", "Round corners on a rounded-rectangle (ratio 0.0-0.5).",
    "mutate set-corner-radius --in X --out Y --shape-id 5 --ratio 0.08")
def cmd_set_corner_radius(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_corner_radius(shape, args.ratio)
    _save(prs, args.out_path)
    LOG.event("set-corner-radius", shape_id=sid, change=change)
    _emit({"op": "set-corner-radius", "shape_id": sid, "change": change})
    return 0


@op("style", "Set fill alpha (0.0-1.0).", "mutate set-opacity --in X --out Y --shape-id 5 --alpha 0.6")
def cmd_set_opacity(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_opacity(shape, args.alpha)
    _save(prs, args.out_path)
    LOG.event("set-opacity", shape_id=sid, change=change)
    _emit({"op": "set-opacity", "shape_id": sid, "change": change})
    return 0


@op("style", "Apply theme card styling (fill + border + corner) in one call.",
    "mutate apply-card-style --in X --out Y --shape-id 5 [--theme path]")
def cmd_apply_card_style(args):
    prs = _open(args.in_path)
    theme = load_theme(args.theme)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    deco = theme.decoration
    fill_role = deco.get("card_fill_role", "background")
    border_role = deco.get("card_border_role", "border")
    border_pt = deco.get("card_border_pt", 0.75)
    corner = deco.get("card_corner_radius_ratio")
    changes = {}
    changes["fill"] = set_solid_fill(shape, theme.color(fill_role))
    changes["line"] = set_line(shape, hex_color=theme.color(border_role), width_pt=border_pt)
    if corner is not None:
        changes["corner"] = set_corner_radius(shape, corner)
    if deco.get("card_shadow"):
        changes["shadow"] = set_shadow(shape)
    _save(prs, args.out_path)
    LOG.event("apply-card-style", shape_id=sid, changes=changes)
    _emit({"op": "apply-card-style", "shape_id": sid, "changes": changes})
    return 0


@op("style", "Apply theme badge styling (primary fill + white text).",
    "mutate apply-badge-style --in X --out Y --shape-id 5 [--theme path]")
def cmd_apply_badge_style(args):
    prs = _open(args.in_path)
    theme = load_theme(args.theme)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    fill_role = theme.decoration.get("badge_fill_role", "primary")
    changes = {
        "fill": set_solid_fill(shape, theme.color(fill_role)),
        "line": set_line(shape, hex_color=theme.color(fill_role), width_pt=0.0),
        "font_size": set_font_size(shape, theme.font_size_pt("badge") or 9),
        "font_bold": set_font_bold(shape, bool(theme.font_bold("badge"))),
        "font_color": set_font_color(shape, theme.color(theme.font_color_role("badge") or "background")),
    }
    _save(prs, args.out_path)
    LOG.event("apply-badge-style", shape_id=sid, changes=changes)
    _emit({"op": "apply-badge-style", "shape_id": sid, "changes": changes})
    return 0


# ============================================================
#  Z-ORDER OPS
# ============================================================

@op("z-order", "Move shape(s) back/front/up/down in z-order.",
    "mutate z-order --in X --out Y --shape-ids 10,11 --position back")
def cmd_z_order(args):
    prs = _open(args.in_path)
    targets = _resolve_targets(prs, args)
    changes = []
    for sid, _slide, shape in targets:
        change = move_z_order(shape, args.position)
        changes.append({"shape_id": sid, "change": change})
    _save(prs, args.out_path)
    LOG.event("z-order", position=args.position, changes=changes)
    _emit({"op": "z-order", "position": args.position, "changes": changes})
    return 0


@op("z-order", "Send shape(s) to back. Shortcut for z-order --position back.",
    "mutate send-to-back --in X --out Y --shape-ids 10,11")
def cmd_send_to_back(args):
    args.position = "back"
    return cmd_z_order(args)


@op("z-order", "Bring shape(s) to front. Shortcut for z-order --position front.",
    "mutate bring-to-front --in X --out Y --shape-ids 5")
def cmd_bring_to_front(args):
    args.position = "front"
    return cmd_z_order(args)


@op("z-order", "Send every connector-like shape (lines, h=0 or w=0) to the back of every slide.",
    "mutate all-connectors-to-back --in X --out Y")
def cmd_all_connectors_to_back(args):
    prs = _open(args.in_path)
    moved = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if is_connector_like(shape):
                if move_z_order(shape, "back") is not None:
                    moved.append({"slide_index": slide_index, "shape_id": int(shape.shape_id)})
    _save(prs, args.out_path)
    LOG.event("all-connectors-to-back", moved=moved, count=len(moved))
    _emit({"op": "all-connectors-to-back", "moved_count": len(moved), "moved": moved})
    return 0


# ============================================================
#  TEXT OPS
# ============================================================

@op("text", "Force a single font family on every text shape in the deck (Latin + East Asian; recurses into groups). Default family: 'Microsoft YaHei'.",
    "mutate unify-font --in X --out Y [--family \"Microsoft YaHei\"]")
def cmd_unify_font(args):
    family = getattr(args, "family", None) or "Microsoft YaHei"
    prs = _open(args.in_path)
    affected: list[int] = []
    skipped_groups = 0

    def _walk_text_shapes(container):
        """Yield every text-bearing shape including those inside groups."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for sp in container.shapes:
            st = getattr(sp, "shape_type", None)
            if st == MSO_SHAPE_TYPE.GROUP:
                yield from _walk_text_shapes(sp)
            elif getattr(sp, "has_text_frame", False):
                yield sp

    for slide in prs.slides:
        for shape in _walk_text_shapes(slide):
            change = set_font_family(shape, family, include_eastasia=True)
            if change:
                affected.append(int(getattr(shape, "shape_id", 0) or 0))

    _save(prs, args.out_path)
    LOG.event("unify-font", family=family, count=len(affected))
    _emit({"op": "unify-font", "family": family, "affected_count": len(affected), "shape_ids": affected})
    return 0


@op("text", "Set font family for one shape, all shapes (--scope all), or a role (--scope role:title). For deck-wide CJK use `unify-font` which is simpler.",
    "mutate set-font-family --in X --out Y --scope all --family \"Microsoft YaHei\"")
def cmd_set_font_family(args):
    prs = _open(args.in_path)
    family = args.family
    affected = []
    if args.shape_id is not None or args.shape_ids:
        targets = _resolve_targets(prs, args)
        for sid, _slide, shape in targets:
            if set_font_family(shape, family):
                affected.append(sid)
    else:
        scope = args.scope or "all"
        for slide in prs.slides:
            for shape in slide.shapes:
                if scope == "all":
                    if set_font_family(shape, family):
                        affected.append(int(shape.shape_id))
                elif scope.startswith("role:"):
                    # Need roles JSON for this; punt to the orchestrator.
                    raise SystemExit("error: scope=role:* requires --roles; use mutate.py via orchestrator")
    _save(prs, args.out_path)
    LOG.event("set-font-family", family=family, count=len(affected))
    _emit({"op": "set-font-family", "family": family, "affected_count": len(affected), "shape_ids": affected})
    return 0


@op("text", "Set font size in pt for a shape.", "mutate set-font-size --in X --out Y --shape-id 5 --size-pt 14")
def cmd_set_font_size(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_font_size(shape, args.size_pt)
    _save(prs, args.out_path)
    LOG.event("set-font-size", shape_id=sid, change=change)
    _emit({"op": "set-font-size", "shape_id": sid, "change": change})
    return 0


@op("text", "Toggle bold.", "mutate set-font-bold --in X --out Y --shape-id 5 --bold true")
def cmd_set_font_bold(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_font_bold(shape, args.bold)
    _save(prs, args.out_path)
    LOG.event("set-font-bold", shape_id=sid, change=change)
    _emit({"op": "set-font-bold", "shape_id": sid, "change": change})
    return 0


@op("text", "Toggle italic.", "mutate set-font-italic --in X --out Y --shape-id 5 --italic false")
def cmd_set_font_italic(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_font_italic(shape, args.italic)
    _save(prs, args.out_path)
    LOG.event("set-font-italic", shape_id=sid, change=change)
    _emit({"op": "set-font-italic", "shape_id": sid, "change": change})
    return 0


@op("text", "Set font color (hex RRGGBB).", "mutate set-font-color --in X --out Y --shape-id 5 --color 161616")
def cmd_set_font_color(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_font_color(shape, args.color)
    _save(prs, args.out_path)
    LOG.event("set-font-color", shape_id=sid, change=change)
    _emit({"op": "set-font-color", "shape_id": sid, "change": change})
    return 0


@op("text", "Horizontal text alignment.", "mutate set-text-align --in X --out Y --shape-id 5 --align left")
def cmd_set_text_align(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_text_align(shape, args.align)
    _save(prs, args.out_path)
    LOG.event("set-text-align", shape_id=sid, change=change)
    _emit({"op": "set-text-align", "shape_id": sid, "change": change})
    return 0


@op("text", "Vertical anchor inside the text frame.",
    "mutate set-text-v-align --in X --out Y --shape-id 5 --anchor middle")
def cmd_set_text_v_align(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_text_v_align(shape, args.anchor)
    _save(prs, args.out_path)
    LOG.event("set-text-v-align", shape_id=sid, change=change)
    _emit({"op": "set-text-v-align", "shape_id": sid, "change": change})
    return 0


@op("text", "Set text-frame inner margins (EMU each).",
    "mutate set-text-margin --in X --out Y --shape-id 5 --left-emu 91440 --right-emu 91440")
def cmd_set_text_margin(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_text_margin(shape, left_emu=args.left_emu, right_emu=args.right_emu, top_emu=args.top_emu, bottom_emu=args.bottom_emu)
    _save(prs, args.out_path)
    LOG.event("set-text-margin", shape_id=sid, change=change)
    _emit({"op": "set-text-margin", "shape_id": sid, "change": change})
    return 0


@op("text", "Set line spacing ratio (1.0 = single).", "mutate set-line-spacing --in X --out Y --shape-id 5 --ratio 1.25")
def cmd_set_line_spacing(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_line_spacing(shape, args.ratio)
    _save(prs, args.out_path)
    LOG.event("set-line-spacing", shape_id=sid, change=change)
    _emit({"op": "set-line-spacing", "shape_id": sid, "change": change})
    return 0


@op("text", "Replace the text content of a shape.",
    "mutate set-text --in X --out Y --shape-id 5 --content \"New text\"")
def cmd_set_text(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_text_content(shape, args.content)
    _save(prs, args.out_path)
    LOG.event("set-text", shape_id=sid, change=change)
    _emit({"op": "set-text", "shape_id": sid, "change": change})
    return 0


@op("text", "Apply a theme typography role (size + bold + color + family) to a shape.",
    "mutate apply-typography --in X --out Y --shape-id 5 --role title [--theme path]")
def cmd_apply_typography(args):
    prs = _open(args.in_path)
    theme = load_theme(args.theme)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    role = args.role
    changes = {}
    if theme.font_family():
        changes["family"] = set_font_family(shape, theme.font_family())
    if theme.font_size_pt(role) is not None:
        changes["size_pt"] = set_font_size(shape, theme.font_size_pt(role))
    if theme.font_bold(role) is not None:
        changes["bold"] = set_font_bold(shape, bool(theme.font_bold(role)))
    color_role = theme.font_color_role(role)
    if color_role:
        changes["color"] = set_font_color(shape, theme.color(color_role))
    _save(prs, args.out_path)
    LOG.event("apply-typography", shape_id=sid, role=role, changes=changes)
    _emit({"op": "apply-typography", "shape_id": sid, "role": role, "changes": changes})
    return 0


# ============================================================
#  SLIDE-LEVEL OPS
# ============================================================

@op("slide", "Delete a shape from its slide.", "mutate delete-shape --in X --out Y --shape-id 33")
def cmd_delete_shape(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = delete_shape(shape)
    _save(prs, args.out_path)
    LOG.event("delete-shape", shape_id=sid, change=change)
    _emit({"op": "delete-shape", "shape_id": sid, "change": change})
    return 0


@op("slide", "Duplicate a shape (new copy appended to slide, returns new shape_id).",
    "mutate duplicate-shape --in X --out Y --shape-id 5")
def cmd_duplicate_shape(args):
    prs = _open(args.in_path)
    by_id = _shapes_by_id(prs)
    sid = int(args.shape_id)
    if sid not in by_id:
        raise SystemExit(f"error: shape_id {sid} not found")
    slide_idx, shape = by_id[sid]
    new_shape, info = duplicate_shape(prs.slides[slide_idx - 1], shape)
    _save(prs, args.out_path)
    LOG.event("duplicate-shape", shape_id=sid, new_shape_id=int(new_shape.shape_id))
    _emit({"op": "duplicate-shape", "shape_id": sid, "new_shape_id": int(new_shape.shape_id), "info": info})
    return 0


@op("slide", "Add a rectangle on a slide.",
    "mutate add-rect --in X --out Y --slide 1 --left 0 --top 0 --width 914400 --height 91440 --fill 0F62FE")
def cmd_add_rect(args):
    from pptx.enum.shapes import MSO_SHAPE

    prs = _open(args.in_path)
    slide = prs.slides[args.slide - 1]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(args.left),
        Emu(args.top),
        Emu(args.width),
        Emu(args.height),
    )
    if args.fill:
        set_solid_fill(shape, args.fill)
    if args.line:
        set_line(shape, hex_color=args.line, width_pt=args.line_width_pt or 0.75)
    else:
        clear_line(shape)
    _save(prs, args.out_path)
    LOG.event("add-rect", slide=args.slide, shape_id=int(shape.shape_id))
    _emit({"op": "add-rect", "slide": args.slide, "new_shape_id": int(shape.shape_id)})
    return 0


@op("slide", "Add a connector line between two points.",
    "mutate add-line --in X --out Y --slide 1 --x1 0 --y1 0 --x2 914400 --y2 0 --color 6F6F6F --width-pt 1")
def cmd_add_line(args):
    from pptx.enum.shapes import MSO_CONNECTOR

    prs = _open(args.in_path)
    slide = prs.slides[args.slide - 1]
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(args.x1),
        Emu(args.y1),
        Emu(args.x2),
        Emu(args.y2),
    )
    if args.color:
        set_line(conn, hex_color=args.color, width_pt=args.width_pt or 1.0)
    _save(prs, args.out_path)
    LOG.event("add-line", slide=args.slide, shape_id=int(conn.shape_id))
    _emit({"op": "add-line", "slide": args.slide, "new_shape_id": int(conn.shape_id)})
    return 0


@op("slide", "Add a text box.",
    "mutate add-text --in X --out Y --slide 1 --left 457200 --top 457200 --width 8000000 --height 400000 --content \"Title\"")
def cmd_add_text(args):
    prs = _open(args.in_path)
    slide = prs.slides[args.slide - 1]
    box = slide.shapes.add_textbox(
        Emu(args.left), Emu(args.top), Emu(args.width), Emu(args.height)
    )
    box.text_frame.text = args.content
    if args.role and args.theme is not None:
        theme = load_theme(args.theme)
        if theme.font_family():
            set_font_family(box, theme.font_family())
        if theme.font_size_pt(args.role) is not None:
            set_font_size(box, theme.font_size_pt(args.role))
        if theme.font_bold(args.role) is not None:
            set_font_bold(box, bool(theme.font_bold(args.role)))
        color_role = theme.font_color_role(args.role)
        if color_role:
            set_font_color(box, theme.color(color_role))
    _save(prs, args.out_path)
    LOG.event("add-text", slide=args.slide, shape_id=int(box.shape_id), content=args.content[:60])
    _emit({"op": "add-text", "slide": args.slide, "new_shape_id": int(box.shape_id)})
    return 0


# ============================================================
#  MODEL-FRIENDLY PLACEMENT OPS
#
#  These ops let an OpenCode agent move shapes by REFERRING to other
#  shapes/cards instead of computing EMU coordinates. Agents are weak at
#  EMU arithmetic but strong at "this icon belongs in card #5 just like
#  the icon in card #1".
# ============================================================

@op("placement", "Move shape NEAR another shape (reference). Anchor: above|below|left-of|right-of|inside-top-left|inside-center|inside-bottom-right.",
    "mutate place-near --in X --out Y --shape-id 38 --reference-shape-id 47 --anchor inside-top-left --gap-emu 91440")
def cmd_place_near(args):
    """Place shape `--shape-id` near `--reference-shape-id` using a named anchor."""
    prs = _open(args.in_path)
    by_id = _shapes_by_id(prs)
    sid = int(args.shape_id)
    ref_sid = int(args.reference_shape_id)
    if sid not in by_id:
        raise SystemExit(f"error: shape_id {sid} not found")
    if ref_sid not in by_id:
        raise SystemExit(f"error: reference shape_id {ref_sid} not found")
    _, shape = by_id[sid]
    _, ref = by_id[ref_sid]

    rL, rT = int(ref.left), int(ref.top)
    rW, rH = int(ref.width), int(ref.height)
    sW, sH = int(shape.width), int(shape.height)
    gap = int(getattr(args, "gap_emu", 0) or 0)
    anchor = args.anchor

    if anchor == "above":
        new_left = rL + (rW - sW) // 2
        new_top = rT - sH - gap
    elif anchor == "below":
        new_left = rL + (rW - sW) // 2
        new_top = rT + rH + gap
    elif anchor == "left-of":
        new_left = rL - sW - gap
        new_top = rT + (rH - sH) // 2
    elif anchor == "right-of":
        new_left = rL + rW + gap
        new_top = rT + (rH - sH) // 2
    elif anchor == "inside-top-left":
        new_left = rL + gap
        new_top = rT + gap
    elif anchor == "inside-top-right":
        new_left = rL + rW - sW - gap
        new_top = rT + gap
    elif anchor == "inside-center":
        new_left = rL + (rW - sW) // 2
        new_top = rT + (rH - sH) // 2
    elif anchor == "inside-bottom-left":
        new_left = rL + gap
        new_top = rT + rH - sH - gap
    elif anchor == "inside-bottom-right":
        new_left = rL + rW - sW - gap
        new_top = rT + rH - sH - gap
    else:
        raise SystemExit(f"error: unsupported anchor {anchor!r}")

    change = set_position(shape, left=new_left, top=new_top)
    _save(prs, args.out_path)
    LOG.event("place-near", shape_id=sid, reference=ref_sid, anchor=anchor, change=change)
    _emit({"op": "place-near", "shape_id": sid, "reference": ref_sid, "anchor": anchor, "change": change})
    return 0


@op("placement", "Mirror a peer's relative position: place shape X at the same offset within its container that --peer-shape-id has within --peer-container-id.",
    "mutate mirror-peer-position --in X --out Y --shape-id 38 --target-container-id 49 --peer-shape-id 53 --peer-container-id 64")
def cmd_mirror_peer_position(args):
    """Place shape so that its offset relative to target_container matches
    peer_shape's offset relative to peer_container. Useful when one card has
    a misplaced icon and you want to copy a sibling card's icon position."""
    prs = _open(args.in_path)
    by_id = _shapes_by_id(prs)
    sid = int(args.shape_id)
    target_cid = int(args.target_container_id)
    peer_sid = int(args.peer_shape_id)
    peer_cid = int(args.peer_container_id)
    for k, name in ((sid, "shape_id"), (target_cid, "target_container_id"),
                    (peer_sid, "peer_shape_id"), (peer_cid, "peer_container_id")):
        if k not in by_id:
            raise SystemExit(f"error: {name}={k} not found")
    _, shape = by_id[sid]
    _, target_card = by_id[target_cid]
    _, peer_shape = by_id[peer_sid]
    _, peer_card = by_id[peer_cid]

    rel_x = int(peer_shape.left) - int(peer_card.left)
    rel_y = int(peer_shape.top) - int(peer_card.top)
    new_left = int(target_card.left) + rel_x
    new_top = int(target_card.top) + rel_y
    change = set_position(shape, left=new_left, top=new_top)
    _save(prs, args.out_path)
    LOG.event("mirror-peer-position", shape_id=sid, target_card=target_cid,
              peer_shape=peer_sid, peer_card=peer_cid, change=change)
    _emit({"op": "mirror-peer-position", "shape_id": sid, "target_card": target_cid,
           "peer_shape": peer_sid, "peer_card": peer_cid, "change": change})
    return 0


@op("placement", "Snap a grid-cell shape: take its top from one peer, its left from another. Useful when ONE cell of a 2D grid (NxM rows/cols) is misaligned and other peers are correct.",
    "mutate align-to-grid-cell --in X --out Y --shape-id 6 --top-from-shape-id 21 --left-from-shape-id 36")
def cmd_align_to_grid_cell(args):
    """Snap shape's top to peer-A's top and/or left to peer-B's left."""
    prs = _open(args.in_path)
    by_id = _shapes_by_id(prs)
    sid = int(args.shape_id)
    if sid not in by_id:
        raise SystemExit(f"error: shape_id {sid} not found")
    _, shape = by_id[sid]
    new_left = None
    new_top = None
    if args.left_from_shape_id is not None:
        ref_sid = int(args.left_from_shape_id)
        if ref_sid not in by_id:
            raise SystemExit(f"error: --left-from-shape-id {ref_sid} not found")
        new_left = int(by_id[ref_sid][1].left)
    if args.top_from_shape_id is not None:
        ref_sid = int(args.top_from_shape_id)
        if ref_sid not in by_id:
            raise SystemExit(f"error: --top-from-shape-id {ref_sid} not found")
        new_top = int(by_id[ref_sid][1].top)
    if new_left is None and new_top is None:
        raise SystemExit("error: provide --top-from-shape-id and/or --left-from-shape-id")
    change = set_position(shape, left=new_left, top=new_top)
    _save(prs, args.out_path)
    LOG.event("align-to-grid-cell", shape_id=sid, change=change)
    _emit({"op": "align-to-grid-cell", "shape_id": sid, "change": change})
    return 0


@op("placement", "Move shape into the bbox of another shape, optionally with a relative offset (0-1 fractional).",
    "mutate move-to-card --in X --out Y --shape-id 38 --container-id 49 --rel-x 0.15 --rel-y 0.10")
def cmd_move_to_card(args):
    """Move shape into a target container at fractional rel position."""
    prs = _open(args.in_path)
    by_id = _shapes_by_id(prs)
    sid = int(args.shape_id)
    cid = int(args.container_id)
    if sid not in by_id:
        raise SystemExit(f"error: shape_id {sid} not found")
    if cid not in by_id:
        raise SystemExit(f"error: container_id {cid} not found")
    _, shape = by_id[sid]
    _, card = by_id[cid]

    rx = float(getattr(args, "rel_x", None) or 0.10)
    ry = float(getattr(args, "rel_y", None) or 0.10)
    if not (0.0 <= rx <= 1.0 and 0.0 <= ry <= 1.0):
        raise SystemExit("error: --rel-x and --rel-y must be in [0, 1]")
    new_left = int(card.left) + int(int(card.width) * rx)
    new_top = int(card.top) + int(int(card.height) * ry)
    change = set_position(shape, left=new_left, top=new_top)
    _save(prs, args.out_path)
    LOG.event("move-to-card", shape_id=sid, container=cid, rel=(rx, ry), change=change)
    _emit({"op": "move-to-card", "shape_id": sid, "container": cid, "rel": [rx, ry], "change": change})
    return 0


# ============================================================
#  CARD REPAIR (peer-card outlier detection)
# ============================================================

@op("repair", "Detect MxN grid layouts (2x3 / 3x2 / NxM dashboards) and fix outlier panels. Recurses into each panel for nested sub-grids when --nested.",
    "mutate repair-grid --in X --out Y [--nested]")
def cmd_repair_grid(args):
    """Pure 2D-grid path. For 1D row layouts use repair-peer-cards."""
    import subprocess, tempfile
    from pathlib import Path as _P

    prs = _open(args.in_path)
    script_dir = _P(__file__).resolve().parent
    nested = bool(getattr(args, "nested", False))
    max_depth = 2 if nested else 1

    with tempfile.TemporaryDirectory() as tmp:
        ins = _P(tmp) / "ins.json"
        subprocess.run(
            [sys.executable, str(script_dir / "inspect_ppt.py"),
             "--input", str(args.in_path), "--output", str(ins)],
            check=True,
        )
        inspection = json.loads(ins.read_text(encoding="utf-8"))

    from _grid_detect import diagnose_grid_repair, detect_grids_nested
    from _card_repair import apply_repair

    actions: list[dict] = []
    summary_grids = []
    for slide_idx, (slide, slide_ins) in enumerate(zip(prs.slides, inspection["slides"]), start=1):
        plan = diagnose_grid_repair(slide_ins, max_depth=max_depth)
        if plan["rows"]:
            apply_repair(slide, slide_ins, plan, actions, slide_idx)
        # Also probe what was detected for the summary
        grids = detect_grids_nested(slide_ins["objects"], max_depth=max_depth)
        for g in grids:
            summary_grids.append({
                "slide_index": slide_idx,
                "depth": g.get("depth", 0),
                "rows": g["rows"], "cols": g["cols"],
                "fill_ratio": g["fill_ratio"],
                "panel_count": len(g["panel_ids"]),
                "outlier_count": len(g["outliers"]),
                "parent_panel_id": g.get("parent_panel_id"),
            })

    _save(prs, args.out_path)
    LOG.event("repair-grid", actions=len(actions), grids=summary_grids, nested=nested)
    _emit({"op": "repair-grid", "nested": nested, "actions_applied": len(actions),
           "grids": summary_grids, "actions": actions})
    return 0


@op("polish",
    "Apply business-grade visual polish without changing content: smart "
    "typography (theme type-scale across detected roles), unified corner "
    "radius and shadow, accent bar above title, footer divider, optional "
    "subtle subtitle background tint. Idempotent: re-runs detect existing "
    "decoration shapes and skip them. --level 1 (subtle) | 2 (standard, "
    "default) | 3 (rich). Auto-picks theme from content if --theme omitted.",
    "mutate polish-business --in X --out Y [--level 2] [--theme path]")
def cmd_polish_business(args):
    """Run the business-polish pass on the deck."""
    import subprocess, tempfile
    from pathlib import Path as _P

    script_dir = _P(__file__).resolve().parent
    prs = _open(args.in_path)

    # Inspect + role-detect for typography + decoration anchoring.
    with tempfile.TemporaryDirectory() as tmp:
        ins_path = _P(tmp) / "ins.json"
        roles_path = _P(tmp) / "roles.json"
        subprocess.run(
            [sys.executable, str(script_dir / "inspect_ppt.py"),
             "--input", str(args.in_path), "--output", str(ins_path)],
            check=True,
        )
        subprocess.run(
            [sys.executable, str(script_dir / "detect_roles.py"),
             "--inspection", str(ins_path), "--output", str(roles_path)],
            check=True,
        )
        inspection = json.loads(ins_path.read_text(encoding="utf-8"))
        role_data = json.loads(roles_path.read_text(encoding="utf-8"))

    # Theme: explicit --theme wins; otherwise auto-pick from content.
    from _business_polish import (
        apply_business_polish, pick_theme,
    )
    if args.theme:
        theme = load_theme(_P(args.theme))
        theme_source = str(args.theme)
    else:
        # Aggregate all text from inspection for theme keyword scoring.
        all_text = " ".join(
            obj.get("text") or ""
            for slide in inspection.get("slides", [])
            for obj in slide.get("objects", [])
        )
        theme_name = pick_theme(all_text, THEMES_DIR)
        theme = load_theme(THEMES_DIR / f"{theme_name}.json")
        theme_source = f"auto:{theme_name}"

    level = int(args.level)
    result = apply_business_polish(prs, theme, role_data, level=level)
    _save(prs, args.out_path)

    LOG.event("polish-business",
              level=level, theme_source=theme_source,
              actions=result["actions_applied"])
    _emit({
        "op": "polish-business",
        "level": level,
        "theme_source": theme_source,
        "actions_applied": result["actions_applied"],
        "actions": result["actions"],
    })
    return 0


@op("repair",
    "Apply size/spacing/alignment uniformity within agent-defined peer groups. "
    "Reads a peer-groups.json the agent wrote after looking at the render. "
    "Children move with their parent (icons + text inside cards stay attached). "
    "Use this when geometric clustering (repair-grid / repair-peer-cards) "
    "would mis-group semantically-distinct shapes that happen to look similar.",
    "mutate repair-peers-smart --in X --out Y --groups peer-groups.json")
def cmd_repair_peers_smart(args):
    from _repair_peers_smart import repair_peers_smart
    prs = _open(args.in_path)
    groups = json.loads(Path(args.groups).read_text(encoding="utf-8"))
    result = repair_peers_smart(prs, groups)
    _save(prs, args.out_path)
    LOG.event("repair-peers-smart",
              groups=result["groups_processed"],
              actions=result["actions_applied"],
              skipped=len(result["skipped"]))
    _emit({
        "op": "repair-peers-smart",
        "groups_processed": result["groups_processed"],
        "actions_applied": result["actions_applied"],
        "actions": result["actions"],
        "skipped": result["skipped"],
    })
    return 0


@op("repair", "Detect and fix peer-card outliers. --scope safe (only box+header) | no-orphans | all (default).",
    "mutate repair-peer-cards --in X --out Y [--scope safe]")
def cmd_repair_peer_cards(args):
    """Use _card_repair to diagnose and apply fixes on every slide.

    For heavily-damaged decks where orphan/displaced relocations cascade
    into wrong assignments, pass --scope safe to apply only the well-
    constrained box-fix and header-strip-fix.
    """
    import subprocess, tempfile
    from pathlib import Path as _P

    scope = getattr(args, "scope", None) or "all"
    if scope not in ("all", "safe", "no-orphans"):
        raise SystemExit(f"error: --scope must be all|safe|no-orphans, got {scope!r}")

    prs = _open(args.in_path)
    script_dir = _P(__file__).resolve().parent
    with tempfile.TemporaryDirectory() as tmp:
        ins = _P(tmp) / "ins.json"
        subprocess.run(
            [sys.executable, str(script_dir / "inspect_ppt.py"),
             "--input", str(args.in_path), "--output", str(ins)],
            check=True,
        )
        inspection = json.loads(ins.read_text(encoding="utf-8"))

    from _card_repair import diagnose_repair, apply_repair
    actions: list[dict] = []
    summary_rows = []
    for slide_idx, (slide, slide_ins) in enumerate(zip(prs.slides, inspection["slides"]), start=1):
        plan = diagnose_repair(slide_ins, scope=scope)
        apply_repair(slide, slide_ins, plan, actions, slide_idx)
        for row in plan["rows"]:
            summary_rows.append({
                "slide_index": slide_idx,
                "card_shape_ids": row["card_shape_ids"],
                "box_fixes": len(row["card_box_fixes"]),
                "header_fixes": len(row["header_strip_fixes"]),
                "orphan_relocations": len(row["orphan_relocations"]),
                "displaced_relocations": len(row.get("displaced_relocations", [])),
            })

    _save(prs, args.out_path)
    LOG.event("repair-peer-cards", actions=len(actions), summary=summary_rows, scope=scope)
    _emit({"op": "repair-peer-cards", "scope": scope, "actions_applied": len(actions), "summary": summary_rows, "actions": actions})
    return 0


# ============================================================
#  CONNECTOR OPS
# ============================================================

@op("connector", "Apply theme connector style (color + width) to one shape.",
    "mutate style-connector --in X --out Y --shape-id 10 [--theme path]")
def cmd_style_connector(args):
    prs = _open(args.in_path)
    theme = load_theme(args.theme)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    color = theme.color(theme.decoration.get("connector_color_role", "text_muted"))
    width_pt = theme.decoration.get("connector_pt", 1.0)
    change = set_line(shape, hex_color=color, width_pt=width_pt)
    _save(prs, args.out_path)
    LOG.event("style-connector", shape_id=sid, change=change)
    _emit({"op": "style-connector", "shape_id": sid, "change": change})
    return 0


@op("connector", "Force connector to be horizontal (h=0) or vertical (w=0).",
    "mutate straighten-connector --in X --out Y --shape-id 10 --axis horizontal")
def cmd_straighten_connector(args):
    prs = _open(args.in_path)
    [(sid, _, shape)] = _resolve_targets(prs, args)
    change = set_size(shape, height=0) if args.axis == "horizontal" else set_size(shape, width=0)
    _save(prs, args.out_path)
    LOG.event("straighten-connector", shape_id=sid, axis=args.axis, change=change)
    _emit({"op": "straighten-connector", "shape_id": sid, "axis": args.axis, "change": change})
    return 0


# ============================================================
#  DISCOVERY
# ============================================================

@op("meta", "List every available op (use --json for machine-readable).", "mutate list-ops --json")
def cmd_list_ops(args):
    if args.json:
        _emit({"ops": OP_CATALOG})
    else:
        by_cat: dict[str, list[dict]] = {}
        for op_meta in OP_CATALOG:
            by_cat.setdefault(op_meta["category"], []).append(op_meta)
        for cat in sorted(by_cat):
            print(f"\n[{cat}]")
            for o in by_cat[cat]:
                print(f"  {o['name']:24s} {o['summary']}")
                print(f"    e.g. {o['example']}")
    return 0


# ============================================================
#  Argparse wiring
# ============================================================

def _add_io(p: argparse.ArgumentParser) -> None:
    p.add_argument("--in", dest="in_path", required=True, type=Path, help="input .pptx")
    p.add_argument("--out", dest="out_path", required=True, type=Path, help="output .pptx")


def _add_targets(p: argparse.ArgumentParser, *, allow_single: bool = True, allow_multi: bool = True) -> None:
    if allow_single:
        p.add_argument("--shape-id", type=int)
    if allow_multi:
        p.add_argument("--shape-ids", help="comma- or space-separated shape ids")


def _add_theme(p: argparse.ArgumentParser) -> None:
    p.add_argument("--theme", type=Path, default=THEMES_DIR / "clean-tech.json")


def _str_to_bool(s: str) -> bool:
    return str(s).lower() in ("1", "true", "yes", "y", "on")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog="mutate", description="Granular pptx mutation operations.")
    sub = parser.add_subparsers(dest="cmd", required=True)

    # Geometry
    p = sub.add_parser("move", help="set absolute position"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--left", type=int); p.add_argument("--top", type=int); p.set_defaults(func=cmd_move)
    p = sub.add_parser("nudge", help="relative offset"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--dx", type=int, default=0); p.add_argument("--dy", type=int, default=0); p.set_defaults(func=cmd_nudge)
    p = sub.add_parser("resize", help="absolute size"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--width", type=int); p.add_argument("--height", type=int); p.set_defaults(func=cmd_resize)
    p = sub.add_parser("align", help="align multiple shapes"); _add_io(p); _add_targets(p, allow_single=False); p.add_argument("--edge", choices=["left", "right", "top", "bottom", "center-h", "center-v"], required=True); p.add_argument("--target", type=int); p.set_defaults(func=cmd_align)
    p = sub.add_parser("distribute", help="evenly distribute shapes"); _add_io(p); _add_targets(p, allow_single=False); p.add_argument("--axis", choices=["horizontal", "vertical"], required=True); p.set_defaults(func=cmd_distribute)
    p = sub.add_parser("equalize-gaps", help="alias of distribute"); _add_io(p); _add_targets(p, allow_single=False); p.add_argument("--axis", choices=["horizontal", "vertical"], required=True); p.set_defaults(func=cmd_equalize_gaps)
    p = sub.add_parser("equalize-size", help="equalize widths/heights to median"); _add_io(p); _add_targets(p, allow_single=False); p.add_argument("--dimension", choices=["width", "height", "both"], default="both"); p.set_defaults(func=cmd_equalize_size)
    p = sub.add_parser("snap-to-grid", help="snap to grid"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--grid-emu", type=int, default=91440); p.set_defaults(func=cmd_snap_to_grid)
    p = sub.add_parser("center-on-slide", help="center on slide"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--axis", choices=["horizontal", "vertical", "both"], default="both"); p.set_defaults(func=cmd_center_on_slide)
    p = sub.add_parser("fit-to-slide", help="pull inside bounds"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--pad-emu", type=int, default=457200); p.set_defaults(func=cmd_fit_to_slide)
    p = sub.add_parser("rotate", help="set rotation degrees"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--degrees", type=float, required=True); p.set_defaults(func=cmd_rotate)

    # Style
    p = sub.add_parser("set-fill", help="solid fill"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--color", required=True); p.set_defaults(func=cmd_set_fill)
    p = sub.add_parser("clear-fill", help="clear fill"); _add_io(p); _add_targets(p, allow_multi=False); p.set_defaults(func=cmd_clear_fill)
    p = sub.add_parser("set-line", help="border color/width"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--color"); p.add_argument("--width-pt", type=float); p.set_defaults(func=cmd_set_line)
    p = sub.add_parser("clear-line", help="remove border"); _add_io(p); _add_targets(p, allow_multi=False); p.set_defaults(func=cmd_clear_line)
    p = sub.add_parser("set-shadow", help="add drop shadow"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--color", default="888888"); p.add_argument("--blur-pt", type=float, default=4.0); p.add_argument("--dist-pt", type=float, default=2.0); p.add_argument("--alpha", type=float, default=0.3); p.set_defaults(func=cmd_set_shadow)
    p = sub.add_parser("clear-shadow", help="remove shadow"); _add_io(p); _add_targets(p, allow_multi=False); p.set_defaults(func=cmd_clear_shadow)
    p = sub.add_parser("set-corner-radius", help="round corners"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--ratio", type=float, required=True); p.set_defaults(func=cmd_set_corner_radius)
    p = sub.add_parser("set-opacity", help="fill alpha 0-1"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--alpha", type=float, required=True); p.set_defaults(func=cmd_set_opacity)
    p = sub.add_parser("apply-card-style", help="theme card style"); _add_io(p); _add_targets(p, allow_multi=False); _add_theme(p); p.set_defaults(func=cmd_apply_card_style)
    p = sub.add_parser("apply-badge-style", help="theme badge style"); _add_io(p); _add_targets(p, allow_multi=False); _add_theme(p); p.set_defaults(func=cmd_apply_badge_style)

    # Z-order
    p = sub.add_parser("z-order", help="back/front/up/down"); _add_io(p); _add_targets(p, allow_single=False); p.add_argument("--position", choices=["back", "front", "up", "down"], required=True); p.set_defaults(func=cmd_z_order)
    p = sub.add_parser("send-to-back", help="shortcut z-order back"); _add_io(p); _add_targets(p, allow_single=False); p.set_defaults(func=cmd_send_to_back)
    p = sub.add_parser("bring-to-front", help="shortcut z-order front"); _add_io(p); _add_targets(p, allow_single=False); p.set_defaults(func=cmd_bring_to_front)
    p = sub.add_parser("all-connectors-to-back", help="sweep all connectors to back"); _add_io(p); p.set_defaults(func=cmd_all_connectors_to_back)

    # Text
    p = sub.add_parser("unify-font", help="force one font family deck-wide (default Microsoft YaHei, Latin + East Asian)")
    _add_io(p); p.add_argument("--family", default="Microsoft YaHei")
    p.set_defaults(func=cmd_unify_font)
    p = sub.add_parser("set-font-family", help="font family"); _add_io(p); _add_targets(p); p.add_argument("--family", required=True); p.add_argument("--scope", choices=["all"]); p.set_defaults(func=cmd_set_font_family)
    p = sub.add_parser("set-font-size", help="font size in pt"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--size-pt", type=float, required=True); p.set_defaults(func=cmd_set_font_size)
    p = sub.add_parser("set-font-bold", help="toggle bold"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--bold", type=_str_to_bool, required=True); p.set_defaults(func=cmd_set_font_bold)
    p = sub.add_parser("set-font-italic", help="toggle italic"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--italic", type=_str_to_bool, required=True); p.set_defaults(func=cmd_set_font_italic)
    p = sub.add_parser("set-font-color", help="font color"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--color", required=True); p.set_defaults(func=cmd_set_font_color)
    p = sub.add_parser("set-text-align", help="horizontal align"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--align", choices=["left", "center", "right", "justify"], required=True); p.set_defaults(func=cmd_set_text_align)
    p = sub.add_parser("set-text-v-align", help="vertical anchor"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--anchor", choices=["top", "middle", "bottom"], required=True); p.set_defaults(func=cmd_set_text_v_align)
    p = sub.add_parser("set-text-margin", help="text frame inner margins"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--left-emu", type=int); p.add_argument("--right-emu", type=int); p.add_argument("--top-emu", type=int); p.add_argument("--bottom-emu", type=int); p.set_defaults(func=cmd_set_text_margin)
    p = sub.add_parser("set-line-spacing", help="paragraph line spacing"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--ratio", type=float, required=True); p.set_defaults(func=cmd_set_line_spacing)
    p = sub.add_parser("set-text", help="replace text content"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--content", required=True); p.set_defaults(func=cmd_set_text)
    p = sub.add_parser("apply-typography", help="apply theme role typography"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--role", required=True); _add_theme(p); p.set_defaults(func=cmd_apply_typography)

    # Slide
    p = sub.add_parser("delete-shape", help="delete shape"); _add_io(p); _add_targets(p, allow_multi=False); p.set_defaults(func=cmd_delete_shape)
    p = sub.add_parser("duplicate-shape", help="duplicate shape"); _add_io(p); p.add_argument("--shape-id", type=int, required=True); p.set_defaults(func=cmd_duplicate_shape)
    p = sub.add_parser("add-rect", help="add rectangle"); _add_io(p); p.add_argument("--slide", type=int, required=True); p.add_argument("--left", type=int, required=True); p.add_argument("--top", type=int, required=True); p.add_argument("--width", type=int, required=True); p.add_argument("--height", type=int, required=True); p.add_argument("--fill"); p.add_argument("--line"); p.add_argument("--line-width-pt", type=float); p.set_defaults(func=cmd_add_rect)
    p = sub.add_parser("add-line", help="add line connector"); _add_io(p); p.add_argument("--slide", type=int, required=True); p.add_argument("--x1", type=int, required=True); p.add_argument("--y1", type=int, required=True); p.add_argument("--x2", type=int, required=True); p.add_argument("--y2", type=int, required=True); p.add_argument("--color"); p.add_argument("--width-pt", type=float); p.set_defaults(func=cmd_add_line)
    p = sub.add_parser("add-text", help="add text box"); _add_io(p); p.add_argument("--slide", type=int, required=True); p.add_argument("--left", type=int, required=True); p.add_argument("--top", type=int, required=True); p.add_argument("--width", type=int, required=True); p.add_argument("--height", type=int, required=True); p.add_argument("--content", required=True); p.add_argument("--role"); _add_theme(p); p.set_defaults(func=cmd_add_text)

    # Placement (model-friendly)
    p = sub.add_parser("place-near", help="move shape near another (anchor-based)")
    _add_io(p); p.add_argument("--shape-id", type=int, required=True)
    p.add_argument("--reference-shape-id", type=int, required=True)
    p.add_argument("--anchor", required=True,
                   choices=["above", "below", "left-of", "right-of",
                            "inside-top-left", "inside-top-right", "inside-center",
                            "inside-bottom-left", "inside-bottom-right"])
    p.add_argument("--gap-emu", type=int, default=0)
    p.set_defaults(func=cmd_place_near)

    p = sub.add_parser("mirror-peer-position", help="copy a peer's relative offset within its card")
    _add_io(p); p.add_argument("--shape-id", type=int, required=True)
    p.add_argument("--target-container-id", type=int, required=True)
    p.add_argument("--peer-shape-id", type=int, required=True)
    p.add_argument("--peer-container-id", type=int, required=True)
    p.set_defaults(func=cmd_mirror_peer_position)

    p = sub.add_parser("align-to-grid-cell", help="grid-cell snap: take top from one peer, left from another")
    _add_io(p); p.add_argument("--shape-id", type=int, required=True)
    p.add_argument("--top-from-shape-id", type=int)
    p.add_argument("--left-from-shape-id", type=int)
    p.set_defaults(func=cmd_align_to_grid_cell)

    p = sub.add_parser("move-to-card", help="drop shape into another container at fractional rel position")
    _add_io(p); p.add_argument("--shape-id", type=int, required=True)
    p.add_argument("--container-id", type=int, required=True)
    p.add_argument("--rel-x", type=float, default=0.10)
    p.add_argument("--rel-y", type=float, default=0.10)
    p.set_defaults(func=cmd_move_to_card)

    # Repair
    p = sub.add_parser("repair-peer-cards", help="auto-fix peer-card outliers"); _add_io(p); p.add_argument("--scope", choices=["all", "safe", "no-orphans"], default="all"); p.set_defaults(func=cmd_repair_peer_cards)
    p = sub.add_parser("repair-grid", help="auto-fix 2D grid layout outliers (NxM dashboards)")
    _add_io(p); p.add_argument("--nested", action="store_true", help="recurse into each panel for sub-grids")
    p.set_defaults(func=cmd_repair_grid)
    p = sub.add_parser("repair-peers-smart",
                        help="apply uniform size/spacing/alignment per agent-defined peer groups")
    _add_io(p)
    p.add_argument("--groups", required=True, type=Path,
                   help="JSON file the agent wrote: {\"groups\": [...]}")
    p.set_defaults(func=cmd_repair_peers_smart)

    # Polish
    p = sub.add_parser("polish-business", help="business-grade visual polish (typography + decoration; doesn't change content)")
    _add_io(p)
    p.add_argument("--level", type=int, choices=[1, 2, 3], default=2,
                   help="1=subtle 2=standard 3=rich")
    p.add_argument("--theme", type=Path, default=None,
                   help="optional theme JSON; if omitted, auto-pick from content")
    p.set_defaults(func=cmd_polish_business)

    # Connector
    p = sub.add_parser("style-connector", help="theme connector style"); _add_io(p); _add_targets(p, allow_multi=False); _add_theme(p); p.set_defaults(func=cmd_style_connector)
    p = sub.add_parser("straighten-connector", help="force horizontal/vertical"); _add_io(p); _add_targets(p, allow_multi=False); p.add_argument("--axis", choices=["horizontal", "vertical"], required=True); p.set_defaults(func=cmd_straighten_connector)

    # Meta
    p = sub.add_parser("list-ops", help="list catalog"); p.add_argument("--json", action="store_true"); p.set_defaults(func=cmd_list_ops)

    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        return args.func(args) or 0
    except SystemExit:
        raise
    except Exception as exc:
        sys.stderr.write(json.dumps({"error": type(exc).__name__, "message": str(exc)}) + "\n")
        return 3


if __name__ == "__main__":
    raise SystemExit(main())
