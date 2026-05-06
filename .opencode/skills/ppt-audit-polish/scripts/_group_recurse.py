"""Recursive group-shape walker with proper coordinate transformation.

PowerPoint groups have their own coordinate space. A child shape's (x, y)
is relative to the group's chOff/chExt, NOT the slide. The transform is:

    slide_x = group.off_x + (child.x - group.chOff_x) * (group.ext_cx / group.chExt_cx)
    slide_y = group.off_y + (child.y - group.chOff_y) * (group.ext_cy / group.chExt_cy)

Nested groups multiply their transforms.

This module returns a flat list of "group_children" with slide-space
coordinates so detectors can opt into them WITHOUT them appearing in the
primary `objects` list (which would double-count for density / peer-row
math against top-level shapes).
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE


_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _read_group_xfrm(group_shape) -> dict | None:
    """Extract off/ext/chOff/chExt from a group's grpSpPr.

    All four sub-elements are required to compute the child transform; if
    any is missing we return None and skip transformation (children get
    their raw coords, which is wrong but no worse than today's behavior).
    """
    el = group_shape._element
    xfrm = None
    for child in el.iter():
        tag = child.tag
        if tag.endswith("}grpSpPr"):
            for sub in child.iter():
                if sub.tag.endswith("}xfrm"):
                    xfrm = sub
                    break
            break
    if xfrm is None:
        return None
    off = xfrm.find(_NS_A + "off")
    ext = xfrm.find(_NS_A + "ext")
    chOff = xfrm.find(_NS_A + "chOff")
    chExt = xfrm.find(_NS_A + "chExt")
    if None in (off, ext, chOff, chExt):
        return None
    try:
        return {
            "off_x": int(off.get("x")), "off_y": int(off.get("y")),
            "ext_cx": int(ext.get("cx")), "ext_cy": int(ext.get("cy")),
            "chOff_x": int(chOff.get("x")), "chOff_y": int(chOff.get("y")),
            "chExt_cx": max(int(chExt.get("cx")), 1),
            "chExt_cy": max(int(chExt.get("cy")), 1),
        }
    except (TypeError, ValueError):
        return None


def _apply(transform: dict, x: int, y: int) -> tuple[int, int]:
    sx = transform["off_x"] + (x - transform["chOff_x"]) * transform["ext_cx"] / transform["chExt_cx"]
    sy = transform["off_y"] + (y - transform["chOff_y"]) * transform["ext_cy"] / transform["chExt_cy"]
    return int(sx), int(sy)


def _scale(transform: dict, cx: int, cy: int) -> tuple[int, int]:
    sx = cx * transform["ext_cx"] / transform["chExt_cx"]
    sy = cy * transform["ext_cy"] / transform["chExt_cy"]
    return int(sx), int(sy)


def walk_group_children(group_shape, group_path: str = "") -> list[dict]:
    """Yield child geometry in SLIDE space.

    Each child is a dict with the same shape contract as inspect_ppt's
    top-level objects, plus:
      "group_path": "Group 1 / Group 2"  (the path of group names)
      "in_group": True
    """
    transform = _read_group_xfrm(group_shape)
    if transform is None:
        return []

    children_out: list[dict] = []
    own_path = (group_path + "/" if group_path else "") + getattr(group_shape, "name", "Group")
    for child in group_shape.shapes:
        ctype = getattr(child, "shape_type", None)
        if ctype == MSO_SHAPE_TYPE.GROUP:
            children_out.extend(walk_group_children(child, own_path))
            continue
        try:
            cx, cy = int(child.left), int(child.top)
            cw, ch = int(child.width), int(child.height)
        except (AttributeError, TypeError):
            continue
        slide_left, slide_top = _apply(transform, cx, cy)
        slide_w, slide_h = _scale(transform, cw, ch)
        children_out.append({
            "shape_id": getattr(child, "shape_id", 0),
            "name": getattr(child, "name", ""),
            "kind": _kind(child),
            "shape_type": int(ctype) if ctype is not None else None,
            "left": slide_left,
            "top": slide_top,
            "width": slide_w,
            "height": slide_h,
            "anomalous": slide_w <= 0 or slide_h <= 0,
            "in_group": True,
            "group_path": own_path,
            "rotation": float(getattr(child, "rotation", 0.0) or 0.0),
            "text": _text_of(child),
            "font_sizes": _font_sizes_of(child),
        })
    return children_out


def _kind(shape) -> str:
    st = getattr(shape, "shape_type", None)
    if st == MSO_SHAPE_TYPE.LINE:
        return "connector"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if getattr(shape, "has_text_frame", False):
        text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
        return "text" if text else "container"
    return "shape"


def _text_of(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def _font_sizes_of(shape) -> list[int]:
    if not getattr(shape, "has_text_frame", False):
        return []
    sizes: list[int] = []
    for p in shape.text_frame.paragraphs:
        run_sizes = [int(r.font.size) for r in p.runs if r.font.size is not None]
        if run_sizes:
            sizes.extend(run_sizes)
        elif p.font.size is not None:
            sizes.append(int(p.font.size))
    return sizes
