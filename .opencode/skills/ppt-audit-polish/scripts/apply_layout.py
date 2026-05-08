"""Universal renderer for agent-designed free-form layouts.

Fixed templates would hard-code layout decisions (always render N
evenly-spaced numbered cards, always 2-column grid, etc.) — useful
as presets but constraining when the content doesn't fit the pattern.
This script flips it:
  * Agent reads content.json + assets-manifest.json + the annotated
    render of the original deck
  * Agent decides the WHOLE layout — number of columns, hierarchy,
    where every shape sits — and writes a `layout.json`
  * This script consumes layout.json + content.json + assets/ and
    renders one slide

Layout JSON schema:

  {
    "background": "0F0F0F",         // optional slide background fill
    "slide_dims": {"width": 12192000, "height": 6858000},  // optional
    "elements": [
      {"kind": "fill", "bbox": [0,0,W,H], "color": "0F0F0F"},
      {"kind": "rect", "bbox": [...], "fill": "...", "line": "...",
       "line_pt": 0.5},
      {"kind": "rounded_rect", "bbox": [...], "fill": "...",
       "border": "...", "border_pt": 0.5, "corner_ratio": 0.04},
      {"kind": "circle", "bbox": [cx-r,cy-r,2r,2r], "fill": "..."},
      {"kind": "line", "bbox": [x1, y1, x2, y2], "color": "...",
       "width_pt": 0.5},
      {"kind": "text",
       "bbox": [...],
       "content": "literal string",  // OR
       "ref": "title",                // dotted path into content.json
       "size_pt": 24, "bold": true, "italic": false,
       "color": "F5F5F5", "font": "Consolas",
       "align": "left|center|right", "v_align": "top|middle|bottom"},
      {"kind": "rich_text",
       "bbox": [...],
       "runs": [...]                  // inline list of run dicts OR
       "runs_ref": "items.0.name_runs",  // dotted path to a runs list
       "default_size_pt": 14, "default_color": "F5F5F5",
       "default_font": "Microsoft YaHei",
       "align": "left", "v_align": "middle"},
      {"kind": "image", "bbox": [...],
       "path": "assets/sid_42.png",   // direct path OR
       "ref": "items.0.image",         // dotted path to image dict
       "fit_mode": "stretch|contain|cover", // default stretch
       "crop": {"left":0.1,"top":0.05,"right":0.05,"bottom":0.1}, // optional
       "z_index": 5},                  // optional, sorts render order
      {"kind": "table", "bbox": [...],
       "cells": [[{"text":"...","fill_hex":"..."}, ...], ...],   // OR
       "ref": "tables.0.cells"},
      {"kind": "chart", "bbox": [...],
       "chart_type": "BAR_CLUSTERED",
       "categories": [...],
       "series": [{"name": "...", "values": [...]}]},
      ...
    ]
  }

  Top-level "background" can be:
    - "FFFFFF"         (solid hex shorthand)
    - {"type": "solid",    "color": "..."}
    - {"type": "gradient", "stops": [{"pos":0,"color":"..."}], "angle": 0}
    - {"type": "picture",  "path": "..."}

Refs are dotted paths: `title`, `subtitle`, `items.0.name`,
`items.3.description`, `items.0.image` (image refs return the dict
{"path": ..., "asset_id": ...} and use the path).

Skips elements with missing refs silently (so the renderer is
forgiving when content.json doesn't have every slot the agent
named).

CLI:
  python apply_layout.py --content content.json \
      --layout layout.json --out fresh.pptx \
      [--assets-base assets/]
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_ROOT = SCRIPT_DIR.parent


def _resolve_ref(content: dict, ref: str):
    """Walk a dotted path through content. Returns None on miss."""
    cur = content
    for part in ref.split("."):
        if isinstance(cur, dict):
            cur = cur.get(part)
        elif isinstance(cur, list):
            try:
                cur = cur[int(part)]
            except (ValueError, IndexError):
                return None
        else:
            return None
        if cur is None:
            return None
    return cur


def render_layout(layout: dict, content: dict, out_path: Path,
                   assets_base: Path | None = None) -> dict:
    from _base import (
        add_blank_slide, add_chart, add_gradient_background,
        add_image_from_path, add_rect, add_rich_text, add_rounded_rect,
        add_table, add_text, wipe_slides,
    )
    # add_circle and add_line aren't exported by _base (line uses connector
    # API directly); use python-pptx primitives where _base helpers don't fit.
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    dims = layout.get("slide_dims", {})
    SW = int(dims.get("width", 12192000))
    SH = int(dims.get("height", 6858000))

    prs = Presentation()
    prs.slide_width = Emu(SW); prs.slide_height = Emu(SH)
    wipe_slides(prs)
    slide = add_blank_slide(prs)

    # Optional background — string for solid color, dict for gradient/picture.
    bg = layout.get("background")
    if isinstance(bg, str):
        add_rect(slide, 0, 0, SW, SH, fill=bg, line="none")
    elif isinstance(bg, dict):
        bg_type = bg.get("type")
        if bg_type == "solid" and bg.get("color"):
            color = bg["color"]
            if not color.startswith("@"):
                add_rect(slide, 0, 0, SW, SH, fill=color, line="none")
        elif bg_type == "gradient":
            stops = bg.get("stops") or []
            stops = [s for s in stops if not str(s.get("color", "")).startswith("@")]
            if stops:
                add_gradient_background(slide, SW, SH, stops,
                                         angle=bg.get("angle"))
        elif bg_type == "picture" and bg.get("path"):
            from pathlib import Path as _P
            path = _P(bg["path"])
            if not path.is_absolute() and assets_base:
                path = assets_base / bg["path"]
            if path.exists():
                add_image_from_path(slide, str(path), 0, 0, SW, SH,
                                     fit_mode="cover")

    skipped: list[dict] = []
    rendered_count = 0

    # If any element has an explicit z_index, render in z-order (back-to-
    # front). Elements without z_index keep their original layout.json
    # position, which becomes their effective z_index.
    elements = list(enumerate(layout.get("elements", [])))
    if any("z_index" in el for _, el in elements):
        elements.sort(key=lambda pair: (
            int(pair[1].get("z_index", pair[0])),
            pair[0],   # stable tie-break by original order
        ))
    elements = [el for _, el in elements]

    for el in elements:
        kind = el.get("kind")
        bbox = el.get("bbox") or [0, 0, 0, 0]

        if kind == "fill":
            add_rect(slide, bbox[0], bbox[1], bbox[2], bbox[3],
                     fill=el.get("color"), line="none")
            rendered_count += 1

        elif kind == "rect":
            add_rect(slide, bbox[0], bbox[1], bbox[2], bbox[3],
                     fill=el.get("fill"),
                     line=el.get("line"),
                     line_pt=el.get("line_pt"))
            rendered_count += 1

        elif kind == "rounded_rect":
            add_rounded_rect(slide, bbox[0], bbox[1], bbox[2], bbox[3],
                             fill=el.get("fill"),
                             line=el.get("border"),
                             line_pt=el.get("border_pt"),
                             corner_ratio=el.get("corner_ratio", 0.06))
            rendered_count += 1

        elif kind == "circle":
            sp = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Emu(bbox[0]), Emu(bbox[1]), Emu(bbox[2]), Emu(bbox[3]),
            )
            fill = el.get("fill")
            if fill:
                sp.fill.solid()
                hexs = fill.lstrip("#")
                sp.fill.fore_color.rgb = RGBColor(
                    int(hexs[0:2], 16), int(hexs[2:4], 16), int(hexs[4:6], 16),
                )
            try:
                sp.line.fill.background()
            except (AttributeError, ValueError):
                pass
            rendered_count += 1

        elif kind == "line":
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(bbox[0]), Emu(bbox[1]), Emu(bbox[2]), Emu(bbox[3]),
            )
            color = el.get("color", "6F6F6F").lstrip("#")
            line.line.color.rgb = RGBColor(
                int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16),
            )
            line.line.width = Pt(float(el.get("width_pt", 0.75)))
            rendered_count += 1

        elif kind == "text":
            # Resolve content: literal string wins, else ref.
            text_value = el.get("content")
            if text_value is None and el.get("ref"):
                resolved = _resolve_ref(content, el["ref"])
                if resolved is None:
                    skipped.append({"reason": "ref-missing", "ref": el["ref"]})
                    continue
                text_value = str(resolved)
            if not text_value:
                skipped.append({"reason": "empty-text", "kind": kind})
                continue
            add_text(slide, bbox[0], bbox[1], bbox[2], bbox[3],
                     text_value,
                     size_pt=float(el.get("size_pt", 11)),
                     bold=bool(el.get("bold", False)),
                     italic=bool(el.get("italic", False)),
                     color=el.get("color", "393939"),
                     family=el.get("font"),
                     align=el.get("align", "left"),
                     v_align=el.get("v_align", "top"),
                     fill=el.get("fill"))
            rendered_count += 1

        elif kind == "rich_text":
            # Multi-run text with mixed formatting.
            # `runs` can be inline (in the layout element) OR resolved
            # via `runs_ref` into content.json (e.g., "items.0.name_runs").
            runs = el.get("runs")
            if not runs and el.get("runs_ref"):
                resolved = _resolve_ref(content, el["runs_ref"])
                if isinstance(resolved, list):
                    runs = resolved
            if not runs:
                skipped.append({"reason": "rich-text-no-runs",
                                "ref": el.get("runs_ref")})
                continue
            add_rich_text(slide, bbox[0], bbox[1], bbox[2], bbox[3], runs,
                          default_size_pt=float(el.get("default_size_pt", 11)),
                          default_color=el.get("default_color", "393939"),
                          default_family=el.get("default_font"),
                          align=el.get("align", "left"),
                          v_align=el.get("v_align", "top"),
                          fill=el.get("fill"))
            rendered_count += 1

        elif kind == "table":
            cells = el.get("cells")
            if not cells and el.get("ref"):
                cells = _resolve_ref(content, el["ref"])
            if not cells:
                skipped.append({"reason": "table-no-cells"})
                continue
            result = add_table(slide, bbox[0], bbox[1], bbox[2], bbox[3],
                                cells,
                                border_color=el.get("border_color", "DDE1E6"),
                                border_pt=float(el.get("border_pt", 0.5)),
                                header_fill=el.get("header_fill"),
                                default_text_color=el.get("color", "393939"),
                                font_family=el.get("font"),
                                font_size_pt=float(el.get("size_pt", 10)))
            if result:
                rendered_count += 1

        elif kind == "chart":
            categories = el.get("categories") or []
            series = el.get("series") or []
            if el.get("ref"):
                resolved = _resolve_ref(content, el["ref"])
                if isinstance(resolved, dict):
                    categories = resolved.get("categories") or categories
                    series = resolved.get("series") or series
            if not series or not categories:
                skipped.append({"reason": "chart-no-data"})
                continue
            result = add_chart(slide, bbox[0], bbox[1], bbox[2], bbox[3],
                                chart_type=el.get("chart_type", "BAR_CLUSTERED"),
                                categories=categories,
                                series=series)
            if result:
                rendered_count += 1

        elif kind == "image":
            # Resolve image path: explicit path or ref into content.
            img_path = el.get("path")
            if not img_path and el.get("ref"):
                resolved = _resolve_ref(content, el["ref"])
                if isinstance(resolved, dict):
                    img_path = resolved.get("path")
                elif isinstance(resolved, str):
                    img_path = resolved
            if not img_path:
                skipped.append({"reason": "image-path-missing",
                                "ref": el.get("ref")})
                continue
            # Resolve relative paths against assets_base if provided,
            # else against the layout file's directory.
            p = Path(img_path)
            if not p.is_absolute() and assets_base:
                p = assets_base / img_path
            if not p.exists():
                skipped.append({"reason": "image-not-found", "path": str(p)})
                continue
            result = add_image_from_path(
                slide, str(p),
                bbox[0], bbox[1], bbox[2], bbox[3],
                fit_mode=el.get("fit_mode", "stretch"),
                crop=el.get("crop"),
            )
            if result is None:
                skipped.append({"reason": "image-render-failed", "path": str(p)})
            else:
                rendered_count += 1

        else:
            skipped.append({"reason": "unknown-kind", "kind": kind})

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {
        "out": str(out_path),
        "rendered": rendered_count,
        "skipped": skipped,
        "elements_total": len(layout.get("elements", [])),
    }


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Render a slide from an agent-designed layout.json + content.json. "
            "The layout is fully agent-driven — no fixed columns, no preset grid."
        ),
    )
    parser.add_argument("--content", required=True, type=Path,
                        help="content.json from extract_content.py")
    parser.add_argument("--layout", required=True, type=Path,
                        help="layout.json the agent designed")
    parser.add_argument("--out", required=True, type=Path,
                        help="output .pptx")
    parser.add_argument("--assets-base", type=Path, default=None,
                        help="folder to resolve relative image paths "
                             "(default: same dir as content.json)")
    args = parser.parse_args()

    content = json.loads(args.content.read_text(encoding="utf-8"))
    layout = json.loads(args.layout.read_text(encoding="utf-8"))
    assets_base = args.assets_base or args.content.parent

    result = render_layout(layout, content, args.out, assets_base=assets_base)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
