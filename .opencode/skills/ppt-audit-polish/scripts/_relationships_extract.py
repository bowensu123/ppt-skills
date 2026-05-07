"""Extract identity + relationship metadata per shape.

Path B's regenerate flow needs to know which shapes have INBOUND
references (other things pointing AT them: connector arrows,
hyperlinks, click actions, comments) and which are placeholders. The
agent uses this to decide PER SHAPE whether to:

  * preserve_identity — keep shape_id, name, placeholder type so the
    references continue to work; only move/resize/restyle in place.
  * recreate — delete the original and create a new primitive
    elsewhere; identity is lost but content (text/binary) can be
    transferred. Safe when no inbound references exist.
  * delete — remove without replacement.

Content-resource extraction (binary) is handled by _asset_extract;
this module focuses purely on identity + relationship metadata so the
agent's preservation decisions are well-informed.

Output schema (relationships.json):

  {
    "input": "...",
    "slide_count": N,
    "slides": [{
      "slide_index": 1,
      "shapes": [{
        "shape_id": 42,
        "name": "VLM panel",
        "placeholder": {
          "type": "TITLE" | "BODY" | "PICTURE" | "OBJECT" | ...,
          "idx": 0,
          "inherited_from": "layout" | "master" | null
        } | null,
        "group_membership": { "group_shape_id": 99, "group_name": "..." } | null,
        "is_referenced_by": [
          {"kind": "connector", "from_shape_id": 88, "endpoint": "begin"},
          {"kind": "hyperlink", "from_shape_id": 101,
           "from_slide_index": 2, "from_run_text": "see VLM model"},
          {"kind": "click_action", "from_shape_id": 77,
           "action_type": "ppaction://hlinksldjump"},
          {"kind": "comment", "comment_id": "c1", "author": "..."}
        ],
        "references_to": [
          {"kind": "hyperlink", "to": "https://..." | "slide:3" | ...,
           "from_run_text": "..."},
          {"kind": "click_action", "target": "next-slide"}
        ],
        "preserve_identity_default": true | false,
        "rationale_default": "..."
      }]
    }]
  }

The agent reads this + render + content/composition manifests, then
writes `relocation.json` which decides per-shape:
  - "preserve_identity" → apply_relocation moves/resizes original shape
  - "recreate"          → apply_relocation deletes + creates new primitive
  - "delete"            → apply_relocation removes
"""
from __future__ import annotations

import argparse
import json
from collections import defaultdict
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# ---- placeholder info ----

def _placeholder_info(shape) -> dict | None:
    try:
        ph = shape.placeholder_format
        if ph is None:
            return None
        ph_type = ph.type
        return {
            "type": ph_type.name if hasattr(ph_type, "name") else str(ph_type),
            "idx": int(ph.idx) if ph.idx is not None else None,
            # Whether it inherits from layout/master is best-effort —
            # python-pptx doesn't expose this directly. Most placeholders
            # do inherit, mark as "layout" by default.
            "inherited_from": "layout",
        }
    except (AttributeError, ValueError, KeyError):
        return None


# ---- group membership ----

def _walk_groups(shapes, parent=None):
    """Yield (shape, parent_group_or_None) recursing into groups."""
    for shape in shapes:
        try:
            stype = shape.shape_type
        except (AttributeError, ValueError):
            stype = None
        if stype == MSO_SHAPE_TYPE.GROUP:
            yield shape, parent  # the group itself
            yield from _walk_groups(shape.shapes, shape)
        else:
            yield shape, parent


# ---- connector relationships ----

def _connector_endpoints(shape) -> tuple[int | None, int | None]:
    """Return (begin_shape_id, end_shape_id). None when free / not connected.

    PPTX cxnSp has nvCxnSpPr/cNvCxnSpPr with stCxn (start) and endCxn (end)
    children whose `id` attribute is the target shape's spId.
    """
    try:
        xml = shape._element
        cxn_sp_pr = xml.find(".//p:nvCxnSpPr/p:cNvCxnSpPr", _NS)
        if cxn_sp_pr is None:
            return (None, None)
        st = cxn_sp_pr.find("a:stCxn", _NS)
        ed = cxn_sp_pr.find("a:endCxn", _NS)
        st_id = int(st.get("id")) if st is not None and st.get("id") else None
        ed_id = int(ed.get("id")) if ed is not None and ed.get("id") else None
        return (st_id, ed_id)
    except (AttributeError, ValueError, etree.LxmlError):
        return (None, None)


def _is_connector(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.LINE
    except (AttributeError, ValueError):
        # Cxn shapes also: check XML tag
        try:
            return etree.QName(shape._element).localname == "cxnSp"
        except (AttributeError, etree.LxmlError):
            return False


# ---- click actions / hyperlinks on shape (not text-run) ----

def _shape_actions(shape) -> list[dict]:
    """Click actions defined on the shape itself (not its text)."""
    out = []
    try:
        xml = shape._element
        # nvSpPr/cNvPr/a:hlinkClick or directly a:hlinkClick under nvSpPr.
        # For pictures: nvPicPr/cNvPr/a:hlinkClick
        for hlink in xml.findall(".//a:hlinkClick", _NS):
            target = (hlink.get(f"{{{_NS['r']}}}id")
                       or hlink.get("action") or "")
            action_type = hlink.get("action") or "ppaction://hlink"
            out.append({
                "action_type": action_type,
                "target_rid_or_action": target,
            })
    except (AttributeError, etree.LxmlError):
        pass
    return out


# ---- text-run hyperlinks (hyperlinks INSIDE text) ----

def _text_run_hyperlinks(shape) -> list[dict]:
    """Hyperlinks attached to individual text runs."""
    out = []
    if not getattr(shape, "has_text_frame", False):
        return out
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    addr = run.hyperlink.address if run.hyperlink else None
                except (AttributeError, ValueError):
                    addr = None
                if addr:
                    out.append({
                        "from_run_text": run.text,
                        "address": addr,
                    })
    except (AttributeError, ValueError):
        pass
    return out


# ---- comments (slide-level comments part) ----

def _slide_comments(slide) -> list[dict]:
    """Best-effort: extract comments anchored to shapes on this slide.

    PPTX comments live in /ppt/comments/comment{N}.xml and are linked to
    slides via the slide rels. Each comment can have an authorId and
    optional anchor coordinates (which we'd need to map back to shapes).
    """
    out = []
    try:
        # python-pptx doesn't have first-class comment support; we walk
        # the slide's XML rels manually.
        slide_part = slide.part
        for rel_id, rel in slide_part.rels.items():
            target = rel.target_partname if hasattr(rel, "target_partname") else ""
            if "comments" in str(target):
                try:
                    comment_xml = etree.fromstring(rel.target_part.blob)
                    for cmt in comment_xml.findall(".//p:cm", _NS):
                        out.append({
                            "comment_id": cmt.get("authorId", ""),
                            "text": "".join(cmt.itertext()).strip()[:200],
                        })
                except (etree.LxmlError, AttributeError):
                    continue
    except (AttributeError, ValueError):
        pass
    return out


# ---- main extractor ----

def extract_relationships(input_path: Path, work_dir: Path) -> dict:
    prs = Presentation(str(input_path))
    work_dir.mkdir(parents=True, exist_ok=True)

    slides_out = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Build inbound-reference map for this slide
        inbound: dict[int, list[dict]] = defaultdict(list)

        # 1. Connectors with begin/end pointing at shapes
        for shape, parent_group in _walk_groups(slide.shapes):
            if _is_connector(shape):
                st, ed = _connector_endpoints(shape)
                if st is not None:
                    inbound[st].append({
                        "kind": "connector",
                        "from_shape_id": shape.shape_id,
                        "endpoint": "begin",
                    })
                if ed is not None:
                    inbound[ed].append({
                        "kind": "connector",
                        "from_shape_id": shape.shape_id,
                        "endpoint": "end",
                    })

        # 2. Click actions on shapes targeting other shapes (hard to
        #    cross-reference without resolving rIds → just record actions)
        # (Captured per shape under references_to, not inbound, since we
        # can't reliably know the target shape_id without rel-resolution.)

        # 3. Comments — anchored to slides, optionally to shapes
        comments = _slide_comments(slide)

        # Per-shape entries
        shape_entries = []
        group_lookup = {}
        for shape, parent_group in _walk_groups(slide.shapes):
            if shape.shape_id is None:
                continue
            sid = int(shape.shape_id)
            entry = {
                "shape_id": sid,
                "name": shape.name or "",
                "placeholder": _placeholder_info(shape),
                "group_membership": (
                    {"group_shape_id": int(parent_group.shape_id),
                     "group_name": parent_group.name or ""}
                    if parent_group is not None
                    else None
                ),
                "is_connector": _is_connector(shape),
                "is_referenced_by": list(inbound.get(sid, [])),
                "references_to": (
                    [
                        {"kind": "click_action", **act}
                        for act in _shape_actions(shape)
                    ]
                    + [
                        {"kind": "hyperlink", **link}
                        for link in _text_run_hyperlinks(shape)
                    ]
                ),
            }

            # Default preservation suggestion (agent can override)
            placeholder = entry["placeholder"]
            inbound_count = len(entry["is_referenced_by"])
            outbound_count = len(entry["references_to"])
            preserve = bool(
                placeholder
                or inbound_count > 0
                or outbound_count > 0
                or entry["is_connector"]   # connectors must keep their refs
            )
            rationale_parts = []
            if placeholder:
                rationale_parts.append(
                    f"placeholder type={placeholder['type']}"
                )
            if inbound_count:
                rationale_parts.append(
                    f"{inbound_count} inbound ref(s) "
                    f"({{kinds}})".format(
                        kinds=",".join(sorted({r["kind"] for r in entry["is_referenced_by"]}))
                    )
                )
            if outbound_count:
                rationale_parts.append(f"{outbound_count} outbound ref(s)")
            if entry["is_connector"]:
                rationale_parts.append("connector — endpoints would break if recreated")
            if not preserve:
                rationale_parts.append(
                    "no relationships, not a placeholder — safe to recreate"
                )
            entry["preserve_identity_default"] = preserve
            entry["rationale_default"] = "; ".join(rationale_parts)

            shape_entries.append(entry)

        slides_out.append({
            "slide_index": slide_idx,
            "shapes": shape_entries,
            "comments": comments,
        })

    manifest = {
        "input": str(input_path),
        "slide_count": len(prs.slides),
        "slides": slides_out,
    }
    (work_dir / "relationships.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Extract per-shape identity (shape_id/name/placeholder) and "
            "relationship metadata (connector endpoints, text-run hyperlinks, "
            "click actions, group membership, comments) so the agent can "
            "decide per-shape whether to preserve_identity, recreate, or "
            "delete during Path B regenerate."
        ),
    )
    parser.add_argument("--in", dest="in_path", required=True, type=Path)
    parser.add_argument("--work-dir", required=True, type=Path)
    args = parser.parse_args()

    manifest = extract_relationships(args.in_path, args.work_dir)
    summary = {
        "slide_count": manifest["slide_count"],
        "total_shapes": sum(len(s["shapes"]) for s in manifest["slides"]),
        "shapes_with_placeholder": sum(
            1 for s in manifest["slides"] for sh in s["shapes"]
            if sh["placeholder"] is not None
        ),
        "shapes_with_inbound_refs": sum(
            1 for s in manifest["slides"] for sh in s["shapes"]
            if sh["is_referenced_by"]
        ),
        "shapes_default_preserve": sum(
            1 for s in manifest["slides"] for sh in s["shapes"]
            if sh["preserve_identity_default"]
        ),
        "manifest": str(args.work_dir / "relationships.json"),
    }
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
