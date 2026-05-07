"""Shared helpers used by every template renderer.

Templates are pure Python modules with a ``render(prs, content, theme)``
function. They build a fresh slide using python-pptx primitives. These
helpers wrap the verbose pptx API.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


def _hex(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def add_rect(slide, left, top, width, height, fill=None, line=None, line_pt=None,
             corner_ratio: float | None = None, shape=MSO_SHAPE.RECTANGLE) -> object:
    sp = slide.shapes.add_shape(shape, Emu(left), Emu(top), Emu(width), Emu(height))
    if fill is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = _hex(fill)
    elif fill is False:
        sp.fill.background()
    if line == "none":
        try:
            sp.line.fill.background()
        except (AttributeError, ValueError):
            pass
    elif line is not None:
        sp.line.color.rgb = _hex(line)
        if line_pt is not None:
            sp.line.width = Pt(line_pt)
    if corner_ratio is not None:
        try:
            sp.adjustments[0] = corner_ratio
        except (IndexError, AttributeError):
            pass
    return sp


def add_text(slide, left, top, width, height, text, *,
             size_pt: float = 11, bold: bool = False, italic: bool = False,
             color: str = "393939", family: str | None = None,
             align: str = "left", v_align: str = "top",
             fill: str | None = None, padding_emu: int | None = None) -> object:
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.text = ""
    tf.word_wrap = True
    if padding_emu is not None:
        tf.margin_left = Emu(padding_emu); tf.margin_right = Emu(padding_emu)
        tf.margin_top = Emu(padding_emu); tf.margin_bottom = Emu(padding_emu)
    if v_align == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_align == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    elif v_align == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP

    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = _hex(fill)
    try:
        box.line.fill.background()
    except (AttributeError, ValueError):
        pass

    para = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    para.alignment = align_map.get(align, PP_ALIGN.LEFT)

    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex(color)
    if family:
        run.font.name = family
    return box


def add_circle(slide, cx, cy, radius, fill, line="none") -> object:
    return add_rect(slide, cx - radius, cy - radius, radius * 2, radius * 2,
                    fill=fill, line=line, shape=MSO_SHAPE.OVAL)


def add_rounded_rect(slide, left, top, width, height, *, fill=None, line=None, line_pt=None, corner_ratio=0.06):
    return add_rect(slide, left, top, width, height,
                    fill=fill, line=line, line_pt=line_pt,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_ratio=corner_ratio)


def horizontal_distribute(slide_w: int, n: int, margin: int) -> list[int]:
    """Return left-edge x positions for n items evenly spaced across the slide.

    Each item's slot width = (slide_w - 2*margin) / n.
    Returns the left edge of each slot (caller decides item internal offset).
    """
    if n <= 0:
        return []
    available = slide_w - 2 * margin
    slot_w = available // n
    return [margin + i * slot_w for i in range(n)]


def truncate(text: str, max_chars: int) -> str:
    return text if len(text) <= max_chars else text[: max_chars - 1] + "…"


def wipe_slides(prs) -> None:
    """Remove every existing slide so a template can write a fresh first slide."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        sld_id_lst.remove(sld_id)


def add_blank_slide(prs):
    """Add a single blank slide using the layout-6 (blank) of the master."""
    blank_layout = None
    for layout in prs.slide_layouts:
        if (layout.name or "").lower().startswith("blank"):
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs.slides.add_slide(blank_layout)


def add_image_from_path(slide, image_path, left, top, width, height,
                         fit_mode: str = "stretch",
                         crop: dict | None = None) -> object | None:
    """Drop a binary image into the slide at the given EMU coordinates.

    fit_mode controls how the image fills the bbox:
      "stretch" — default, distort to fill (matches python-pptx default)
      "contain" — preserve aspect, letterbox (image fits inside bbox)
      "cover"   — preserve aspect, crop edges to fill bbox
      "crop"    — apply explicit `crop` rect (left/top/right/bottom in
                  fractions 0..1) before placing

    crop applies the source-rectangle from the original PPTX. Format:
      {"left": 0.10, "top": 0.05, "right": 0.05, "bottom": 0.10}
    Values are fractions of the original image to chop off each side.

    Returns the picture shape, or None on failure.
    """
    from pathlib import Path
    p = Path(image_path)
    if not p.exists():
        return None
    try:
        # Step 1: place the picture (default stretches to bbox)
        if fit_mode in ("contain", "cover"):
            # Need to know image native dims to compute aspect-fit box.
            try:
                from PIL import Image
                with Image.open(p) as img:
                    iw, ih = img.size
            except Exception:
                iw, ih = width, height  # graceful fallback
            ox, oy, fw, fh = aspect_fit_box(iw, ih, width, height)
            if fit_mode == "contain":
                left += ox; top += oy
                width = fw; height = fh
            elif fit_mode == "cover":
                # Reverse: scale so image FILLS bbox, crop excess.
                src_ratio = iw / max(ih, 1)
                box_ratio = width / max(height, 1)
                if src_ratio > box_ratio:
                    # too wide — make height fit, width overflow
                    new_w = int(height * src_ratio)
                    pic = slide.shapes.add_picture(
                        str(p),
                        Emu(left - (new_w - width) // 2), Emu(top),
                        Emu(new_w), Emu(height),
                    )
                    return pic
                else:
                    new_h = int(width / src_ratio)
                    pic = slide.shapes.add_picture(
                        str(p),
                        Emu(left), Emu(top - (new_h - height) // 2),
                        Emu(width), Emu(new_h),
                    )
                    return pic

        pic = slide.shapes.add_picture(
            str(p), Emu(left), Emu(top), Emu(width), Emu(height),
        )

        # Step 2: apply explicit crop via srcRect XML (always applies if set,
        # regardless of fit_mode — the agent can combine fit_mode=stretch +
        # crop=<original PPTX srcRect> to faithfully reproduce the source).
        if crop:
            _apply_picture_crop(pic, crop)
        return pic
    except (ValueError, OSError, KeyError):
        return None


def add_table(slide, left, top, width, height, cells: list[list[dict]], *,
              border_color: str = "DDE1E6", border_pt: float = 0.5,
              header_fill: str | None = None,
              default_text_color: str = "393939",
              font_family: str | None = None,
              font_size_pt: float = 10) -> object | None:
    """Render a table from a 2D list of cell dicts.

    cells[row][col] = {"text": str, "fill_hex": str|None, "color_hex": str|None}

    Returns the table shape, or None if cells is empty.
    """
    if not cells or not cells[0]:
        return None
    n_rows = len(cells); n_cols = len(cells[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Emu(left), Emu(top), Emu(width), Emu(height),
    )
    tbl = table_shape.table
    for r, row in enumerate(cells):
        for c, cell_data in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = cell_data.get("text", "") or ""
            fill = cell_data.get("fill_hex") or (header_fill if r == 0 else None)
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex(fill)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size_pt)
                    color = cell_data.get("color_hex") or default_text_color
                    if color:
                        run.font.color.rgb = _hex(color)
                    if font_family:
                        run.font.name = font_family
    return table_shape


def add_chart(slide, left, top, width, height, *,
              chart_type: str = "BAR_CLUSTERED",
              categories: list,
              series: list[dict]) -> object | None:
    """Render a chart from a category list + series data.

    chart_type: short name matching pptx.enum.chart.XL_CHART_TYPE
                ("BAR_CLUSTERED", "LINE", "PIE", "COLUMN_CLUSTERED", etc.)
    series: [{"name": str, "values": [number, ...]}]

    Returns the chart shape, or None if data is empty.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    if not series or not categories:
        return None
    data = CategoryChartData()
    data.categories = categories
    for s in series:
        vals = [v if v is not None else 0 for v in s.get("values", [])]
        data.add_series(s.get("name") or "", vals)
    type_enum = getattr(XL_CHART_TYPE, chart_type, XL_CHART_TYPE.BAR_CLUSTERED)
    chart_frame = slide.shapes.add_chart(
        type_enum, Emu(left), Emu(top), Emu(width), Emu(height), data,
    )
    return chart_frame


def add_gradient_background(slide, slide_w: int, slide_h: int,
                              stops: list[dict], angle: float | None = None) -> object:
    """Render a gradient as the full-slide background using a rect with
    a:gradFill XML. python-pptx doesn't expose gradient fill directly,
    so we build the XML.

    stops: [{"pos": int (0..100000), "color": "RRGGBB"}, ...]
    angle: degrees (optional). PPTX uses 60000ths of a degree.
    """
    from lxml import etree
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Emu(0), Emu(0),
                                    Emu(slide_w), Emu(slide_h))
    # Replace solidFill with gradFill via XML. spPr lives in the `p:`
    # namespace for autoshapes (presentationml), not `a:` (drawingml).
    p_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    sp_pr = rect._element.find(f"{p_ns}spPr")
    if sp_pr is None:
        # Fallback: any spPr regardless of namespace.
        for el in rect._element.iter():
            if el.tag.endswith("}spPr"):
                sp_pr = el
                break
    if sp_pr is None:
        return rect  # graceful: bare rect rendered, no gradient applied
    # Strip existing fills (drawingml namespace for fill children).
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag in ("solidFill", "gradFill", "blipFill", "noFill"):
        for el in sp_pr.findall(f"{ns}{tag}"):
            sp_pr.remove(el)
    grad = etree.SubElement(sp_pr, f"{ns}gradFill",
                              {"flip": "none", "rotWithShape": "1"})
    gs_lst = etree.SubElement(grad, f"{ns}gsLst")
    for stop in stops:
        gs = etree.SubElement(gs_lst, f"{ns}gs",
                                {"pos": str(int(stop.get("pos", 0)))})
        srgb = etree.SubElement(gs, f"{ns}srgbClr",
                                  {"val": stop.get("color", "FFFFFF").lstrip("#")})
    if angle is not None:
        etree.SubElement(grad, f"{ns}lin",
                          {"ang": str(int(angle * 60000)), "scaled": "0"})
    # Send to back so content renders above.
    try:
        sp_tree = rect._element.getparent()
        sp_tree.remove(rect._element)
        sp_tree.insert(2, rect._element)
    except (AttributeError, ValueError):
        pass
    # No border on background.
    try:
        rect.line.fill.background()
    except (AttributeError, ValueError):
        pass
    return rect


def _apply_picture_crop(pic, crop: dict) -> None:
    """Set srcRect on a picture shape via XML.

    crop values are fractions 0..1; PPTX stores them as 1000ths of a
    percent (i.e. fraction × 100000) in the XML.
    """
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    blipFill = pic._element.find(".//a:blipFill", nsmap)
    if blipFill is None:
        return
    # Remove any existing srcRect to avoid duplicates.
    for old in blipFill.findall("a:srcRect", nsmap):
        blipFill.remove(old)
    src = etree.SubElement(blipFill,
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect")
    for side in ("left", "top", "right", "bottom"):
        v = crop.get(side, 0.0)
        if v:
            src.set(side, str(int(round(float(v) * 100000))))
    # blipFill in pptx requires srcRect BEFORE the stretch element. Move it
    # to the front so OOXML stays well-ordered.
    blipFill.insert(1, src)  # after blip, before stretch


def aspect_fit_box(content_w: int, content_h: int,
                    box_w: int, box_h: int) -> tuple[int, int, int, int]:
    """Center an image of (content_w, content_h) inside (box_w, box_h),
    preserving aspect. Returns (offset_x, offset_y, fit_w, fit_h) in
    the same units as inputs.
    """
    if content_w <= 0 or content_h <= 0 or box_w <= 0 or box_h <= 0:
        return (0, 0, box_w, box_h)
    src_ratio = content_w / content_h
    box_ratio = box_w / box_h
    if src_ratio > box_ratio:
        # Width-bound
        fit_w = box_w
        fit_h = int(box_w / src_ratio)
    else:
        # Height-bound
        fit_h = box_h
        fit_w = int(box_h * src_ratio)
    return ((box_w - fit_w) // 2, (box_h - fit_h) // 2, fit_w, fit_h)


def add_rich_text(slide, left, top, width, height, runs: list[dict], *,
                  default_size_pt: float = 11,
                  default_color: str = "393939",
                  default_family: str | None = None,
                  align: str = "left", v_align: str = "top",
                  fill: str | None = None) -> object:
    """Render a text frame with mixed-format runs.

    Each run is a dict with keys:
      text, font_family, size_pt, bold, italic, color_hex,
      paragraph_index (optional), run_index (optional)

    Runs sharing the same paragraph_index land in the same paragraph.
    Missing format fields fall back to the `default_*` parameters. This
    preserves bold/italic/colored words in titles or descriptions that
    Path B would otherwise flatten to plain text.
    """
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.text = ""
    tf.word_wrap = True
    if v_align == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_align == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = _hex(fill)
    try:
        box.line.fill.background()
    except (AttributeError, ValueError):
        pass

    align_map = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
    }
    pp_align = align_map.get(align, PP_ALIGN.LEFT)

    # Group runs by paragraph_index. Default to a single paragraph.
    by_para: dict[int, list[dict]] = {}
    for run in runs:
        p_idx = int(run.get("paragraph_index", 0) or 0)
        by_para.setdefault(p_idx, []).append(run)
    if not by_para:
        return box

    para_indices = sorted(by_para.keys())
    for i, p_idx in enumerate(para_indices):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = pp_align
        for run_data in by_para[p_idx]:
            r = para.add_run()
            r.text = run_data.get("text", "")
            r.font.size = Pt(float(run_data.get("size_pt") or default_size_pt))
            if run_data.get("bold") is not None:
                r.font.bold = bool(run_data["bold"])
            if run_data.get("italic") is not None:
                r.font.italic = bool(run_data["italic"])
            color = run_data.get("color_hex") or default_color
            if color:
                r.font.color.rgb = _hex(color)
            family = run_data.get("font_family") or default_family
            if family:
                r.font.name = family
            link = run_data.get("hyperlink")
            if link:
                try:
                    r.hyperlink.address = link
                except (AttributeError, ValueError):
                    pass
    return box
