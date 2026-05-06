from __future__ import annotations

from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT_DIR = Path(r"C:\Users\42080\Desktop\ppt-skills-main\fix-output")
PPTX_PATH = OUT_DIR / "autods_agent_architecture_beautified.pptx"
PNG_PATH = OUT_DIR / "autods_agent_architecture_beautified_preview.png"

W, H = 13.333, 7.5

COLORS = {
    "bg": "F6F8FC",
    "ink": "102033",
    "muted": "617089",
    "soft": "D9E3F2",
    "line": "AEBBD0",
    "navy": "19324D",
    "blue": "2F6FEB",
    "cyan": "19A7CE",
    "teal": "0E8F7D",
    "green": "3A9C68",
    "amber": "D98924",
    "slate": "EEF3FA",
    "white": "FFFFFF",
    "lav": "EDEBFF",
    "blue_soft": "EAF2FF",
    "cyan_soft": "E8F8FC",
    "green_soft": "EAF7F0",
    "amber_soft": "FFF4E4",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_box(slide, x, y, w, h, fill, line="FFFFFF", radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(0.8)
    return shp


def add_text(slide, text, x, y, w, h, size=14, color="102033", bold=False, align=PP_ALIGN.LEFT,
             font="Microsoft YaHei", valign=MSO_ANCHOR.TOP, margin=0.03):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_para(slide, title, body, x, y, w, h, fill, accent, number=None):
    add_box(slide, x, y, w, h, fill, line="D8E1EF")
    if number:
        badge = add_box(slide, x + w - 0.46, y + 0.13, 0.30, 0.30, accent, line=accent)
        add_text(slide, str(number), x + w - 0.46, y + 0.13, 0.30, 0.30, 10, "FFFFFF", True,
                 PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        title_x = x + 0.16
        title_w = w - 0.66
    else:
        title_x = x + 0.16
        title_w = w - 0.32
    add_text(slide, title, title_x, y + 0.12, title_w, 0.22, 9.4, "102033", True)
    add_text(slide, body, x + 0.16, y + 0.45, w - 0.32, h - 0.53, 8.9, "556579")


def add_arrow(slide, x1, y1, x2, y2, color="8CA0BA", width=1.4, dashed=False):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = rgb(color)
    c.line.width = Pt(width)
    if dashed:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    ln = c._element.spPr.ln
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    ln.append(tail)
    return c


def build_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_box(slide, 0, 0, W, H, COLORS["bg"], line=COLORS["bg"], radius=False)
    add_box(slide, 0, 0, W, 0.14, COLORS["blue"], line=COLORS["blue"], radius=False)
    add_text(slide, "AutoDS-Agent", 0.55, 0.34, 2.0, 0.35, 18, COLORS["blue"], True, font="Aptos Display")
    add_text(slide, "可靠表格机器学习自动化架构", 2.42, 0.30, 5.9, 0.45, 24, COLORS["ink"], True)
    add_text(
        slide,
        "从任务理解到执行修复，再到泄露校验、预算控制与可复现实验记录的一页式可编辑流程图",
        0.58,
        0.78,
        8.4,
        0.28,
        10.5,
        COLORS["muted"],
    )
    add_box(slide, 10.55, 0.36, 2.12, 0.38, COLORS["navy"], line=COLORS["navy"])
    add_text(slide, "EDITABLE  •  ONE PAGE", 10.67, 0.44, 1.88, 0.18, 8.6, "FFFFFF", True,
             PP_ALIGN.CENTER, font="Aptos")

    # Side inputs and outputs.
    add_text(slide, "Input", 0.62, 1.22, 1.6, 0.25, 13.5, COLORS["blue"], True, font="Aptos Display")
    add_box(slide, 0.55, 1.55, 2.08, 3.42, COLORS["white"], line="D6DFEC")
    add_box(slide, 0.55, 1.55, 2.08, 0.12, COLORS["blue"], line=COLORS["blue"], radius=False)
    add_text(slide, "输入", 0.76, 1.86, 1.35, 0.28, 16, COLORS["ink"], True)
    add_text(slide, "数据集 / 表格\n任务类型：classification\n评价指标\n约束 / 预算", 0.76, 2.30, 1.58, 1.55, 10.2, COLORS["muted"])

    add_text(slide, "Output", 10.88, 1.22, 1.6, 0.25, 13.5, COLORS["teal"], True, font="Aptos Display")
    add_box(slide, 10.58, 1.55, 2.18, 3.42, COLORS["white"], line="D6DFEC")
    add_box(slide, 10.58, 1.55, 2.18, 0.12, COLORS["teal"], line=COLORS["teal"], radius=False)
    add_text(slide, "输出", 10.81, 1.86, 1.45, 0.28, 16, COLORS["ink"], True)
    add_text(slide, "best pipeline\npredictive score\ncost summary\nleakage-free report", 10.81, 2.30, 1.62, 1.55, 10.4, COLORS["muted"])

    # Central staged architecture.
    add_text(slide, "A. 任务理解", 2.95, 1.22, 2.5, 0.24, 12, COLORS["blue"], True)
    add_box(slide, 2.95, 1.50, 6.95, 0.03, COLORS["blue"], line=COLORS["blue"], radius=False)
    add_para(slide, "Schema Inspector", "推断 schema / 类型 / 缺失 / 基数", 2.95, 1.67, 2.45, 0.88,
             COLORS["blue_soft"], COLORS["blue"], 1)
    add_para(slide, "Planning Agent", "拆解任务并选择预处理与模型策略", 5.70, 1.67, 2.65, 0.88,
             COLORS["blue_soft"], COLORS["blue"], 2)

    add_text(slide, "B. 执行与自修复循环", 2.95, 2.83, 2.9, 0.24, 12, COLORS["cyan"], True)
    add_box(slide, 2.95, 3.11, 6.95, 0.03, COLORS["cyan"], line=COLORS["cyan"], radius=False)
    x0, gap, cw = 2.95, 0.18, 1.60
    cards = [
        ("Generator", "生成 sklearn-compatible pipeline code", 3),
        ("Executor", "沙箱运行并收集 outputs", 4),
        ("Classifier", "映射到 13 类失败", 5),
        ("Repair Agent", "定向 repair / regenerate code", 6),
    ]
    for idx, (title, body, n) in enumerate(cards):
        add_para(slide, title, body, x0 + idx * (cw + gap), 3.28, cw, 0.92,
                 COLORS["cyan_soft"], COLORS["cyan"], n)
        if idx < 3:
            add_arrow(slide, x0 + idx * (cw + gap) + cw + 0.02, 3.74, x0 + (idx + 1) * (cw + gap) - 0.05, 3.74,
                      color="6FAFC1", width=1.1)
    add_arrow(slide, 8.95, 4.08, 4.25, 4.08, color="6FAFC1", width=1.1, dashed=True)
    add_text(slide, "failure → repair / regenerate", 5.18, 4.12, 2.25, 0.2, 8.4, COLORS["cyan"], True,
             PP_ALIGN.CENTER, font="Aptos")
    add_text(slide, "Runtime feedback：tracebacks / metrics / artifacts / cost", 6.55, 2.66, 3.02, 0.22,
             8.6, COLORS["muted"], font="Aptos")
    add_arrow(slide, 7.95, 2.92, 7.95, 3.26, color="AABACD", width=1.0, dashed=True)

    add_text(slide, "C. 可靠性、控制与可复现", 2.95, 4.62, 3.3, 0.24, 12, COLORS["green"], True)
    add_box(slide, 2.95, 4.90, 6.95, 0.03, COLORS["green"], line=COLORS["green"], radius=False)
    add_para(slide, "Leakage Validator", "9 项泄露检查：split-before-transform / no target leakage / no test-set tuning",
             2.95, 5.05, 2.05, 0.82, COLORS["green_soft"], COLORS["green"], 7)
    add_para(slide, "Budget Controller", "跟踪 LLM calls / executions / time / trials / repairs；决定 continue / stop / terminate",
             5.28, 5.05, 2.17, 0.82, COLORS["green_soft"], COLORS["green"], 8)
    add_para(slide, "Experiment Logger", "记录 trajectories / errors / fixes / metrics / decisions",
             7.72, 5.05, 2.18, 0.82, COLORS["green_soft"], COLORS["green"], 9)

    add_arrow(slide, 2.52, 3.22, 2.88, 3.22, color="8093AD", width=1.2)
    add_arrow(slide, 8.37, 2.11, 10.50, 2.11, color="8093AD", width=1.2)
    add_arrow(slide, 9.90, 5.46, 10.50, 3.62, color="8093AD", width=1.2)

    # Bottom assurance layer.
    add_box(slide, 0.55, 6.23, 12.22, 0.86, "E9EEF7", line="D4DEED")
    bottom = [
        ("长期记忆", "case history for future tasks；experience for future runs", COLORS["blue"]),
        ("信号捕获", "errors / tracebacks / validation score / cost / pipeline artifacts", COLORS["amber"]),
        ("设计原则", "structured planning / targeted self-healing / leakage-aware / budget-aware", COLORS["green"]),
    ]
    for i, (title, body, accent) in enumerate(bottom):
        bx = 0.83 + i * 4.00
        add_box(slide, bx, 6.43, 0.08, 0.42, accent, line=accent, radius=False)
        add_text(slide, title, bx + 0.18, 6.34, 0.92, 0.22, 11.2, COLORS["ink"], True)
        add_text(slide, body, bx + 0.18, 6.60, 3.35, 0.36, 7.8, COLORS["muted"], font="Aptos")

    prs.save(PPTX_PATH)


def load_font(size: int, bold: bool = False):
    font_dir = Path(r"C:\Windows\Fonts")
    for name in (["msyhbd.ttc", "msyh.ttc"] if bold else ["msyh.ttc", "segoeui.ttf"]):
        path = font_dir / name
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def render_preview() -> None:
    scale = 144
    img = Image.new("RGB", (int(W * scale), int(H * scale)), "#" + COLORS["bg"])
    draw = ImageDraw.Draw(img)

    def xywh(x, y, w, h):
        return [int(x * scale), int(y * scale), int((x + w) * scale), int((y + h) * scale)]

    def rect(x, y, w, h, fill, outline=None, r=12):
        draw.rounded_rectangle(xywh(x, y, w, h), radius=r, fill="#" + fill, outline="#" + (outline or fill), width=2)

    def text(t, x, y, w, size, fill, bold=False, align="left", max_lines=None):
        font = load_font(size, bold)
        max_chars = max(4, int(w * 11.5 * (15 / max(size, 1))))
        lines = []
        for raw in t.split("\n"):
            lines.extend(wrap(raw, max_chars) or [""])
        if max_lines:
            lines = lines[:max_lines]
        yy = int(y * scale)
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=font)
            if align == "center":
                xx = int(x * scale + (w * scale - (bbox[2] - bbox[0])) / 2)
            else:
                xx = int(x * scale)
            draw.text((xx, yy), line, fill="#" + fill, font=font)
            yy += int(size * 1.35)

    # Recreate the same design at preview fidelity.
    draw.rectangle(xywh(0, 0, W, 0.14), fill="#" + COLORS["blue"])
    text("AutoDS-Agent", 0.55, 0.34, 2.0, 26, COLORS["blue"], True)
    text("可靠表格机器学习自动化架构", 2.42, 0.30, 5.9, 34, COLORS["ink"], True)
    text("从任务理解到执行修复，再到泄露校验、预算控制与可复现实验记录的一页式可编辑流程图",
         0.58, 0.78, 8.4, 16, COLORS["muted"])
    rect(10.55, 0.36, 2.12, 0.38, COLORS["navy"], COLORS["navy"], 14)
    text("EDITABLE  •  ONE PAGE", 10.67, 0.44, 1.88, 12, COLORS["white"], True, "center")

    for title, x, accent, items in [
        ("Input", 0.55, COLORS["blue"], ["输入", "数据集 / 表格", "任务类型：classification", "评价指标", "约束 / 预算"]),
        ("Output", 10.58, COLORS["teal"], ["输出", "best pipeline", "predictive score", "cost summary", "leakage-free report"]),
    ]:
        text(title, x + 0.07, 1.22, 1.6, 19, accent, True)
        panel_w = 2.18 if title == "Output" else 2.08
        rect(x, 1.55, panel_w, 3.42, COLORS["white"], "D6DFEC", 14)
        draw.rectangle(xywh(x, 1.55, panel_w, 0.12), fill="#" + accent)
        text(items[0], x + 0.21, 1.86, 1.35, 22, COLORS["ink"], True)
        text("\n".join(items[1:]), x + 0.21, 2.30, panel_w - 0.48, 14, COLORS["muted"])

    def preview_card(x, y, w, h, title, body, fill, accent, n):
        rect(x, y, w, h, fill, "D8E1EF", 12)
        rect(x + w - 0.46, y + 0.13, 0.30, 0.30, accent, accent, 9)
        text(str(n), x + w - 0.46, y + 0.155, 0.30, 13, COLORS["white"], True, "center")
        text(title, x + 0.16, y + 0.12, w - 0.66, 13, COLORS["ink"], True, max_lines=1)
        text(body, x + 0.16, y + 0.45, w - 0.32, 12, COLORS["muted"], max_lines=3)

    sections = [
        ("A. 任务理解", 1.22, COLORS["blue"]),
        ("B. 执行与自修复循环", 2.83, COLORS["cyan"]),
        ("C. 可靠性、控制与可复现", 4.62, COLORS["green"]),
    ]
    for label, y, accent in sections:
        text(label, 2.95, y, 3.3, 16, accent, True)
        draw.rectangle(xywh(2.95, y + 0.28, 6.95, 0.03), fill="#" + accent)

    preview_card(2.95, 1.67, 2.45, 0.88, "Schema Inspector", "推断 schema / 类型 / 缺失 / 基数", COLORS["blue_soft"], COLORS["blue"], 1)
    preview_card(5.70, 1.67, 2.65, 0.88, "Planning Agent", "拆解任务并选择预处理与模型策略", COLORS["blue_soft"], COLORS["blue"], 2)
    for idx, (title, body, n) in enumerate([
        ("Generator", "生成 sklearn-compatible pipeline code", 3),
        ("Executor", "沙箱运行并收集 outputs", 4),
        ("Classifier", "映射到 13 类失败", 5),
        ("Repair Agent", "定向 repair / regenerate code", 6),
    ]):
        preview_card(2.95 + idx * 1.75, 3.28, 1.50, 0.92, title, body, COLORS["cyan_soft"], COLORS["cyan"], n)
    text("Runtime feedback：tracebacks / metrics / artifacts / cost", 6.55, 2.66, 3.02, 12, COLORS["muted"])
    text("failure → repair / regenerate", 5.18, 4.12, 2.25, 12, COLORS["cyan"], True, "center")
    preview_card(2.95, 5.05, 2.05, 0.82, "Leakage Validator", "9 项泄露检查：split-before-transform / no target leakage / no test-set tuning", COLORS["green_soft"], COLORS["green"], 7)
    preview_card(5.28, 5.05, 2.17, 0.82, "Budget Controller", "跟踪 LLM calls / executions / time / trials / repairs；决定 continue / stop / terminate", COLORS["green_soft"], COLORS["green"], 8)
    preview_card(7.72, 5.05, 2.18, 0.82, "Experiment Logger", "记录 trajectories / errors / fixes / metrics / decisions", COLORS["green_soft"], COLORS["green"], 9)

    rect(0.55, 6.23, 12.22, 0.86, "E9EEF7", "D4DEED", 14)
    for i, (title, body, accent) in enumerate([
        ("长期记忆", "case history for future tasks；experience for future runs", COLORS["blue"]),
        ("信号捕获", "errors / tracebacks / validation score / cost / pipeline artifacts", COLORS["amber"]),
        ("设计原则", "structured planning / targeted self-healing / leakage-aware / budget-aware", COLORS["green"]),
    ]):
        bx = 0.83 + i * 4.00
        draw.rectangle(xywh(bx, 6.43, 0.08, 0.42), fill="#" + accent)
        text(title, bx + 0.18, 6.34, 0.92, 15, COLORS["ink"], True)
        text(body, bx + 0.18, 6.60, 3.35, 11, COLORS["muted"], max_lines=2)

    img.save(PNG_PATH)


if __name__ == "__main__":
    build_deck()
    render_preview()
    print(PPTX_PATH)
    print(PNG_PATH)
